"""
simulate_brain.py — Полная симуляция автономного Агента

Покрывает ВСЕ 15 инструментов + глубокая интроспекция каждого шага Brain:

  ЦИКЛ Brain.think() (9 шагов):
    1. Memory.store(input)
    2. ContextBuilder.build()  <- показывает history/goals/facts
    3. Fast path check
    4. LLM.call()              <- ScriptedLLM очередь
    5. Interpreter.interpret() <- action + confidence + reasoning
    6. UncertaintyEstimator    <- calibrated_confidence, threshold, should_act
    7. Explainer               <- risk_level, summary, reasoning_chain
    8. Memory.store(response)
    9. GoalStack.update()      <- активные цели

  ИНСТРУМЕНТЫ (15):
    calculator / datetime / kv_set / kv_get / kv_delete
    text_summarizer / diagram_generator / memory_search
    pdf_reader / docx_reader / xlsx_reader / pptx_reader / odt_reader
    image_reader / media_reader (error path)

Запуск:
    python simulate_brain.py                  # полная симуляция
    python simulate_brain.py --part 1         # только базовые инструменты
    python simulate_brain.py --part 2         # только файловые ридеры
    python simulate_brain.py --part 3         # только спецсценарии
    python simulate_brain.py --no-deep        # без глубокой интроспекции Brain
"""

from __future__ import annotations

import argparse
import collections
import logging
import sys
import tempfile
import time
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from rich import box
from rich.console import Console
from rich.panel import Panel
from rich.rule import Rule
from rich.table import Table

from brain.core import Brain, ThinkResult
from brain.interfaces.llm_interface import LLMInterface
from brain.interfaces.memory_interface import MemoryInterface

from tools import ToolExecutor, ToolRegistry
from tools.builtins import (
    CalculatorTool,
    DateTimeTool,
    DiagramGeneratorTool,
    DocxReaderTool,
    ImageReaderTool,
    KeyValueDeleteTool,
    KeyValueGetTool,
    KeyValueSetTool,
    KeyValueStore,
    MediaReaderTool,
    MemorySearchTool,
    OdtReaderTool,
    PdfReaderTool,
    PptxReaderTool,
    TextSummarizerTool,
    XlsxReaderTool,
)


# =======================================================================
#  SCRIPTED LLM
# =======================================================================

class ScriptedLLM(LLMInterface):
    """Детерминированная очередь ответов — вместо реального LLM API."""

    def __init__(self) -> None:
        self._queue: collections.deque[dict] = collections.deque()
        self.total_calls: int = 0
        self.last_context: dict = {}

    def enqueue(self, *responses: dict) -> "ScriptedLLM":
        for r in responses:
            self._queue.append(r)
        return self

    def call(self, context: dict) -> dict:
        self.total_calls += 1
        self.last_context = context
        if self._queue:
            return self._queue.popleft()
        return {"action": "respond",
                "content": "[ScriptedLLM: очередь пуста]",
                "confidence": 0.5,
                "reasoning": "No more scripted responses"}

    def is_available(self) -> bool:
        return True

    @property
    def model_name(self) -> str:
        return "ScriptedLLM-v1 (симуляция)"


# =======================================================================
#  SIMPLE MEMORY
# =======================================================================

class SimpleMemory(MemoryInterface):
    def __init__(self) -> None:
        self._history: dict[str, list[dict]] = {}

    def recall_history(self, session_id: str, limit: int = 10) -> list[dict]:
        return self._history.get(session_id, [])[-limit:]

    def recall_facts(self, query: str, top_k: int = 5) -> list[dict]:
        return []

    def store(self, session_id: str, role: str, content: str) -> None:
        self._history.setdefault(session_id, []).append(
            {"role": role, "content": content}
        )

    def forget(self, session_id: str) -> None:
        self._history.pop(session_id, None)

    def retrieve(self, session_id: str) -> list[dict]:
        """Используется MemorySearchTool."""
        return self._history.get(session_id, [])

    def count(self, session_id: str) -> int:
        return len(self._history.get(session_id, []))


# =======================================================================
#  ФАЙЛОВАЯ ФАБРИКА
# =======================================================================

class FileFactory:
    """Создаёт минимальные тестовые файлы для каждого формата."""

    def __init__(self, tmp_dir: Path) -> None:
        self.tmp = tmp_dir
        self.files: dict[str, Path] = {}
        self.errors: dict[str, str] = {}

    def build_all(self) -> None:
        for name, maker in [
            ("pdf",  self._make_pdf),
            ("docx", self._make_docx),
            ("xlsx", self._make_xlsx),
            ("pptx", self._make_pptx),
            ("odt",  self._make_odt),
            ("png",  self._make_image),
        ]:
            try:
                self.files[name] = maker()
            except Exception as exc:  # noqa: BLE001
                self.errors[name] = str(exc)

    def _make_pdf(self) -> Path:
        from pypdf import PdfWriter  # type: ignore
        w = PdfWriter()
        w.add_blank_page(width=595, height=842)
        w.add_metadata({"/Title": "Тест PDF — Симуляция Агента",
                        "/Author": "Brain Simulator"})
        p = self.tmp / "test.pdf"
        with open(p, "wb") as f:
            w.write(f)
        return p

    def _make_docx(self) -> Path:
        from docx import Document  # type: ignore
        doc = Document()
        doc.add_heading("Симуляция автономного агента", level=1)
        doc.add_paragraph(
            "Brain — когнитивное ядро агента. "
            "Он принимает решения и делегирует задачи инструментам."
        )
        doc.add_paragraph(
            "Компоненты: ContextBuilder, Interpreter, UncertaintyEstimator, Explainer."
        )
        p = self.tmp / "test.docx"
        doc.save(str(p))
        return p

    def _make_xlsx(self) -> Path:
        from openpyxl import Workbook  # type: ignore
        wb = Workbook()
        ws = wb.active
        ws.title = "Метрики"
        ws["A1"], ws["B1"], ws["C1"] = "Инструмент", "Вызовов", "OK"
        for row, (tool, calls, ok) in enumerate([
            ("calculator", 42, 42), ("datetime", 15, 15),
            ("kv_set", 8, 8), ("memory_search", 5, 4),
        ], start=2):
            ws.cell(row=row, column=1, value=tool)
            ws.cell(row=row, column=2, value=calls)
            ws.cell(row=row, column=3, value=ok)
        p = self.tmp / "test.xlsx"
        wb.save(str(p))
        return p

    def _make_pptx(self) -> Path:
        from pptx import Presentation  # type: ignore
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Архитектура автономного агента"
        slide.placeholders[1].text = "Brain → LLM → Tools → Memory"
        p = self.tmp / "test.pptx"
        prs.save(str(p))
        return p

    def _make_odt(self) -> Path:
        from odf.opendocument import OpenDocumentText  # type: ignore
        from odf.text import P  # type: ignore
        doc = OpenDocumentText()
        doc.text.addElement(
            P(text="Brain читает OpenDocument файл через odt_reader.")
        )
        p = self.tmp / "test.odt"
        doc.save(str(p))
        return p

    def _make_image(self) -> Path:
        from PIL import Image, ImageDraw  # type: ignore
        img = Image.new("RGB", (300, 200), color=(20, 30, 60))
        draw = ImageDraw.Draw(img)
        draw.rectangle([10, 10, 289, 189], outline=(0, 200, 150), width=3)
        draw.text((50, 85), "BRAIN SIMULATION", fill=(0, 255, 200))
        p = self.tmp / "test.png"
        img.save(str(p))
        return p


# =======================================================================
#  СЦЕНАРИИ
# =======================================================================

@dataclass
class ScenarioStep:
    user_input: str
    llm_responses: list[dict]
    expected_action: str = "respond"
    expected_tool: str | None = None
    allow_tool_fail: bool = False

@dataclass
class Scenario:
    name: str
    icon: str
    description: str
    part: int
    steps: list[ScenarioStep]
    tags: list[str] = field(default_factory=list)


def _r(content: str, conf: float = 0.97) -> dict:
    return {"action": "respond", "content": content, "confidence": conf,
            "reasoning": "Получен результат от инструмента — формулирую ответ"}

def _tc(tool: str, params: dict, conf: float = 0.95, why: str = "") -> dict:
    return {"action": "tool_call",
            "content": {"tool_name": tool, "params": params},
            "confidence": conf,
            "reasoning": why or f"Делегирую задачу инструменту {tool}"}


LONG_TEXT = (
    "Автономный агент — это программная система, способная принимать решения самостоятельно. "
    "В отличие от традиционных программ, агент воспринимает окружение и адаптируется. "
    "Ключевые компоненты: восприятие, планирование, исполнение. "
    "LLM используется как модуль рассуждений, но Brain сам решает когда его вызвать. "
    "Такой подход минимизирует галлюцинации и делает систему предсказуемой."
)

MERMAID = (
    "graph TD\n"
    "  IN[Ввод] --> B{Brain.think}\n"
    "  B --> CB[ContextBuilder]\n"
    "  CB --> LLM[LLM.call]\n"
    "  LLM --> I[Interpreter]\n"
    "  I --> UG{Uncertainty Gate}\n"
    "  UG -->|pass| A{Action}\n"
    "  A -->|respond| OUT[Ответ]\n"
    "  A -->|tool_call| TE[ToolExecutor]\n"
    "  TE --> TR[ToolResult]\n"
    "  TR --> B"
)


def build_scenarios(files: dict[str, Path]) -> list[Scenario]:
    fp = lambda key: str(files.get(key, Path(f"/tmp/test.{key}")))

    return [

        # ── ЧАСТЬ 1: БАЗОВЫЕ ИНСТРУМЕНТЫ ─────────────────────────────

        Scenario("Математика", "🧮", "calculator: (15*7)+42", 1, [ScenarioStep(
            "Посчитай: (15 * 7) + 42",
            [_tc("calculator", {"expression": "(15 * 7) + 42"},
                 why="Математическое выражение — делегирую calculator"),
             _r("(15 × 7) + 42 = **147**. Вычислено через calculator.")],
            expected_tool="calculator",
        )], tags=["calculator"]),

        Scenario("Дата и время", "🕐", "datetime: текущее время UTC", 1, [ScenarioStep(
            "Который сейчас час? Какой день недели? UTC.",
            [_tc("datetime", {"action": "now"},
                 why="Запрос текущего времени — вызываю datetime"),
             _r("Время и дата получены актуально!")],
            expected_tool="datetime",
        )], tags=["datetime"]),

        Scenario("Память: SET", "💾", "kv_set: сохранить значение", 1, [ScenarioStep(
            "Запомни: agent_name = BrainBot 3000",
            [_tc("kv_set", {"key": "agent_name", "value": "BrainBot 3000"},
                 why="Нужно сохранить данные в памяти — kv_set"),
             _r("Запомнил! agent_name = 'BrainBot 3000' сохранён.")],
            expected_tool="kv_set",
        )], tags=["kv_set", "memory"]),

        Scenario("Память: GET", "📖", "kv_get: прочитать значение", 1, [ScenarioStep(
            "Как меня зовут? Вспомни agent_name.",
            [_tc("kv_get", {"key": "agent_name"},
                 why="Пользователь хочет вспомнить — kv_get"),
             _r("Вас зовут: **BrainBot 3000** (из памяти)")],
            expected_tool="kv_get",
        )], tags=["kv_get", "memory"]),

        Scenario("Память: DELETE", "🗑️", "kv_delete: удалить ключ", 1, [ScenarioStep(
            "Забудь agent_name — удали из памяти.",
            [_tc("kv_delete", {"key": "agent_name"},
                 why="Пользователь просит удалить ключ — kv_delete"),
             _r("Ключ 'agent_name' удалён из памяти.")],
            expected_tool="kv_delete",
        )], tags=["kv_delete", "memory"]),

        Scenario("Суммаризация", "📝", "text_summarizer: сжать текст", 1, [ScenarioStep(
            f"Сократи в 2 предложения:\n{LONG_TEXT}",
            [_tc("text_summarizer", {"text": LONG_TEXT, "max_sentences": 2},
                 why="Нужна суммаризация — text_summarizer"),
             _r("Краткое изложение готово!")],
            expected_tool="text_summarizer",
        )], tags=["text_summarizer"]),

        Scenario("Mermaid-диаграмма", "📊", "diagram_generator: flowchart", 1, [ScenarioStep(
            "Нарисуй Mermaid-схему цикла Brain->Tool->Brain",
            [_tc("diagram_generator",
                 {"diagram_type": "mermaid", "content": MERMAID, "title": "Brain-Tool цикл"},
                 why="Создаю диаграмму через diagram_generator"),
             _r("Mermaid-диаграмма создана! .mmd и .html сохранены.")],
            expected_tool="diagram_generator",
        )], tags=["diagram_generator", "mermaid"]),

        Scenario("Поиск в памяти", "🔍", "memory_search: keyword в истории", 1, [ScenarioStep(
            "Найди в памяти всё связанное с 'агент'",
            [_tc("memory_search",
                 {"session_id": "sim_session", "query": "агент", "limit": 5},
                 why="Поиск по памяти — memory_search"),
             _r("Поиск завершён. Найдены релевантные записи.")],
            expected_tool="memory_search",
        )], tags=["memory_search"]),

        # ── ЧАСТЬ 2: ФАЙЛОВЫЕ РИДЕРЫ ──────────────────────────────────

        Scenario("PDF Reader", "📄", "pdf_reader: читать PDF", 2, [ScenarioStep(
            f"Прочитай PDF метаданные: {fp('pdf')}",
            [_tc("pdf_reader", {"file_path": fp("pdf"), "action": "get_metadata"},
                 why="Нужно прочитать PDF — pdf_reader"),
             _r("PDF прочитан! Метаданные: заголовок, автор, страниц.")],
            expected_tool="pdf_reader",
        )], tags=["pdf_reader", "files"]),

        Scenario("DOCX Reader", "📃", "docx_reader: читать Word", 2, [ScenarioStep(
            f"Извлеки текст из Word: {fp('docx')}",
            [_tc("docx_reader", {"file_path": fp("docx")},
                 why="Нужен текст из .docx — docx_reader"),
             _r("Word документ прочитан! Текст извлечён.")],
            expected_tool="docx_reader",
        )], tags=["docx_reader", "files"]),

        Scenario("XLSX Reader", "📊", "xlsx_reader: читать Excel", 2, [ScenarioStep(
            f"Прочитай Excel листы: {fp('xlsx')}",
            [_tc("xlsx_reader", {"file_path": fp("xlsx"), "action": "get_sheets"},
                 why="Нужно прочитать .xlsx — xlsx_reader"),
             _r("Excel прочитан! Лист 'Метрики' найден.")],
            expected_tool="xlsx_reader",
        )], tags=["xlsx_reader", "files"]),

        Scenario("PPTX Reader", "📑", "pptx_reader: читать PowerPoint", 2, [ScenarioStep(
            f"Перечисли слайды: {fp('pptx')}",
            [_tc("pptx_reader", {"file_path": fp("pptx"), "action": "list_slides"},
                 why="Нужны слайды из .pptx — pptx_reader"),
             _r("Презентация прочитана! 1 слайд найден.")],
            expected_tool="pptx_reader",
        )], tags=["pptx_reader", "files"]),

        Scenario("ODT Reader", "📋", "odt_reader: читать OpenDocument", 2, [ScenarioStep(
            f"Извлеки текст из ODT: {fp('odt')}",
            [_tc("odt_reader", {"file_path": fp("odt")},
                 why="Читаю .odt через odt_reader"),
             _r("OpenDocument прочитан! Текст извлечён.")],
            expected_tool="odt_reader",
        )], tags=["odt_reader", "files"]),

        Scenario("Image Reader", "🖼️", "image_reader: анализ PNG", 2, [ScenarioStep(
            f"Проанализируй изображение: {fp('png')}",
            [_tc("image_reader", {"file_path": fp("png"), "action": "get_metadata"},
                 why="Анализирую изображение через image_reader"),
             _r("Изображение: PNG 300x200px RGB.")],
            expected_tool="image_reader",
        )], tags=["image_reader", "files"]),

        # ── ЧАСТЬ 3: СПЕЦСЦЕНАРИИ ─────────────────────────────────────

        Scenario("Media: error path", "🎵",
                 "media_reader: Brain обрабатывает ошибку", 3, [ScenarioStep(
            "Прочитай аудио: /tmp/nonexistent_audio.mp3",
            [_tc("media_reader", {"file_path": "/tmp/nonexistent_audio.mp3"},
                 why="Читаю метаданные аудио через media_reader"),
             _r("Файл не найден. Проверьте путь и повторите.")],
            expected_tool="media_reader",
            allow_tool_fail=True,
        )], tags=["media_reader", "error_handling"]),

        Scenario("Многошаговый план", "🔗",
                 "calculator -> kv_set -> respond: цепочка инструментов", 3, [ScenarioStep(
            "Посчитай периметр прямоугольника 15x7, сохрани и скажи мне",
            [_tc("calculator", {"expression": "2*(15+7)"},
                 why="Сначала вычислю периметр через calculator"),
             _tc("kv_set", {"key": "perimeter", "value": "44"},
                 why="Сохраняю результат через kv_set"),
             _r("Периметр = 44 (2*(15+7)). Сохранён в памяти под 'perimeter'.")],
            expected_tool="calculator",
        )], tags=["multi-step", "calculator", "kv_set"]),
    ]


# =======================================================================
#  РЕГИСТРАЦИЯ ИНСТРУМЕНТОВ (все 15)
# =======================================================================

def setup_registry(memory: MemoryInterface) -> tuple[ToolRegistry, ToolExecutor]:
    store = KeyValueStore()
    registry = ToolRegistry()
    for tool in [
        CalculatorTool(),
        DateTimeTool(),
        KeyValueGetTool(store=store),
        KeyValueSetTool(store=store),
        KeyValueDeleteTool(store=store),
        TextSummarizerTool(),
        DiagramGeneratorTool(),
        MemorySearchTool(memory=memory),
        PdfReaderTool(),
        DocxReaderTool(),
        XlsxReaderTool(),
        PptxReaderTool(),
        OdtReaderTool(),
        ImageReaderTool(),
        MediaReaderTool(),
    ]:
        registry.register(tool)
    return registry, ToolExecutor(registry)


# =======================================================================
#  ГЛУБОКАЯ ИНТРОСПЕКЦИЯ BRAIN (9 шагов)
# =======================================================================

RISK_COLOR = {"low": "green", "medium": "yellow", "high": "red", "critical": "bold red"}
RISK_ICON  = {"low": "[OK]", "medium": "[!!]", "high": "[HI]", "critical": "[!!]"}


def _step(console: Console, n: int, name: str, detail: str = "", ok: bool = True) -> None:
    col = "cyan" if ok else "dim red"
    det = f"  [dim]{detail}[/]" if detail else ""
    console.print(f"   [{col}]{'●' if ok else 'x'} {n:>2}. {name:<24}[/]{det}")


def transparent_think(
    brain: Brain,
    raw_input: str,
    session_id: str,
    console: Console,
    deep: bool = True,
    turn_num: int = 1,
    is_feedback: bool = False,
) -> ThinkResult:
    """Запускает Brain.think() с пошаговым показом всех 9 внутренних этапов."""
    prefix = "📨 tool_result → Brain" if is_feedback else f"🧠 Brain.think() #{turn_num}"
    t_total = time.perf_counter()

    console.print()
    inp_preview = raw_input[:80] + ("…" if len(raw_input) > 80 else "")
    console.print(f"  [bold dim]{prefix}[/]  [italic dim]{inp_preview}[/]")

    if not deep:
        result = brain.think(raw_input, session_id)
        _print_result_short(console, result, (time.perf_counter() - t_total) * 1000)
        return result

    # ──── ШАГ 1: Сохранить ввод ─────────────────────────────────────
    brain._memory.store(session_id, "tool" if is_feedback else "user", raw_input)
    mem_n = brain._memory.count(session_id)
    _step(console, 1, "memory.store(input)", f"total={mem_n} записей")

    # ──── ШАГ 2: Построить контекст ─────────────────────────────────
    goals_list = brain._goals.current()
    context = brain._context_builder.build(
        raw_input=raw_input, goals=goals_list, session_id=session_id,
    )
    h_n = len(context.get("history", []))
    g_n = len(goals_list)
    f_n = len(context.get("facts", []))
    _step(console, 2, "context_builder.build()", f"history={h_n}  goals={g_n}  facts={f_n}")

    # ──── ШАГ 3: Fast path ──────────────────────────────────────────
    fast = brain._try_fast_path(raw_input, context)
    if fast:
        _step(console, 3, "fast_path CHECK", f"TRIGGERED -> action={fast.action}")
        fast.explanation = brain._explainer.explain(fast, context, None)
        _print_result_short(console, fast, (time.perf_counter() - t_total) * 1000)
        return fast
    _step(console, 3, "fast_path CHECK", "skip -> LLM needed")

    # ──── ШАГ 4: LLM.call() ─────────────────────────────────────────
    t0 = time.perf_counter()
    llm_out = brain._llm.call(context)
    llm_ms  = (time.perf_counter() - t0) * 1000
    a_raw   = llm_out.get("action", "?")
    c_raw   = llm_out.get("confidence", 0.0)
    tn_raw  = ""
    if isinstance(llm_out.get("content"), dict):
        tn_raw = llm_out["content"].get("tool_name", "")
    _step(console, 4, "LLM.call()",
          f"action={a_raw}  conf={c_raw:.0%}  tool={tn_raw or '—'}  ({llm_ms:.0f}ms)")

    # ──── ШАГ 5: Interpret ──────────────────────────────────────────
    result = brain._interpreter.interpret(llm_out, context)
    content_desc = (
        f"tool={result.content.get('tool_name','?')}" if isinstance(result.content, dict)
        else repr(str(result.content)[:40])
    )
    _step(console, 5, "interpreter.interpret()", f"action={result.action}  {content_desc}")

    # ──── ШАГ 6: Uncertainty gate ───────────────────────────────────
    uncertainty = brain._uncertainty.estimate(
        llm_confidence=result.confidence, context=context,
    )
    gate_ok = uncertainty.should_act
    gate_col = "green" if gate_ok else "red"
    _step(console, 6, "uncertainty.estimate()",
          f"llm={result.confidence:.0%} -> cal={uncertainty.calibrated_confidence:.0%}  "
          f"thr={uncertainty.threshold_used:.0%}  [{gate_col}]act={gate_ok}[/]")

    if not gate_ok:
        _step(console, 7, "BLOCKED", uncertainty.reasoning[:60], ok=False)
        blocked = ThinkResult(
            action="wait", content="Уверенность ниже порога.",
            confidence=uncertainty.calibrated_confidence,
            reasoning=uncertainty.reasoning, needs_human_approval=True,
        )
        blocked.explanation = brain._explainer.explain(blocked, context, uncertainty)
        _print_result_short(console, blocked, (time.perf_counter() - t_total) * 1000)
        return blocked

    # ──── ШАГ 7: Explain ────────────────────────────────────────────
    result.explanation = brain._explainer.explain(result, context, uncertainty)
    risk = result.explanation.risk_level.value
    rc   = RISK_COLOR.get(risk, "white")
    _step(console, 7, "explainer.explain()",
          f"risk=[{rc}]{risk}[/]  summary={result.explanation.summary[:55]}…")

    # ──── ШАГ 8: Сохранить ответ ────────────────────────────────────
    if result.action in {"respond", "clarify"} and result.content:
        brain._memory.store(session_id, "assistant", str(result.content))
    mem_after = brain._memory.count(session_id)
    _step(console, 8, "memory.store(response)", f"total={mem_after} записей")

    # ──── ШАГ 9: Обновить цели ──────────────────────────────────────
    brain._goals.update(result)
    g_after = len(brain._goals.current())
    _step(console, 9, "goal_stack.update()", f"active goals={g_after}")

    elapsed = (time.perf_counter() - t_total) * 1000
    _print_result_full(console, result, elapsed, uncertainty)
    return result


def _print_result_short(console: Console, result: ThinkResult, ms: float) -> None:
    ac = {"respond": "green", "tool_call": "yellow", "wait": "dim", "stop": "red"}
    col = ac.get(result.action, "white")
    risk = result.explanation.risk_level.value if result.explanation else "?"
    rc = RISK_COLOR.get(risk, "dim")
    extra = ""
    if result.action == "tool_call" and isinstance(result.content, dict):
        extra = f"  -> [yellow]{result.content.get('tool_name', '?')}[/]"
    elif result.action == "respond":
        extra = f"  -> [green]\"{str(result.content)[:50]}\"[/]"
    console.print(
        f"   [dim]->[/] [{col}]{result.action}[/]  "
        f"conf={result.confidence:.0%}  risk=[{rc}]{risk}[/]  "
        f"{ms:.0f}ms{extra}"
    )


def _print_result_full(
    console: Console, result: ThinkResult, ms: float, uncertainty: Any
) -> None:
    col = {"respond": "green", "tool_call": "yellow", "wait": "dim", "stop": "red"}.get(
        result.action, "white"
    )
    risk = result.explanation.risk_level.value if result.explanation else "?"
    rc   = RISK_COLOR.get(risk, "white")

    lines = [
        f"[{col}]action:[/]       {result.action}",
        f"[dim]confidence:[/]   {result.confidence:.0%}  "
        f"(calibrated: {uncertainty.calibrated_confidence:.0%}  "
        f"threshold: {uncertainty.threshold_used:.0%}  "
        f"uncertainty_score: {uncertainty.uncertainty_score:.2f})",
        f"[dim]reasoning:[/]    {result.reasoning}",
        f"[{rc}]risk:[/]         {risk.upper()}",
        f"[dim]elapsed:[/]      {ms:.1f}ms",
    ]
    if result.explanation and result.explanation.summary:
        lines.append(f"[dim]summary:[/]     {result.explanation.summary[:90]}")

    if result.action == "tool_call" and isinstance(result.content, dict):
        tn = result.content.get("tool_name", "?")
        ps = result.content.get("params", {})
        ps_str = ", ".join(f"{k}={repr(v)[:30]}" for k, v in ps.items())
        lines.append(f"[yellow]=> tool_call:[/]  [bold]{tn}[/]({ps_str})")
    elif result.action == "respond":
        lines.append(f"[green]=> content:[/]   {str(result.content)[:120]}")

    console.print(Panel("\n".join(lines), border_style=col, padding=(0, 1)))


# =======================================================================
#  ОТОБРАЖЕНИЕ ЭЛЕМЕНТОВ
# =======================================================================

def print_header(console: Console, n_tools: int, n_scenarios: int) -> None:
    console.print()
    console.rule("[bold cyan]🧠 ПОЛНАЯ СИМУЛЯЦИЯ — Brain + ВСЕ ИНСТРУМЕНТЫ[/]")
    console.print()
    t = Table(box=box.ROUNDED, show_header=False, padding=(0, 1))
    t.add_column("", style="dim", width=16)
    t.add_column("")
    t.add_row("Brain цикл",   "[cyan]9 шагов: Memory -> Context -> FastPath -> LLM -> Interpret -> Uncertainty -> Explain -> Store -> Goals[/]")
    t.add_row("Инструментов", f"[green]{n_tools}[/] зарегистрировано")
    t.add_row("Сценариев",    f"[yellow]{n_scenarios}[/]  (Часть 1: базовые | Часть 2: файловые ридеры | Часть 3: спецсценарии)")
    t.add_row("LLM",          "[dim]ScriptedLLM — детерминированные ответы без API-ключа[/]")
    console.print(t)
    console.print()


def print_tool_list(console: Console, registry: ToolRegistry) -> None:
    names = sorted(registry.names())
    console.print(f"  [dim]Зарегистрировано: [bold]{len(names)}[/] инструментов[/]")
    console.print("  " + "  ".join(f"[cyan]{n}[/]" for n in names))
    console.print()


def print_part_banner(console: Console, part: int) -> None:
    labels = {
        1: "ЧАСТЬ 1 — БАЗОВЫЕ ИНСТРУМЕНТЫ",
        2: "ЧАСТЬ 2 — ФАЙЛОВЫЕ РИДЕРЫ",
        3: "ЧАСТЬ 3 — СПЕЦСЦЕНАРИИ",
    }
    console.print()
    console.rule(f"[bold cyan]{labels.get(part, '')}[/]")


def print_scenario_header(
    console: Console, idx: int, total: int, sc: Scenario
) -> None:
    console.print()
    console.rule(
        f"[bold yellow]{sc.icon}  #{idx}/{total}  {sc.name.upper()}[/]  "
        f"[dim]{sc.description}[/]"
    )
    if sc.tags:
        console.print("  " + "  ".join(f"[cyan]#{t}[/]" for t in sc.tags))
    console.print()


def print_user_input(console: Console, text: str) -> None:
    display = text[:200] + ("…" if len(text) > 200 else "")
    console.print(Panel(f"[bold white]{display}[/]",
                        title="[blue]👤 ВВОД[/]", border_style="blue", padding=(0, 1)))


def print_tool_exec(
    console: Console, tool: str, params: dict,
    success: bool, output: Any, error: str | None, ms: float,
) -> None:
    status = "[green]SUCCESS[/]" if success else "[red]FAIL[/]"
    border = "green" if success else "red"
    ps = {k: (str(v)[:50] + "…" if len(str(v)) > 50 else v) for k, v in params.items()}
    lines = [
        f"[bold]{tool}[/]  {status}  [dim]({ms:.1f}ms)[/]",
        f"[dim]params:[/] {ps}",
    ]
    if success:
        out = str(output)
        lines.append(f"[green]output:[/] {out[:200]}{'…' if len(out) > 200 else ''}")
    else:
        lines.append(f"[red]error:[/]  {error}")
    console.print(Panel("\n".join(lines), title="[yellow]🔧 TOOL EXECUTOR[/]",
                        border_style=border, padding=(0, 1)))


def print_feedback(console: Console, content: str) -> None:
    prev = content[:100] + ("…" if len(content) > 100 else "")
    console.print(f"\n  [dim]<- feedback -> Brain:[/] [italic dim]{prev}[/]\n")


def print_final_response(console: Console, content: str) -> None:
    console.print(Panel(f"[bold green]{str(content)[:300]}[/]",
                        title="[green]ОТВЕТ АГЕНТА[/]",
                        border_style="green", padding=(0, 1)))


def print_memory_state(console: Console, mem: SimpleMemory, session: str) -> None:
    n = mem.count(session)
    console.print(f"  [dim]  Память сессии: {n} записей[/]")


def print_step_result(
    console: Console, passed: bool, reason: str,
    ms: float, llm_n: int, tool: str,
) -> None:
    status = "[bold green]PASS[/]" if passed else "[bold red]FAIL[/]"
    console.print(
        f"\n  {status}  [dim]tool=[/][yellow]{tool}[/]  "
        f"[dim]llm_calls={llm_n}  time={ms:.0f}ms[/]"
    )
    if not passed:
        console.print(f"  [red]-> {reason}[/]")
    console.print()


# =======================================================================
#  ЯДРО СИМУЛЯЦИИ
# =======================================================================

SESSION = "sim_session"


@dataclass
class TurnResult:
    final_action: str
    final_content: Any
    tool_name: str | None
    tool_success: bool | None
    llm_calls: int
    elapsed_ms: float
    passed: bool
    fail_reason: str = ""


def run_turn(
    brain: Brain,
    executor: ToolExecutor,
    llm: ScriptedLLM,
    memory: SimpleMemory,
    user_input: str,
    step: ScenarioStep,
    console: Console,
    deep: bool,
    turn_base: int = 1,
) -> TurnResult:
    print_user_input(console, user_input)

    llm_before  = llm.total_calls
    first_tool: str | None = None
    tool_success: bool | None = None
    t_start = time.perf_counter()
    current = user_input
    final: ThinkResult | None = None

    for i in range(1, 10):
        is_fb = (i > 1)
        result = transparent_think(
            brain, current, SESSION, console,
            deep=deep, turn_num=turn_base + i - 1, is_feedback=is_fb,
        )

        if result.action == "respond":
            final = result
            print_final_response(console, result.content)
            break

        if result.action == "tool_call":
            content = result.content
            if not isinstance(content, dict):
                final = result
                break
            tn     = content.get("tool_name", "")
            params = content.get("params", {})
            if first_tool is None:
                first_tool = tn

            t0 = time.perf_counter()
            tr = executor.run(tn, params, allow_destructive=True)
            tool_elapsed = (time.perf_counter() - t0) * 1000
            tool_success = tr.success

            print_tool_exec(console, tn, params, tr.success, tr.output, tr.error, tool_elapsed)

            current = (
                f"[tool_result] {tn} succeeded: {tr.output}"
                if tr.success else
                f"[tool_result] {tn} failed: {tr.error}"
            )
            print_feedback(console, current)
            continue

        final = result
        break

    elapsed   = (time.perf_counter() - t_start) * 1000
    llm_calls = llm.total_calls - llm_before

    # Валидация
    passed, reason = True, ""
    if final is None:
        passed, reason = False, "Brain не вернул финальный ответ"
    elif final.action != step.expected_action:
        passed, reason = False, f"Ожидался '{step.expected_action}', получен '{final.action}'"
    elif step.expected_tool and first_tool != step.expected_tool:
        passed, reason = False, f"Ожидался '{step.expected_tool}', вызван '{first_tool}'"
    elif step.expected_tool and tool_success is False and not step.allow_tool_fail:
        passed, reason = False, f"Инструмент '{first_tool}' завершился с ошибкой"

    return TurnResult(
        final_action=final.action if final else "none",
        final_content=final.content if final else None,
        tool_name=first_tool,
        tool_success=tool_success,
        llm_calls=llm_calls,
        elapsed_ms=elapsed,
        passed=passed,
        fail_reason=reason,
    )


# =======================================================================
#  RUNNER
# =======================================================================

@dataclass
class SimRecord:
    scenario: str
    step: int
    passed: bool
    tool: str
    action: str
    llm_calls: int
    ms: float
    fail_reason: str = ""
    allow_tool_fail: bool = False


def run_simulation(
    scenarios: list[Scenario],
    brain: Brain,
    executor: ToolExecutor,
    llm: ScriptedLLM,
    memory: SimpleMemory,
    console: Console,
    deep: bool,
) -> list[SimRecord]:
    records: list[SimRecord] = []
    cur_part = 0
    turn_base = 1

    for idx, sc in enumerate(scenarios, 1):
        if sc.part != cur_part:
            cur_part = sc.part
            print_part_banner(console, cur_part)

        print_scenario_header(console, idx, len(scenarios), sc)

        for s_idx, step in enumerate(sc.steps):
            llm.enqueue(*step.llm_responses)

            turn = run_turn(
                brain=brain, executor=executor, llm=llm, memory=memory,
                user_input=step.user_input, step=step,
                console=console, deep=deep, turn_base=turn_base,
            )
            turn_base += turn.llm_calls

            print_step_result(
                console, turn.passed, turn.fail_reason,
                turn.elapsed_ms, turn.llm_calls, turn.tool_name or "—",
            )
            print_memory_state(console, memory, SESSION)

            records.append(SimRecord(
                scenario=sc.name, step=s_idx + 1,
                passed=turn.passed, tool=turn.tool_name or "—",
                action=turn.final_action, llm_calls=turn.llm_calls,
                ms=turn.elapsed_ms, fail_reason=turn.fail_reason,
                allow_tool_fail=step.allow_tool_fail,
            ))

    return records


# =======================================================================
#  ИТОГОВЫЙ ОТЧЁТ
# =======================================================================

def print_summary(console: Console, records: list[SimRecord]) -> None:
    console.print()
    console.rule("[bold cyan]ИТОГИ — ВСЕ ИНСТРУМЕНТЫ[/]")
    console.print()

    t = Table(title="Результаты по шагам", box=box.ROUNDED,
              show_lines=True, header_style="bold")
    t.add_column("Сценарий",   style="cyan", max_width=20)
    t.add_column("Ш", justify="center", width=3)
    t.add_column("Action",     width=10)
    t.add_column("Инструмент", style="yellow", max_width=18)
    t.add_column("LLM", justify="center", width=5)
    t.add_column("ms",  justify="right",  width=7)
    t.add_column("Статус", justify="center", width=10)

    for r in records:
        if r.passed:
            status = "[green]PASS[/]"
        elif r.allow_tool_fail:
            status = "[yellow]PASS*[/]"
        else:
            status = "[red]FAIL[/]"
        t.add_row(r.scenario, str(r.step), r.action, r.tool,
                  str(r.llm_calls), f"{r.ms:.0f}", status)

    console.print(t)
    console.print("[dim]  * PASS* = ошибка инструмента ожидалась (error_path тест)[/]")
    console.print()

    total   = len(records)
    ok      = sum(1 for r in records if r.passed or r.allow_tool_fail)
    failed  = total - ok
    tot_ms  = sum(r.ms for r in records)
    tot_llm = sum(r.llm_calls for r in records)
    tot_tc  = sum(1 for r in records if r.tool != "—")

    s = Table(box=box.SIMPLE, show_header=False, padding=(0, 2))
    s.add_column("", style="dim")
    s.add_column("")
    s.add_row("Шагов всего",   str(total))
    s.add_row("Прошло",        f"[green]{ok}[/] / {total}  ({ok/total*100:.0f}%)")
    s.add_row("Провалилось",   f"[red]{failed}[/]" if failed else "[green]0[/]")
    s.add_row("Tool-вызовов",  f"[yellow]{tot_tc}[/]")
    s.add_row("LLM-вызовов",   str(tot_llm))
    s.add_row("Время всего",   f"{tot_ms:.0f}ms")
    s.add_row("Среднее/шаг",   f"{tot_ms/total:.0f}ms" if total else "—")
    console.print(s)

    # Статистика по инструментам
    tool_st: dict[str, list[bool]] = {}
    for r in records:
        if r.tool != "—":
            tool_st.setdefault(r.tool, []).append(r.passed or r.allow_tool_fail)
    if tool_st:
        t2 = Table(title="По инструментам", box=box.SIMPLE_HEAVY, header_style="bold dim")
        t2.add_column("Инструмент", style="cyan")
        t2.add_column("Вызовов", justify="center")
        t2.add_column("OK", justify="center")
        for tn, res in sorted(tool_st.items()):
            ok_n = sum(res)
            col = "green" if ok_n == len(res) else "red"
            t2.add_row(tn, str(len(res)), f"[{col}]{ok_n}[/]")
        console.print()
        console.print(t2)

    console.print()
    verdict = (
        "[bold green]ВСЕ ПРОШЛИ — Brain корректно работает со всеми инструментами![/]"
        if not failed else
        f"[bold red]{failed} провалилось — проверьте таблицу[/]"
    )
    console.print(Panel(verdict, border_style="green" if not failed else "red"))
    console.print()


# =======================================================================
#  ENTRY POINT
# =======================================================================

def main() -> int:
    p = argparse.ArgumentParser(description="Полная симуляция Brain + все инструменты")
    p.add_argument("--part", type=int, choices=[1, 2, 3],
                   help="Только одна часть (1=базовые, 2=ридеры, 3=спец)")
    p.add_argument("--tool", help="Только сценарии с этим инструментом")
    p.add_argument("--no-deep", action="store_true",
                   help="Быстрый вывод без пошагового Brain")
    p.add_argument("--no-color", action="store_true")
    p.add_argument("--verbose",  action="store_true")
    args = p.parse_args()

    console = Console(highlight=False, markup=True, no_color=args.no_color)
    if not args.verbose:
        logging.disable(logging.CRITICAL)

    # Создать тестовые файлы
    tmp = Path(tempfile.mkdtemp(prefix="brain_sim_"))
    factory = FileFactory(tmp)
    factory.build_all()
    for fmt, err in factory.errors.items():
        console.print(f"  [yellow]Не удалось создать {fmt}: {err}[/]")

    # Построить и отфильтровать сценарии
    scenarios = build_scenarios(factory.files)
    if args.part:
        scenarios = [s for s in scenarios if s.part == args.part]
    elif args.tool:
        tl = args.tool.lower()
        scenarios = [s for s in scenarios if any(tl in tag for tag in s.tags)]

    if not scenarios:
        console.print("[red]Нет сценариев для запуска.[/]")
        return 1

    # Настройка Brain + Tools
    llm    = ScriptedLLM()
    memory = SimpleMemory()
    brain  = Brain(llm=llm, memory=memory)
    registry, executor = setup_registry(memory)

    print_header(console, len(registry.names()), len(scenarios))
    print_tool_list(console, registry)

    records = run_simulation(
        scenarios=scenarios, brain=brain, executor=executor,
        llm=llm, memory=memory, console=console, deep=not args.no_deep,
    )

    print_summary(console, records)

    failed = sum(1 for r in records if not r.passed and not r.allow_tool_fail)
    return 0 if not failed else 1


if __name__ == "__main__":
    sys.exit(main())
