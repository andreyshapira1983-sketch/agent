# Document Parser — подсистема Perception Layer (Слой 1)
# Архитектура автономного AI-агента
# Чтение PDF, DOCX, TXT, MD, CSV, JSON, HTML-файлов.
# pylint: disable=broad-except

from __future__ import annotations

import os
import re
import json


class ParsedDocument:
    """Результат разбора документа."""

    def __init__(self, path: str, text: str, metadata: dict | None = None,
                 pages: int | None = None, doc_type: str | None = None):
        self.path = path
        self.text = text
        self.metadata = metadata or {}
        self.pages = pages
        self.doc_type = doc_type or os.path.splitext(path)[1].lower().lstrip('.')
        self.word_count = len(text.split())
        self.char_count = len(text)

    def to_dict(self) -> dict:
        return {
            'path': self.path,
            'doc_type': self.doc_type,
            'pages': self.pages,
            'word_count': self.word_count,
            'char_count': self.char_count,
            'metadata': self.metadata,
            'text_preview': self.text[:300],
        }


class DocumentParser:
    """
    Document Parser — подсистема PerceptionLayer.

    Поддерживает:
        .pdf        — через pypdf (pip install pypdf)
        .docx/.doc  — через python-docx (pip install python-docx)
        .xlsx/.xls  — через openpyxl / xlrd (pip install openpyxl)
        .pptx/.ppt  — через python-pptx (pip install python-pptx)
        .txt/.md    — встроено (markdown stripped)
        .csv        — встроено
        .json       — встроено
        .html/.htm  — через BeautifulSoup или regex
        .py/.js/.ts и др. — как текст
        .odt/.ods/.odp — через odfpy (pip install odfpy)
        .rtf        — через striprtf (pip install striprtf)
        .xml/.svg   — встроенный xml.etree
        .eml        — встроенный email
        .msg        — через extract-msg (pip install extract-msg)
        .db/.sqlite — встроенный sqlite3
        .zip/.tar/.gz/.7z/.rar — архивы (список + текстовые файлы)
        .mp3/.wav/.ogg/.flac/.aac — метаданные через mutagen (pip install mutagen)
        .mp4/.avi/.mkv/.mov/.webm — метаданные через mutagen или ffprobe
        .exe/.dll/.bin — строки/метаданные из бинарника

    При отсутствии зависимости — пробует fallback (pdfminer, PyPDF2, docx2txt).
    """

    SUPPORTED = {
        '.pdf', '.docx', '.doc', '.txt', '.md',
        '.csv', '.json', '.html', '.htm', '.py',
        '.js', '.ts', '.yaml', '.yml', '.toml', '.ini', '.cfg',
        # Таблицы
        '.xlsx', '.xls',
        # Презентации
        '.pptx', '.ppt',
        # LibreOffice
        '.odt', '.ods', '.odp',
        # RTF
        '.rtf',
        # XML / SVG
        '.xml', '.svg',
        # Email
        '.eml', '.msg',
        # Базы данных
        '.db', '.sqlite', '.sqlite3',
        # Архивы
        '.zip', '.tar', '.gz', '.tgz', '.bz2', '.7z', '.rar',
        # Аудио
        '.mp3', '.wav', '.ogg', '.flac', '.aac', '.m4a', '.wma',
        # Видео
        '.mp4', '.avi', '.mkv', '.mov', '.webm', '.flv', '.wmv',
        # Бинарные
        '.exe', '.dll', '.so', '.bin',
    }

    def __init__(self, monitoring=None):
        self.monitoring = monitoring

    # ── Основной метод ────────────────────────────────────────────────────────

    def parse(self, path: str) -> ParsedDocument | None:
        """
        Разбирает документ по пути.

        Returns:
            ParsedDocument или None если файл не поддерживается / не читается.
        """
        if not os.path.exists(path):
            self._log(f'Файл не найден: {path}', level='error')
            return None

        ext = os.path.splitext(path)[1].lower()

        try:
            if ext == '.pdf':
                return self._parse_pdf(path)
            elif ext in ('.docx', '.doc'):
                return self._parse_docx(path)
            elif ext in ('.html', '.htm'):
                return self._parse_html(path)
            elif ext == '.csv':
                return self._parse_csv(path)
            elif ext == '.json':
                return self._parse_json(path)
            elif ext in ('.xlsx', '.xls'):
                return self._parse_xlsx(path)
            elif ext in ('.pptx', '.ppt'):
                return self._parse_pptx(path)
            elif ext in ('.odt', '.ods', '.odp'):
                return self._parse_odf(path)
            elif ext == '.rtf':
                return self._parse_rtf(path)
            elif ext in ('.xml', '.svg'):
                return self._parse_xml(path)
            elif ext == '.eml':
                return self._parse_eml(path)
            elif ext == '.msg':
                return self._parse_msg(path)
            elif ext in ('.db', '.sqlite', '.sqlite3'):
                return self._parse_sqlite(path)
            elif ext in ('.zip',):
                return self._parse_zip(path)
            elif ext in ('.tar', '.gz', '.tgz', '.bz2'):
                return self._parse_tar(path)
            elif ext == '.7z':
                return self._parse_7z(path)
            elif ext == '.rar':
                return self._parse_rar(path)
            elif ext in ('.mp3', '.wav', '.ogg', '.flac', '.aac', '.m4a', '.wma'):
                return self._parse_audio(path)
            elif ext in ('.mp4', '.avi', '.mkv', '.mov', '.webm', '.flv', '.wmv'):
                return self._parse_video(path)
            elif ext in ('.exe', '.dll', '.so', '.bin'):
                return self._parse_binary(path)
            elif ext in self.SUPPORTED or ext in ('.txt', '.md', '.py',
                                                   '.js', '.ts', '.yaml',
                                                   '.yml', '.toml', '.ini', '.cfg'):
                return self._parse_text(path)
            else:
                self._log(f'Неподдерживаемый тип: {ext}', level='warning')
                return None
        except Exception as e:
            self._log(f'Ошибка парсинга {path}: {e}', level='error')
            return None

    def parse_directory(self, directory: str,
                        recursive: bool = True,
                        extensions: set | None = None) -> list[ParsedDocument]:
        """Разбирает все документы в директории."""
        exts = extensions or self.SUPPORTED
        docs = []
        walker = os.walk(directory) if recursive else [(directory, [], os.listdir(directory))]
        for root, _, files in walker:
            for fname in files:
                if os.path.splitext(fname)[1].lower() in exts:
                    doc = self.parse(os.path.join(root, fname))
                    if doc:
                        docs.append(doc)
        self._log(f'Разобрано {len(docs)} документов из {directory}')
        return docs

    # ── PDF ───────────────────────────────────────────────────────────────────

    def _parse_pdf(self, path: str) -> ParsedDocument:
        text = ''
        pages = 0
        metadata = {}

        # Попытка 1: pypdf (современный)
        try:
            import pypdf
            reader = pypdf.PdfReader(path)
            pages = len(reader.pages)
            info = reader.metadata or {}
            metadata = {
                'title': str(info.get('/Title', '')),
                'author': str(info.get('/Author', '')),
                'subject': str(info.get('/Subject', '')),
            }
            parts = []
            for page in reader.pages:
                parts.append(page.extract_text() or '')
            text = '\n'.join(parts)
            self._log(f'PDF (pypdf): {pages} страниц, {len(text)} символов')
            return ParsedDocument(path, text, metadata, pages, 'pdf')
        except ImportError:
            pass

        # Попытка 2: PyPDF2 (старый)
        try:
            import PyPDF2
            reader = PyPDF2.PdfReader(path)
            pages = len(reader.pages)
            parts = [reader.pages[i].extract_text() or ''
                     for i in range(pages)]
            text = '\n'.join(parts)
            return ParsedDocument(path, text, {}, pages, 'pdf')
        except ImportError:
            pass

        # Попытка 3: pdfplumber
        try:
            import pdfplumber
            with pdfplumber.open(path) as pdf:
                pages = len(pdf.pages)
                text = '\n'.join(p.extract_text() or '' for p in pdf.pages)
            return ParsedDocument(path, text, {}, pages, 'pdf')
        except ImportError:
            pass

        raise ImportError(
            'Для чтения PDF установи: pip install pypdf  '
            '(или PyPDF2, pdfplumber)'
        )

    # ── DOCX ──────────────────────────────────────────────────────────────────

    def _parse_docx(self, path: str) -> ParsedDocument:
        # Попытка 1: python-docx
        try:
            import docx
            doc = docx.Document(path)
            paragraphs = [p.text for p in doc.paragraphs]
            # Таблицы
            for table in doc.tables:
                for row in table.rows:
                    paragraphs.append('\t'.join(c.text for c in row.cells))
            text = '\n'.join(p for p in paragraphs if p.strip())
            props = doc.core_properties
            metadata = {
                'title': props.title or '',
                'author': props.author or '',
                'subject': props.subject or '',
            }
            return ParsedDocument(path, text, metadata, None, 'docx')
        except ImportError:
            pass

        # Попытка 2: docx2txt
        try:
            import docx2txt
            text = docx2txt.process(path) or ''
            return ParsedDocument(path, text, {}, None, 'docx')
        except ImportError:
            pass

        raise ImportError('Для чтения DOCX установи: pip install python-docx')

    # ── HTML ──────────────────────────────────────────────────────────────────

    def _parse_html(self, path: str) -> ParsedDocument:
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            raw = f.read()

        # Попытка BS4
        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(raw, 'lxml')
            for tag in soup(['script', 'style', 'nav', 'footer']):
                tag.decompose()
            title = soup.title.string if soup.title else ''
            text = soup.get_text(separator='\n', strip=True)
            return ParsedDocument(path, text, {'title': title}, None, 'html')
        except ImportError:
            pass

        # Fallback: regex
        text = re.sub(r'<[^>]+>', ' ', raw)
        text = re.sub(r'\s+', ' ', text).strip()
        return ParsedDocument(path, text, {}, None, 'html')

    # ── CSV ───────────────────────────────────────────────────────────────────

    def _parse_csv(self, path: str) -> ParsedDocument:
        import csv
        rows = []
        try:
            with open(path, 'r', encoding='utf-8', errors='replace', newline='') as f:
                reader = csv.reader(f)
                for row in reader:
                    rows.append('\t'.join(row))
        except Exception as e:
            return ParsedDocument(path, f'Ошибка CSV: {e}', {}, None, 'csv')
        text = '\n'.join(rows)
        return ParsedDocument(path, text,
                               {'rows': len(rows)}, None, 'csv')

    # ── JSON ──────────────────────────────────────────────────────────────────

    def _parse_json(self, path: str) -> ParsedDocument:
        try:
            with open(path, 'r', encoding='utf-8', errors='replace') as f:
                data = json.load(f)
            text = json.dumps(data, ensure_ascii=False, indent=2)
        except Exception as e:
            text = f'Ошибка JSON: {e}'
        return ParsedDocument(path, text, {}, None, 'json')

    # ── Текстовые форматы ─────────────────────────────────────────────────────

    def _parse_text(self, path: str) -> ParsedDocument:
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            text = f.read()
        ext = os.path.splitext(path)[1].lower()
        # Markdown — убираем синтаксис для clean-text
        if ext == '.md':
            text_clean = re.sub(r'#{1,6}\s', '', text)
            text_clean = re.sub(r'\*{1,2}([^*]+)\*{1,2}', r'\1', text_clean)
            text_clean = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text_clean)
            return ParsedDocument(path, text_clean, {}, None, 'md')
        return ParsedDocument(path, text, {}, None, ext.lstrip('.'))

    # ── XLSX / XLS ────────────────────────────────────────────────────────────

    def _parse_xlsx(self, path: str) -> ParsedDocument:
        ext = os.path.splitext(path)[1].lower()
        # Попытка 1: openpyxl (xlsx)
        try:
            import openpyxl
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                parts.append(f'=== Лист: {sheet_name} ===')
                for row in ws.iter_rows(values_only=True):
                    row_strs = [str(v) if v is not None else '' for v in row]
                    parts.append('\t'.join(row_strs))
            wb.close()
            text = '\n'.join(parts)
            metadata = {'sheets': len(wb.sheetnames), 'sheet_names': list(wb.sheetnames)}
            self._log(f'XLSX (openpyxl): {len(wb.sheetnames)} листов, {len(text)} символов')
            return ParsedDocument(path, text, metadata, None, ext.lstrip('.'))
        except ImportError:
            pass
        # Попытка 2: xlrd (xls / старый xlsx)
        try:
            import xlrd
            wb = xlrd.open_workbook(path)
            parts = []
            for sheet in wb.sheets():
                parts.append(f'=== Лист: {sheet.name} ===')
                for row_idx in range(sheet.nrows):
                    row = sheet.row_values(row_idx)
                    parts.append('\t'.join(str(v) for v in row))
            text = '\n'.join(parts)
            return ParsedDocument(path, text, {'sheets': wb.nsheets}, None, ext.lstrip('.'))
        except ImportError:
            pass
        raise ImportError('Для чтения XLSX установи: pip install openpyxl')

    # ── PowerPoint ────────────────────────────────────────────────────────────

    def _parse_pptx(self, path: str) -> ParsedDocument:
        ext = os.path.splitext(path)[1].lower()
        try:
            from pptx import Presentation
            from pptx.util import Pt
            prs = Presentation(path)
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                parts.append(f'=== Слайд {i} ===')
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        tf = getattr(shape, 'text_frame', None)
                        if tf:
                            t = tf.text.strip()
                            if t:
                                parts.append(t)
                    # Таблицы внутри слайда
                    if shape.has_table:
                        tbl = getattr(shape, 'table', None)
                        if tbl:
                            for row in tbl.rows:
                                row_text = '\t'.join(
                                    cell.text.strip() for cell in row.cells
                                )
                                if row_text.strip():
                                    parts.append(row_text)
            text = '\n'.join(parts)
            metadata = {
                'slides': len(prs.slides),
                'title': (prs.core_properties.title or '') if prs.core_properties else '',
                'author': (prs.core_properties.author or '') if prs.core_properties else '',
            }
            self._log(f'PPTX (python-pptx): {len(prs.slides)} слайдов, {len(text)} символов')
            return ParsedDocument(path, text, metadata, None, ext.lstrip('.'))
        except ImportError:
            pass
        raise ImportError('Для чтения PowerPoint установи: pip install python-pptx')

    # ── LibreOffice ODF (ODT / ODS / ODP) ──────────────────────────────────────

    def _parse_odf(self, path: str) -> ParsedDocument:
        ext = os.path.splitext(path)[1].lower()
        try:
            from odf.opendocument import load as odf_load
            from odf.text import P, H
            from odf.table import TableCell
            from odf import teletype
            doc = odf_load(path)
            parts = []
            # Текст из ODT/ODP: абзацы + заголовки
            doc_text = getattr(doc, 'text', None)
            for el in (doc_text.childNodes if doc_text is not None else []):
                t = teletype.extractText(el)
                if t.strip():
                    parts.append(t)
            # ODS: таблицы
            doc_sheet = getattr(doc, 'spreadsheet', None)
            if doc_sheet:
                for sheet in doc_sheet.childNodes:
                    name = sheet.getAttribute('name') or ''
                    if name:
                        parts.append(f'=== Лист: {name} ===')
                    for row in sheet.childNodes:
                        cells = [teletype.extractText(c) for c in row.childNodes]
                        line = '\t'.join(cells)
                        if line.strip():
                            parts.append(line)
            text = '\n'.join(parts)
            meta = getattr(doc, 'meta', None)
            metadata = {
                'title': str(getattr(meta, 'title', '') or '') if meta else '',
                'creator': str(getattr(meta, 'creator', '') or '') if meta else '',
            }
            self._log(f'ODF ({ext}): {len(text)} символов')
            return ParsedDocument(path, text, metadata, None, ext.lstrip('.'))
        except ImportError:
            pass
        raise ImportError('Для чтения ODT/ODS/ODP установи: pip install odfpy')

    # ── RTF ──────────────────────────────────────────────────────────────────────

    def _parse_rtf(self, path: str) -> ParsedDocument:
        try:
            from striprtf.striprtf import rtf_to_text
            with open(path, 'r', encoding='utf-8', errors='replace') as f:
                raw = f.read()
            text = rtf_to_text(raw)
            self._log(f'RTF: {len(text)} символов')
            return ParsedDocument(path, text, {}, None, 'rtf')
        except ImportError:
            pass
        # Fallback: убрать RTF-теги regex
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            raw = f.read()
        text = re.sub(r'\\[a-z]+\d*\s?|[\\{}]', ' ', raw)
        text = re.sub(r'\s+', ' ', text).strip()
        return ParsedDocument(path, text, {}, None, 'rtf')

    # ── XML / SVG ──────────────────────────────────────────────────────────────

    def _parse_xml(self, path: str) -> ParsedDocument:
        import xml.etree.ElementTree as ET
        ext = os.path.splitext(path)[1].lower()
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            raw = f.read()
        try:
            root = ET.fromstring(raw)
            # Собираем весь текст из всех нодов
            texts = [t.strip() for t in root.itertext() if t.strip()]
            text = '\n'.join(texts)
            tag = root.tag.split('}')[-1] if '}' in root.tag else root.tag
            metadata = {'root_tag': tag, 'format': ext.lstrip('.')}
        except ET.ParseError:
            # Fallback: чистый текст без тегов
            text = re.sub(r'<[^>]+>', ' ', raw)
            text = re.sub(r'\s+', ' ', text).strip()
            metadata = {'format': ext.lstrip('.')}
        self._log(f'XML/SVG: {len(text)} символов')
        return ParsedDocument(path, text, metadata, None, ext.lstrip('.'))

    # ── Email EML ─────────────────────────────────────────────────────────────

    def _parse_eml(self, path: str) -> ParsedDocument:
        import email as _email
        from email import policy
        with open(path, 'rb') as f:
            msg = _email.message_from_binary_file(f, policy=policy.default)
        parts = []
        metadata: dict = {
            'from':    str(msg.get('From', '')),
            'to':      str(msg.get('To', '')),
            'subject': str(msg.get('Subject', '')),
            'date':    str(msg.get('Date', '')),
        }
        parts.append(f'От: {metadata["from"]}')
        parts.append(f'Кому: {metadata["to"]}')
        parts.append(f'Тема: {metadata["subject"]}')
        parts.append(f'Дата: {metadata["date"]}')
        parts.append('')
        if msg.is_multipart():
            for part in msg.walk():
                ct = part.get_content_type()
                if ct == 'text/plain':
                    body = part.get_content()
                    parts.append(str(body))
        else:
            parts.append(str(msg.get_content()))
        # Вложения
        attachments: list[str] = [
            p.get_filename() for p in msg.walk()
            if p.get_content_disposition() == 'attachment' and p.get_filename() is not None
        ]  # type: ignore[misc]
        if attachments:
            parts.append(f'\nВложения: {', '.join(attachments)}')
            metadata['attachments'] = attachments
        text = '\n'.join(parts)
        self._log(f'EML: {metadata["subject"]}')
        return ParsedDocument(path, text, metadata, None, 'eml')

    # ── Outlook MSG ──────────────────────────────────────────────────────────

    def _parse_msg(self, path: str) -> ParsedDocument:
        try:
            import extract_msg
            with extract_msg.openMsg(path) as msg:
                sender  = getattr(msg, 'sender', '') or ''
                to      = getattr(msg, 'to', '') or ''
                subject = getattr(msg, 'subject', '') or ''
                date    = getattr(msg, 'date', '') or ''
                body    = getattr(msg, 'body', '') or ''
                parts = [
                    f'От: {sender}',
                    f'Кому: {to}',
                    f'Тема: {subject}',
                    f'Дата: {date}',
                    '',
                    body,
                ]
                raw_att = getattr(msg, 'attachments', []) or []
                attachments = [
                    getattr(a, 'longFilename', None) or getattr(a, 'shortFilename', '')
                    for a in raw_att
                ]
                if attachments:
                    parts.append(f'\nВложения: {", ".join(str(a) for a in attachments)}')
                metadata = {
                    'from': str(sender),
                    'to': str(to),
                    'subject': str(subject),
                    'attachments': [str(a) for a in attachments],
                }
                text = '\n'.join(parts)
            self._log(f'MSG (Outlook): {metadata["subject"]}')
            return ParsedDocument(path, text, metadata, None, 'msg')
        except ImportError:
            pass
        raise ImportError('Для чтения .msg установи: pip install extract-msg')

    # ── SQLite ─────────────────────────────────────────────────────────────────

    def _parse_sqlite(self, path: str) -> ParsedDocument:
        import sqlite3
        parts = [f'SQLite база: {os.path.basename(path)}']
        metadata: dict = {}
        try:
            conn = sqlite3.connect(f'file:{path}?mode=ro', uri=True)
            cur = conn.cursor()
            cur.execute("SELECT name FROM sqlite_master WHERE type='table' ORDER BY name")
            tables = [row[0] for row in cur.fetchall()]
            parts.append(f'Таблицы ({len(tables)}): {", ".join(tables)}')
            metadata['tables'] = tables
            for tbl in tables[:10]:  # первые 10 таблиц
                # Санитизация имени таблицы (SQL injection prevention)
                safe_tbl = ''.join(c for c in tbl if c.isalnum() or c == '_')
                if not safe_tbl or safe_tbl != tbl:
                    continue
                try:
                    cur.execute(f'PRAGMA table_info("{safe_tbl}")')
                    cols = [row[1] for row in cur.fetchall()]
                    cur.execute(f'SELECT COUNT(*) FROM "{safe_tbl}"')
                    cnt = cur.fetchone()[0]
                    parts.append(f'\nТаблица "{safe_tbl}" ({cnt} строк): {", ".join(cols)}')
                    # Первые 5 строк
                    cur.execute(f'SELECT * FROM "{safe_tbl}" LIMIT 5')
                    for row in cur.fetchall():
                        parts.append('\t'.join(str(v) for v in row))
                except Exception:
                    pass
            conn.close()
        except Exception as e:
            parts.append(f'Ошибка открытия: {e}')
        text = '\n'.join(parts)
        self._log(f'SQLite: {len(metadata.get("tables", []))} таблиц')
        return ParsedDocument(path, text, metadata, None, 'sqlite')

    # ── TAR / GZ / BZ2 ─────────────────────────────────────────────────────

    def _parse_tar(self, path: str) -> ParsedDocument:
        import tarfile
        _TEXT_EXTS = {
            '.txt', '.md', '.py', '.js', '.ts', '.json', '.yaml', '.yml',
            '.csv', '.html', '.htm', '.xml', '.ini', '.cfg', '.toml', '.log',
            '.sh', '.bat', '.rb', '.go', '.java', '.c', '.cpp', '.h',
        }
        parts = []
        names = []
        try:
            with tarfile.open(path, 'r:*') as tf:
                members = tf.getmembers()
                names = [m.name for m in members]
                parts.append(f'TAR-архив содержит {len(names)} файлов:')
                for m in members:
                    parts.append(f'  {m.name}  ({m.size:,} байт)')
                parts.append('')
                read_count = 0
                for m in members:
                    if read_count >= 20:
                        parts.append(f'[...ещё {len(members) - read_count} не прочитано]')
                        break
                    if not m.isfile():
                        continue
                    ext_inner = os.path.splitext(m.name)[1].lower()
                    if ext_inner in _TEXT_EXTS and m.size > 0:
                        try:
                            f = tf.extractfile(m)
                            if f:
                                content = f.read(51_200).decode('utf-8', errors='replace')
                                parts.append(f'--- {m.name} ---')
                                parts.append(content)
                                read_count += 1
                        except Exception:
                            pass
        except Exception as e:
            return ParsedDocument(path, f'Ошибка TAR: {e}', {}, None, 'tar')
        text = '\n'.join(parts)
        self._log(f'TAR: {len(names)} файлов')
        return ParsedDocument(path, text, {'files': len(names)}, None, 'tar')

    # ── 7Z ────────────────────────────────────────────────────────────────────────

    def _parse_7z(self, path: str) -> ParsedDocument:
        _TEXT_EXTS = {
            '.txt', '.md', '.py', '.js', '.json', '.yaml', '.yml',
            '.csv', '.html', '.xml', '.ini', '.cfg', '.toml', '.log',
        }
        try:
            import py7zr
            parts = []
            with py7zr.SevenZipFile(path, mode='r') as z:
                names = z.getnames()
                parts.append(f'7Z-архив содержит {len(names)} файлов:')
                for name in names:
                    parts.append(f'  {name}')
                parts.append('')
                text_files = [n for n in names
                              if os.path.splitext(n)[1].lower() in _TEXT_EXTS]
                if text_files:
                    read_data: dict = z.read(text_files[:20]) or {}
                    for name, bio in read_data.items():
                        content = bio.read(51_200).decode('utf-8', errors='replace')
                        parts.append(f'--- {name} ---')
                        parts.append(content)
            text = '\n'.join(parts)
            self._log(f'7Z: {len(names)} файлов')
            return ParsedDocument(path, text, {'files': len(names)}, None, '7z')
        except ImportError:
            pass
        raise ImportError('Для чтения 7Z установи: pip install py7zr')

    # ── RAR ────────────────────────────────────────────────────────────────────────

    def _parse_rar(self, path: str) -> ParsedDocument:
        _TEXT_EXTS = {
            '.txt', '.md', '.py', '.js', '.json', '.yaml', '.yml',
            '.csv', '.html', '.xml', '.ini', '.cfg', '.toml', '.log',
        }
        try:
            import rarfile
            parts = []
            with rarfile.RarFile(path) as rf:
                names = rf.namelist()
                parts.append(f'RAR-архив содержит {len(names)} файлов:')
                for info in rf.infolist():
                    parts.append(f'  {info.filename}  ({info.file_size:,} байт)')
                parts.append('')
                read_count = 0
                for name in names:
                    if read_count >= 20:
                        break
                    if os.path.splitext(name)[1].lower() in _TEXT_EXTS:
                        try:
                            content = rf.read(name)[:51_200].decode('utf-8', errors='replace')
                            parts.append(f'--- {name} ---')
                            parts.append(content)
                            read_count += 1
                        except Exception:
                            pass
            text = '\n'.join(parts)
            self._log(f'RAR: {len(names)} файлов')
            return ParsedDocument(path, text, {'files': len(names)}, None, 'rar')
        except ImportError:
            pass
        raise ImportError('Для чтения RAR установи: pip install rarfile  (требуется unrar в PATH)')

    # ── ZIP ───────────────────────────────────────────────────────────────────

    def _parse_zip(self, path: str) -> ParsedDocument:
        import zipfile
        _TEXT_EXTS = {
            '.txt', '.md', '.py', '.js', '.ts', '.json', '.yaml', '.yml',
            '.csv', '.html', '.htm', '.xml', '.ini', '.cfg', '.toml', '.log',
            '.sh', '.bat', '.ps1', '.rb', '.go', '.java', '.c', '.cpp', '.h',
        }
        parts = []
        names = []
        try:
            with zipfile.ZipFile(path, 'r') as zf:
                names = zf.namelist()
                parts.append(f'ZIP-архив содержит {len(names)} файлов:')
                for name in names:
                    try:
                        info = zf.getinfo(name)
                        parts.append(f'  {name}  ({info.file_size:,} байт)')
                    except Exception:
                        parts.append(f'  {name}')
                parts.append('')
                # Читаем текстовые файлы внутри (до 50 КБ каждый)
                read_count = 0
                for name in names:
                    if read_count >= 20:
                        parts.append(f'[...ещё {len(names) - read_count} файлов не прочитано]')
                        break
                    ext_inner = os.path.splitext(name)[1].lower()
                    if ext_inner in _TEXT_EXTS:
                        try:
                            with zf.open(name) as f:
                                content = f.read(51_200).decode('utf-8', errors='replace')
                            parts.append(f'--- {name} ---')
                            parts.append(content)
                            read_count += 1
                        except Exception:
                            pass
        except zipfile.BadZipFile as e:
            return ParsedDocument(path, f'Ошибка ZIP: {e}', {}, None, 'zip')
        text = '\n'.join(parts)
        self._log(f'ZIP: {len(names)} файлов')
        return ParsedDocument(path, text, {'files': len(names)}, None, 'zip')

    # ── Аудио ─────────────────────────────────────────────────────────────────

    def _parse_audio(self, path: str) -> ParsedDocument:
        ext = os.path.splitext(path)[1].lower()
        fname = os.path.basename(path)
        parts = [f'Аудиофайл: {fname}']
        metadata: dict = {}
        try:
            from mutagen import File as _MutagenFile  # type: ignore[attr-defined]
            audio = _MutagenFile(path)
            if audio is not None:
                info = getattr(audio, 'info', None)
                if info:
                    duration = getattr(info, 'length', None)
                    if duration:
                        m, s = divmod(int(duration), 60)
                        parts.append(f'Длительность: {duration:.1f} сек ({m}:{s:02d} мин)')
                        metadata['duration_sec'] = round(duration, 1)
                    bitrate = getattr(info, 'bitrate', None)
                    if bitrate:
                        parts.append(f'Битрейт: {bitrate} кбит/с')
                        metadata['bitrate'] = bitrate
                    sample_rate = getattr(info, 'sample_rate', None)
                    if sample_rate:
                        parts.append(f'Частота: {sample_rate} Гц')
                # Теги
                for key in ('title', 'artist', 'album', 'date', 'comment',
                            'TIT2', 'TPE1', 'TALB', 'TDRC'):
                    val = audio.get(key)
                    if val:
                        v = val[0] if hasattr(val, '__getitem__') else str(val)
                        parts.append(f'{key}: {v}')
        except ImportError:
            parts.append('(mutagen не установлен — pip install mutagen)')
        except Exception as e:
            parts.append(f'(ошибка метаданных: {e})')
        size = os.path.getsize(path)
        parts.append(f'Размер файла: {size / 1024:.1f} КБ')
        metadata['size_bytes'] = size
        text = '\n'.join(parts)
        self._log(f'Аудио ({ext}): {fname}')
        return ParsedDocument(path, text, metadata, None, ext.lstrip('.'))

    # ── Видео ─────────────────────────────────────────────────────────────────

    def _parse_video(self, path: str) -> ParsedDocument:
        ext = os.path.splitext(path)[1].lower()
        fname = os.path.basename(path)
        parts = [f'Видеофайл: {fname}']
        metadata: dict = {}
        parsed = False
        try:
            from mutagen import File as _MutagenFile  # type: ignore[attr-defined]
            video = _MutagenFile(path)
            if video is not None:
                info = getattr(video, 'info', None)
                if info:
                    duration = getattr(info, 'length', None)
                    if duration:
                        m, s = divmod(int(duration), 60)
                        parts.append(f'Длительность: {duration:.1f} сек ({m}:{s:02d} мин)')
                        metadata['duration_sec'] = round(duration, 1)
                    parsed = True
        except ImportError:
            pass
        except Exception:
            pass
        if not parsed:
            # Fallback: ffprobe
            try:
                import subprocess
                import json as _json
                result = subprocess.run(
                    ['ffprobe', '-v', 'quiet', '-print_format', 'json',
                     '-show_format', '-show_streams', path],
                    capture_output=True, text=True, timeout=15,
                    check=False,
                )
                if result.returncode == 0:
                    probe = _json.loads(result.stdout)
                    fmt = probe.get('format', {})
                    duration = float(fmt.get('duration', 0) or 0)
                    if duration:
                        m, s = divmod(int(duration), 60)
                        parts.append(f'Длительность: {duration:.1f} сек ({m}:{s:02d} мин)')
                        metadata['duration_sec'] = round(duration, 1)
                    br = fmt.get('bit_rate')
                    if br:
                        parts.append(f'Битрейт: {int(br) // 1000} кбит/с')
                    for stream in probe.get('streams', [])[:3]:
                        ctype = stream.get('codec_type', '')
                        codec = stream.get('codec_name', '')
                        if ctype == 'video':
                            w = stream.get('width', '')
                            h = stream.get('height', '')
                            fps = stream.get('r_frame_rate', '')
                            parts.append(f'Видео: {codec} {w}x{h} @ {fps}')
                        elif ctype == 'audio':
                            parts.append(f'Аудио: {codec}')
            except FileNotFoundError:
                parts.append('(ffprobe не найден — установите ffmpeg для метаданных видео)')
            except Exception as e:
                parts.append(f'(ошибка ffprobe: {e})')
        size = os.path.getsize(path)
        parts.append(f'Размер файла: {size / 1024 / 1024:.2f} МБ')
        metadata['size_bytes'] = size
        text = '\n'.join(parts)
        self._log(f'Видео ({ext}): {fname}')
        return ParsedDocument(path, text, metadata, None, ext.lstrip('.'))

    # ── Бинарные файлы ────────────────────────────────────────────────────────

    def _parse_binary(self, path: str) -> ParsedDocument:
        ext = os.path.splitext(path)[1].lower()
        fname = os.path.basename(path)
        parts = [f'Бинарный файл: {fname}']
        metadata: dict = {}
        size = os.path.getsize(path)
        parts.append(f'Размер: {size / 1024:.1f} КБ')
        metadata['size_bytes'] = size
        # PE-метаданные для .exe / .dll
        if ext in ('.exe', '.dll'):
            try:
                import pefile
                import datetime
                pe = pefile.PE(path)
                fh = getattr(pe, 'FILE_HEADER', None)
                ts = getattr(fh, 'TimeDateStamp', None)
                if ts is not None:
                    dt = datetime.datetime.fromtimestamp(ts)
                    parts.append(f'Дата компиляции: {dt.strftime("%Y-%m-%d %H:%M:%S")}')
                for fi_list in (getattr(pe, 'FileInfo', None) or []):
                    entries = fi_list if isinstance(fi_list, list) else [fi_list]
                    for entry in entries:
                        for st in getattr(entry, 'StringTable', []):
                            for key, value in st.entries.items():
                                k = key.decode('utf-8', errors='replace')
                                v = value.decode('utf-8', errors='replace').strip()
                                if v:
                                    parts.append(f'{k}: {v}')
                pe.close()
            except ImportError:
                pass
            except Exception as e:
                parts.append(f'(ошибка PE: {e})')
        # Извлекаем строки из бинарника (первые 1 МБ)
        try:
            with open(path, 'rb') as f:
                data = f.read(min(size, 1_048_576))
            printable = re.findall(rb'[\x20-\x7E]{5,}', data)
            strings = [s.decode('ascii', errors='replace') for s in printable[:300]]
            if strings:
                parts.append(f'\nСтроки в файле (первые {len(strings)} из {len(printable)}):')  
                parts.extend(strings)
        except Exception as e:
            parts.append(f'(ошибка чтения строк: {e})')
        text = '\n'.join(parts)
        self._log(f'Бинарник ({ext}): {fname}')
        return ParsedDocument(path, text, metadata, None, ext.lstrip('.'))

    # ── Helpers ───────────────────────────────────────────────────────────────

    def _log(self, message: str, level: str = 'info'):
        if self.monitoring:
            getattr(self.monitoring, level, self.monitoring.info)(
                message, source='document_parser'
            )
        else:
            print(f'[DocumentParser] {message}')
