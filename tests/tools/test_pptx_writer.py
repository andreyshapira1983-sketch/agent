"""Tests for tools/builtins/pptx_writer.py — PptxWriterTool."""

from __future__ import annotations

from pathlib import Path

import pytest

from tools.builtins.pptx_writer import PptxWriterTool


# ════════════════════════════════════════════════════════════════════
# Helpers
# ════════════════════════════════════════════════════════════════════

def _read_pptx_titles(path: Path) -> list[str]:
    from pptx import Presentation
    prs = Presentation(str(path))
    titles = []
    for slide in prs.slides:
        if slide.shapes.title is not None:
            titles.append(slide.shapes.title.text)
        else:
            titles.append("")
    return titles


# ════════════════════════════════════════════════════════════════════
# Validation
# ════════════════════════════════════════════════════════════════════

class TestValidation:

    def test_missing_file_path(self):
        result = PptxWriterTool().execute(slides=[{"title": "x"}])
        assert not result.success

    def test_wrong_extension(self, tmp_path):
        result = PptxWriterTool().execute(
            file_path=str(tmp_path / "deck.txt"),
            slides=[{"title": "x"}],
        )
        assert not result.success
        assert ".pptx" in result.error

    def test_empty_slides(self, tmp_path):
        result = PptxWriterTool().execute(
            file_path=str(tmp_path / "deck.pptx"),
            slides=[],
        )
        assert not result.success

    def test_too_many_slides(self, tmp_path):
        slides = [{"title": f"s{i}"} for i in range(150)]
        result = PptxWriterTool().execute(
            file_path=str(tmp_path / "deck.pptx"),
            slides=slides,
        )
        assert not result.success

    def test_refuses_existing_file_without_overwrite(self, tmp_path):
        p = tmp_path / "deck.pptx"
        p.write_bytes(b"existing")
        result = PptxWriterTool().execute(
            file_path=str(p), slides=[{"title": "x"}],
        )
        assert not result.success
        assert "exists" in result.error


# ════════════════════════════════════════════════════════════════════
# Happy path
# ════════════════════════════════════════════════════════════════════

class TestCreate:

    def test_title_only(self, tmp_path):
        out = tmp_path / "deck.pptx"
        result = PptxWriterTool().execute(
            file_path=str(out),
            slides=[{"title": "Cover Slide"}],
        )
        assert result.success
        assert out.exists()
        assert result.output["slides"] == 1
        titles = _read_pptx_titles(out)
        assert titles == ["Cover Slide"]

    def test_title_subtitle(self, tmp_path):
        out = tmp_path / "deck.pptx"
        result = PptxWriterTool().execute(
            file_path=str(out),
            slides=[{"title": "T", "subtitle": "S"}],
        )
        assert result.success
        titles = _read_pptx_titles(out)
        assert titles == ["T"]

    def test_bullets(self, tmp_path):
        out = tmp_path / "deck.pptx"
        result = PptxWriterTool().execute(
            file_path=str(out),
            slides=[{"title": "Plans", "bullets": ["A", "B", "C"]}],
        )
        assert result.success
        from pptx import Presentation
        prs = Presentation(str(out))
        slide = prs.slides[0]
        body_text = slide.placeholders[1].text_frame.text
        assert "A" in body_text and "B" in body_text and "C" in body_text

    def test_body(self, tmp_path):
        out = tmp_path / "deck.pptx"
        result = PptxWriterTool().execute(
            file_path=str(out),
            slides=[{"title": "Body", "body": "Long paragraph of content."}],
        )
        assert result.success

    def test_overwrite(self, tmp_path):
        out = tmp_path / "deck.pptx"
        out.write_bytes(b"old")
        result = PptxWriterTool().execute(
            file_path=str(out),
            slides=[{"title": "Replaced"}],
            overwrite=True,
        )
        assert result.success
        assert _read_pptx_titles(out) == ["Replaced"]


# ════════════════════════════════════════════════════════════════════
# Limits
# ════════════════════════════════════════════════════════════════════

class TestLimits:

    def test_long_text_truncated(self, tmp_path):
        long_title = "x" * 10_000
        out = tmp_path / "deck.pptx"
        result = PptxWriterTool().execute(
            file_path=str(out),
            slides=[{"title": long_title}],
        )
        assert result.success
        title = _read_pptx_titles(out)[0]
        # Title should be clipped to MAX_CHARS_PER_TEXT
        assert len(title) < 10_000

    def test_too_many_bullets(self, tmp_path):
        out = tmp_path / "deck.pptx"
        result = PptxWriterTool().execute(
            file_path=str(out),
            slides=[{"title": "x", "bullets": [str(i) for i in range(50)]}],
        )
        assert not result.success
