"""
tools/builtins/pptx_writer.py — Build PowerPoint (.pptx) presentations.

Companion to PptxReaderTool. Lets the agent generate a deck from a list
of slide dicts:

    slides=[
        {"title": "Q3 Report", "subtitle": "Customer growth"},
        {"title": "Highlights", "bullets": ["Up 12%", "New region"]},
        {"title": "Next steps", "body": "Free-form paragraph text..."},
    ]

Currently supported slide shapes:
    title      — title only (layout 0)
    title+sub  — title + subtitle (layout 0)
    bullets    — title + bulleted body (layout 1)
    body       — title + free-text body (layout 1)
    blank      — title only on a blank layout (layout 6)

Safety:
    is_destructive=True — writes to disk; refuses to overwrite unless
                          overwrite=true is passed.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Pt

from ..base import ToolBase, ToolResult, ToolSpec


_MAX_SLIDES = 100
_MAX_BULLETS_PER_SLIDE = 20
_MAX_CHARS_PER_TEXT = 5_000


class PptxWriterTool(ToolBase):
    """Create a fresh `.pptx` file from a list of slide dicts."""

    @property
    def spec(self) -> ToolSpec:
        return ToolSpec(
            name="pptx_writer",
            description=(
                "Создаёт PowerPoint презентацию (.pptx) из списка слайдов. "
                "Каждый слайд — dict с полями title / subtitle / bullets / body. "
                "Пишет на диск — destructive."
            ),
            parameters={
                "file_path": "str — путь к итоговому .pptx",
                "slides":    "list[dict] — слайды (title/subtitle/bullets/body)",
                "title":     "str (optional) — общий заголовок презентации",
                "overwrite": "bool (optional, default false) — разрешить перезапись",
            },
            requires_approval=False,
            is_destructive=True,
        )

    # ────────────────────────────────────────────────────────────────

    def execute(self, **params: Any) -> ToolResult:
        file_path = str(params.get("file_path", "")).strip()
        slides = params.get("slides") or []
        overwrite = bool(params.get("overwrite", False))

        if not file_path:
            return self._fail("'file_path' is required")
        if not file_path.lower().endswith(".pptx"):
            return self._fail("file_path must end with .pptx")
        if not isinstance(slides, list) or not slides:
            return self._fail("'slides' must be a non-empty list of dicts")
        if len(slides) > _MAX_SLIDES:
            return self._fail(f"too many slides: {len(slides)} (max {_MAX_SLIDES})")

        path = Path(file_path).expanduser()
        if path.exists() and not overwrite:
            return self._fail(f"file already exists: {path} (pass overwrite=true)")

        try:
            path.parent.mkdir(parents=True, exist_ok=True)
        except OSError as exc:
            return self._fail(f"cannot create parent dir: {exc}")

        prs = Presentation()
        for idx, raw in enumerate(slides):
            if not isinstance(raw, dict):
                return self._fail(f"slide {idx}: must be a dict, got {type(raw).__name__}")
            slide_err = self._add_slide(prs, raw, idx)
            if slide_err is not None:
                return self._fail(slide_err)

        try:
            prs.save(str(path))
        except OSError as exc:
            return self._fail(f"failed to save: {exc}")

        size = path.stat().st_size
        return self._ok({
            "path":   str(path),
            "slides": len(slides),
            "size_bytes": size,
        })

    # ────────────────────────────────────────────────────────────────

    def _add_slide(self, prs: Presentation, data: dict, idx: int) -> str | None:
        title    = _clip(str(data.get("title", "")).strip())
        subtitle = _clip(str(data.get("subtitle", "")).strip())
        body     = _clip(str(data.get("body", "")).strip())
        bullets  = data.get("bullets") or []
        if bullets and not isinstance(bullets, list):
            return f"slide {idx}: 'bullets' must be a list"
        if len(bullets) > _MAX_BULLETS_PER_SLIDE:
            return f"slide {idx}: too many bullets ({len(bullets)} > {_MAX_BULLETS_PER_SLIDE})"

        # Pick layout based on what content the caller provided
        if bullets or body:
            layout = prs.slide_layouts[1]   # Title + Content
        elif subtitle:
            layout = prs.slide_layouts[0]   # Title + Subtitle
        else:
            layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0]

        slide = prs.slides.add_slide(layout)

        # Title
        if slide.shapes.title is not None:
            slide.shapes.title.text = title or f"Slide {idx + 1}"

        # Subtitle on layout 0
        if subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle

        # Bullets or body on layout 1
        if (bullets or body) and len(slide.placeholders) > 1:
            ph = slide.placeholders[1]
            tf = ph.text_frame
            tf.clear()

            if bullets:
                for i, b in enumerate(bullets):
                    text = _clip(str(b))
                    if i == 0:
                        tf.text = text
                    else:
                        p = tf.add_paragraph()
                        p.text = text
                        p.level = 0
            elif body:
                tf.text = body
                # No font setting beyond defaults — themes will style it
        return None


# ════════════════════════════════════════════════════════════════════
# Helpers
# ════════════════════════════════════════════════════════════════════

def _clip(text: str) -> str:
    if len(text) <= _MAX_CHARS_PER_TEXT:
        return text
    return text[: _MAX_CHARS_PER_TEXT - 1] + "…"
