"""
tools/builtins/pptx_reader.py — Чтение презентаций PowerPoint (.pptx)

Использует python-pptx (чистый Python, без LibreOffice/COM/Office).

Actions:
    extract_text   — весь текст слайдов (заголовки + контент)
    get_metadata   — автор, название, дата, кол-во слайдов, размер
    list_slides    — список слайдов с заголовками и количеством фигур
"""

from __future__ import annotations

import time
from pathlib import Path
from typing import Any

from ..base import ToolBase, ToolResult, ToolSpec

_MAX_FILE_SIZE_MB = 100
_MAX_CHARS_OUTPUT = 100_000


class PptxReaderTool(ToolBase):
    """
    Читает PowerPoint презентации (.pptx).

    params:
        file_path  (str): Путь к .pptx файлу
        action     (str, optional): extract_text | get_metadata | list_slides
                                    (по умолчанию: extract_text)
        max_chars  (int, optional): Максимум символов (по умолчанию 20000)
        slide      (int, optional): Номер слайда (1-based) для extract_text.
                                    Если не указан — все слайды.
    """

    @property
    def spec(self) -> ToolSpec:
        return ToolSpec(
            name="pptx_reader",
            description="Читает текст и метаданные из PowerPoint презентаций (.pptx)",
            parameters={
                "file_path": "str — путь к .pptx файлу",
                "action":    "str (optional) — extract_text | get_metadata | list_slides",
                "max_chars": "int (optional, default 20000) — максимум символов вывода",
                "slide":     "int (optional) — номер слайда (1-based), по умолчанию все",
            },
            requires_approval=False,
            is_destructive=False,
        )

    def execute(self, **params: Any) -> ToolResult:
        t0 = time.perf_counter()
        try:
            import pptx  # noqa: F401  (python-pptx)
        except ImportError:
            return self._fail("python-pptx не установлен. Запустите: pip install python-pptx")

        file_path = params.get("file_path", "")
        if not file_path:
            return self._fail("Параметр 'file_path' обязателен")

        path = Path(str(file_path))
        if not path.exists():
            return self._fail(f"Файл не найден: {file_path}")
        if path.suffix.lower() != ".pptx":
            return self._fail(
                f"Ожидается .pptx, получено: {path.suffix!r}. "
                "Для .ppt-файлов сначала пересохраните как .pptx в PowerPoint."
            )

        size_mb = path.stat().st_size / (1024 * 1024)
        if size_mb > _MAX_FILE_SIZE_MB:
            return self._fail(
                f"Файл слишком большой: {size_mb:.1f} МБ (лимит {_MAX_FILE_SIZE_MB} МБ)"
            )

        action = str(params.get("action", "extract_text")).lower()
        if action not in ("extract_text", "get_metadata", "list_slides"):
            return self._fail(
                f"Неизвестное действие: {action!r}. "
                "Доступно: extract_text | get_metadata | list_slides"
            )

        try:
            from pptx import Presentation  # type: ignore
            prs = Presentation(str(path))

            dur_ms = (time.perf_counter() - t0) * 1000

            if action == "extract_text":
                return self._extract_text(prs, path, params, dur_ms)
            elif action == "get_metadata":
                return self._get_metadata(prs, path, dur_ms)
            else:
                return self._list_slides(prs, dur_ms)

        except Exception as exc:
            return self._fail(f"Ошибка чтения .pptx: {exc}")

    # ------------------------------------------------------------------
    def _extract_text(self, prs: Any, path: Path, params: dict, dur_ms: float) -> ToolResult:
        max_chars = int(params.get("max_chars", 20_000))
        max_chars = min(max(1, max_chars), _MAX_CHARS_OUTPUT)

        slide_filter = params.get("slide")
        if slide_filter is not None:
            slide_filter = int(slide_filter)

        parts: list[str] = []
        total_slides = len(prs.slides)

        for idx, slide in enumerate(prs.slides, 1):
            if slide_filter is not None and idx != slide_filter:
                continue

            slide_texts: list[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        slide_texts.append(line)

            if slide_texts:
                header = f"\n═══ Слайд {idx}/{total_slides} ═══"
                parts.append(header)
                parts.extend(slide_texts)

        text = "\n".join(parts).strip()
        truncated = False
        if len(text) > max_chars:
            text = text[:max_chars]
            truncated = True

        return self._ok(
            text,
            duration_ms=round(dur_ms, 2),
            slides_total=total_slides,
            slide_filter=slide_filter,
            chars_returned=len(text),
            truncated=truncated,
            file_size_kb=round(path.stat().st_size / 1024, 1),
        )

    # ------------------------------------------------------------------
    def _get_metadata(self, prs: Any, path: Path, dur_ms: float) -> ToolResult:
        props = prs.core_properties

        def _s(v: Any) -> str | None:
            return str(v) if v else None

        stat = path.stat()
        return self._ok(
            {
                "title":         _s(props.title),
                "author":        _s(props.author),
                "subject":       _s(props.subject),
                "description":   _s(props.description),
                "created":       _s(props.created),
                "modified":      _s(props.modified),
                "slides":        len(prs.slides),
                "slide_width":   round(prs.slide_width.inches, 2) if prs.slide_width else None,
                "slide_height":  round(prs.slide_height.inches, 2) if prs.slide_height else None,
                "file_size_kb":  round(stat.st_size / 1024, 1),
            },
            duration_ms=round(dur_ms, 2),
        )

    # ------------------------------------------------------------------
    def _list_slides(self, prs: Any, dur_ms: float) -> ToolResult:
        slides_info: list[dict] = []
        for idx, slide in enumerate(prs.slides, 1):
            title_text = ""
            shapes_count = len(slide.shapes)
            notes_text = ""

            for shape in slide.shapes:
                # Заголовок слайда — первый shape с типом TITLE или первый текстовый
                if shape.has_text_frame:
                    name = shape.name.lower()
                    if "title" in name and not title_text:
                        title_text = shape.text_frame.text.strip()

            # Заметки докладчика
            if slide.has_notes_slide:
                tf = slide.notes_slide.notes_text_frame
                notes_text = tf.text.strip()[:200] if tf else ""

            slides_info.append({
                "slide":  idx,
                "title":  title_text or f"(Слайд {idx})",
                "shapes": shapes_count,
                "notes":  notes_text,
            })

        return self._ok(
            slides_info,
            duration_ms=round(dur_ms, 2),
            total_slides=len(slides_info),
        )
