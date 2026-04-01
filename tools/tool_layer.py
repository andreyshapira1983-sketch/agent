# Tool Layer (инструменты) — Слой 5
# Архитектура автономного AI-агента
# Система инструментов для действий: браузер, терминал, файловая система,
# Python runtime, GitHub API, Docker, пакетные менеджеры, БД, облако, поиск.
#
# SECURITY: Все инструменты защищены — локальный мозг контролирует LLM, не наоборот.
# pylint: disable=broad-except,arguments-differ,import-error,no-name-in-module
# pylint: disable=unspecified-encoding,keyword-arg-before-vararg,using-constant-test
# pylint: disable=redefined-outer-name,unused-variable,unused-argument,subprocess-run-check

from __future__ import annotations

import subprocess
import os
import re
import shlex
import socket
import threading
import time
from abc import ABC, abstractmethod
from typing import Any


# ── Базовый класс инструмента ─────────────────────────────────────────────────

class BaseTool(ABC):
    """Базовый класс для всех инструментов.
    
    Подклассы могут переопределять run() с любой сигнатурой параметров.
    Это позволяет каждому инструменту иметь специфичный API.
    """

    def __init__(self, name: str, description: str):
        self.name = name
        self.description = description
        self._last_result = None

    @abstractmethod
    def run(self, *args, **kwargs) -> dict:
        """Выполняет инструмент. Возвращает результат.
        
        Подклассы переопределяют с конкретными параметрами.
        Например:
            def run(self, command: str, timeout: int | None = None) -> dict:
        """

    def __call__(self, *args, **kwargs):
        """Вызов инструмента через (). Сохраняет последний результат."""
        result = self.run(*args, **kwargs)
        self._last_result = result
        return result

    @property
    def last_result(self):
        """Последний результат выполнения."""
        return self._last_result


# ── Security: Command Validator ────────────────────────────────────────────────

class CommandValidator:
    """
    Валидатор команд для TerminalTool.
    Whitelist-подход: разрешено только то, что явно указано.
    Локальный мозг решает — LLM не может обойти.
    """

    # Разрешённые исполняемые файлы (без аргументов)
    ALLOWED_EXECUTABLES = frozenset({
        # Python
        'python', 'python3', 'pip', 'pip3', 'poetry', 'pipenv', 'venv',
        # Node
        'node', 'npm', 'npx', 'yarn', 'pnpm',
        # Git
        'git',
        # Системные (безопасные, read-only)
        'ls', 'dir', 'cat', 'head', 'tail', 'echo', 'pwd', 'cd',
        'find', 'grep', 'wc', 'sort', 'uniq', 'diff', 'which', 'where',
        'type', 'file', 'stat', 'whoami', 'hostname', 'date', 'env',
        'printenv', 'uname',
        # Docker (с ограничениями)
        'docker', 'docker-compose',
        # Тестирование
        'pytest', 'unittest', 'jest', 'mocha',
        # Сборка
        'make', 'cmake', 'cargo', 'go', 'javac', 'java', 'gcc', 'g++',
    })

    # Запрещённые паттерны в аргументах (даже для разрешённых exe)
    BLOCKED_PATTERNS = [
        re.compile(r';\s*'),               # цепочка команд через ;
        re.compile(r'\|\|'),                # логический OR
        re.compile(r'&&'),                  # логический AND
        re.compile(r'\|'),                  # pipe
        re.compile(r'`'),                   # backtick подстановка
        re.compile(r'\$\('),               # command substitution
        re.compile(r'>\s*'),               # перенаправление вывода
        re.compile(r'<\s*'),               # перенаправление ввода
        re.compile(r'\.\./\.\./'),         # глубокий path traversal
    ]

    # Полностью запрещённые команды
    BLOCKED_COMMANDS = frozenset({
        'rm', 'rmdir', 'del', 'format', 'mkfs', 'dd',
        'shutdown', 'reboot', 'halt', 'poweroff', 'init',
        'sudo', 'su', 'runas', 'doas',
        'chmod', 'chown', 'chgrp',
        'kill', 'killall', 'pkill', 'taskkill',
        'curl', 'wget', 'nc', 'ncat', 'netcat', 'nmap',
        'ssh', 'scp', 'rsync', 'ftp', 'sftp', 'telnet',
        'eval', 'exec', 'source', 'bash', 'sh', 'cmd', 'powershell', 'pwsh',
        'reg', 'regedit', 'sc', 'net', 'netsh', 'wmic',
        'crontab', 'at', 'schtasks',
    })

    @classmethod
    def validate(cls, command: str) -> tuple[bool, str]:
        """
        Проверяет команду. Возвращает (разрешена, причина).
        Whitelist: если exe нет в списке — блокируется.
        """
        if not command or not command.strip():
            return False, "Пустая команда"

        command = command.strip()

        # Парсим первое слово — имя исполняемого файла
        try:
            parts = shlex.split(command)
        except ValueError:
            return False, "Невалидный синтаксис команды"

        if not parts:
            return False, "Пустая команда"

        exe = os.path.basename(parts[0]).lower()
        # На Windows убираем .exe
        if exe.endswith('.exe'):
            exe = exe[:-4]

        # Проверка блокированных команд
        if exe in cls.BLOCKED_COMMANDS:
            return False, f"Команда '{exe}' запрещена"

        # Проверка whitelist
        if exe not in cls.ALLOWED_EXECUTABLES:
            return False, f"Команда '{exe}' не в списке разрешённых"

        # Проверка опасных паттернов в аргументах
        args_str = command[len(parts[0]):] if len(command) > len(parts[0]) else ''
        for pattern in cls.BLOCKED_PATTERNS:
            if pattern.search(args_str):
                return False, f"Обнаружен запрещённый паттерн в аргументах: {pattern.pattern}"

        # SECURITY (SSRF): блокируем `python -c` с сетевыми модулями
        if exe in ('python', 'python3') and '-c' in parts:
            inline_code = ' '.join(parts[parts.index('-c') + 1:])
            _NET_KEYWORDS = ('socket', 'http.', 'urllib.request', 'requests',
                             'httpx', 'aiohttp', 'urlopen', 'urlretrieve')
            for _nk in _NET_KEYWORDS:
                if _nk in inline_code:
                    return False, f"Сетевой код через python -c запрещён ('{_nk}')"

        # SECURITY (SSRF): блокируем `python script.py` если скрипт содержит сетевой код
        if exe in ('python', 'python3') and '-c' not in parts and '-m' not in parts:
            for _arg in parts[1:]:
                if _arg.startswith('-'):
                    continue
                # Первый позиционный аргумент — это имя скрипта
                if _arg.endswith('.py') and os.path.isfile(_arg):
                    try:
                        with open(_arg, 'r', encoding='utf-8', errors='replace') as _f:
                            _script = _f.read(64_000)  # читаем макс. 64KB
                        _ok, _reason = CodeValidator.validate(_script)
                        if not _ok:
                            return False, f"Скрипт '{_arg}' не прошёл проверку: {_reason}"
                    except OSError:
                        pass
                break  # проверяем только первый позиционный

        return True, "OK"


# ── Security: Path Validator ────────────────────────────────────────────────────

class PathValidator:
    """
    Валидатор путей для FileSystemTool.
    Jail-подход: все пути должны быть внутри разрешённой директории.
    Блокирует чувствительные файлы.
    """

    # Файлы, к которым доступ запрещён всегда
    BLOCKED_FILES = frozenset({
        '.env', '.env.local', '.env.production', '.env.staging',
        'credentials.json', 'secrets.json', 'secrets.yaml', 'secrets.yml',
        'id_rsa', 'id_ed25519', 'id_ecdsa',
        '.ssh', '.gnupg', '.aws', '.azure',
        'shadow', 'passwd',
        '.git/config',
    })

    BLOCKED_EXTENSIONS = frozenset({
        '.pem', '.key', '.p12', '.pfx', '.keystore',
    })

    @classmethod
    def validate(cls, path: str, base_dir: str, action: str = 'read') -> tuple[bool, str]:
        """
        Проверяет, что путь безопасен.
        Возвращает (разрешён, причина).
        """
        if not path:
            return False, "Пустой путь"

        # Нормализуем пути
        try:
            resolved = os.path.realpath(os.path.normpath(
                path if os.path.isabs(path) else os.path.join(base_dir, path)
            ))
            base_resolved = os.path.realpath(os.path.normpath(base_dir))
        except (ValueError, OSError) as e:
            return False, f"Ошибка в пути: {e}"

        # Jail check: путь должен быть внутри base_dir
        if not resolved.startswith(base_resolved + os.sep) and resolved != base_resolved:
            return False, "Путь выходит за пределы рабочей директории"

        # Проверка имени файла
        filename = os.path.basename(resolved).lower()
        _, ext = os.path.splitext(filename)

        if filename in cls.BLOCKED_FILES:
            return False, f"Доступ к '{filename}' запрещён (чувствительный файл)"

        if ext in cls.BLOCKED_EXTENSIONS:
            return False, f"Доступ к файлам '{ext}' запрещён (ключи/сертификаты)"

        # Проверка компонентов пути на чувствительные директории
        parts = resolved.lower().split(os.sep)
        for blocked_dir in ('.ssh', '.gnupg', '.aws', '.azure'):
            if blocked_dir in parts:
                return False, f"Доступ к директории '{blocked_dir}' запрещён"

        # Для записи/удаления — дополнительная проверка
        if action in ('write', 'append', 'delete'):
            # Нельзя удалять/менять критические файлы агента
            if '.agent_memory' in resolved and action == 'delete':
                return False, "Удаление файлов памяти запрещено"

        return True, "OK"


# ── Security: Code Validator ────────────────────────────────────────────────────

class CodeValidator:
    """
    Валидатор Python-кода для PythonRuntimeTool.
    Блокирует опасные импорты, вызовы и конструкции.
    """

    BLOCKED_IMPORTS = frozenset({
        # os/pathlib/shutil разрешены (доступ к ФС ограничен _safe_open + PathValidator)
        # sys разрешён через SAFE_MODULES — нужен многим пакетам внутри
        'subprocess',
        'socket', 'http', 'urllib', 'requests', 'httpx', 'aiohttp',
        'ctypes', 'multiprocessing', 'threading',
        'signal', 'resource', 'pty', 'fcntl', 'termios',
        'importlib', 'pkgutil', 'zipimport',
        'pickle', 'shelve', 'marshal',
        'code', 'codeop', 'compile', 'compileall',
        'webbrowser', 'antigravity',
        'builtins', '__builtin__',
    })

    BLOCKED_CALLS = [
        re.compile(r'\bexec\s*\('),             # exec()
        re.compile(r'\beval\s*\('),             # eval()
        re.compile(r'\bcompile\s*\('),          # compile()
        re.compile(r'\b__import__\s*\('),       # __import__()
        re.compile(r'\bgetattr\s*\('),          # getattr()
        re.compile(r'\bsetattr\s*\('),          # setattr()
        re.compile(r'\bdelattr\s*\('),          # delattr()
        re.compile(r'\bglobals\s*\('),          # globals()
        re.compile(r'\blocals\s*\('),           # locals()
        re.compile(r'\bvars\s*\('),             # vars()
        re.compile(r'\bdir\s*\('),              # dir() — утечка namespace
        re.compile(r'\bbreakpoint\s*\('),       # breakpoint()
        re.compile(r'\bexit\s*\('),             # exit()
        re.compile(r'\bquit\s*\('),             # quit()
        # Опасные dunder-атрибуты (но НЕ __name__, __main__, __init__, __doc__ и т.д.)
        re.compile(r'__subclasses__'),          # __subclasses__() — побег из sandbox
        re.compile(r'__builtins__'),            # __builtins__ — доступ к builtins
        re.compile(r'__globals__'),             # __globals__ — утечка глобалов
        re.compile(r'__code__'),               # __code__ — модификация кода
        re.compile(r'__mro__'),                # __mro__ — обход иерархии
        re.compile(r'__bases__'),              # __bases__ — изменение наследования
        re.compile(r'__class__\s*\.\s*__'),    # __class__.__x__ — цепочки побега
        # Опасные вызовы os — блокируем явно даже при разрешённом импорте
        re.compile(r'\bos\s*\.\s*system\s*\('),   # os.system()
        re.compile(r'\bos\s*\.\s*popen\s*\('),    # os.popen()
        re.compile(r'\bos\s*\.\s*exec[vlpe]+\s*\('),  # os.execv/execl/etc.
        re.compile(r'\bos\s*\.\s*fork\s*\('),     # os.fork()
        re.compile(r'\bos\s*\.\s*spawn[vlpe]+\s*\('), # os.spawnv/etc.
    ]

    # Подмодули, разрешённые явно несмотря на блокировку корневого пакета
    ALLOWED_SUBMODULES = frozenset({
        'urllib.parse',       # парсинг URL, без сети
    })

    @classmethod
    def validate(cls, code: str) -> tuple[bool, str]:
        """Проверяет Python-код. Возвращает (безопасен, причина)."""
        if not code or not code.strip():
            return False, "Пустой код"

        # Проверяем синтаксис через ast.parse до любого выполнения
        import ast as _ast
        try:
            _ast.parse(code)
        except IndentationError:
            # IndentationError (неверные отступы) исправляется автоматически
            # в PythonRuntimeTool.run() — пропускаем здесь, не блокируем код
            pass
        except SyntaxError as e:
            # Даём LLM строку с ошибкой и номер строки чтобы исправить
            bad_line = ''
            dsl_hint = ''
            if e.lineno is not None:
                lines = code.splitlines()
                if 0 < e.lineno <= len(lines):
                    raw = lines[e.lineno - 1].strip()
                    bad_line = f" | строка {e.lineno}: {raw[:60]}"
                    # Частая ошибка: агент вставляет DSL-команду (SEARCH:/READ:/WRITE:) внутрь python-блока
                    if re.match(r'^(SEARCH|READ|WRITE|bash|BASH)\s*:', raw, re.IGNORECASE):
                        dsl_hint = (' ПОДСКАЗКА: командa DSL (SEARCH:/READ:/WRITE:) '
                                    'не является Python — используй отдельный блок действия, '
                                    'не вставляй внутрь python-кода.')
            # Частая ошибка: код обрезан по лимиту токенов (незакрытые строки/скобки)
            truncation_msgs = (
                'unterminated string literal',
                "was never closed",
                'EOF while scanning',
                'unexpected EOF',
            )
            msg_text = str(e.msg) if e.msg else ''
            if any(t in msg_text for t in truncation_msgs) and not dsl_hint:
                dsl_hint = (' ПОДСКАЗКА: код, вероятно, обрезан по лимиту токенов. '
                            'Пиши компактный код — не более 50 строк, без длинных '
                            'вложенных словарей/строк с кириллицей в качестве ключей.')
            return False, f"Синтаксическая ошибка в строке {e.lineno}{bad_line}: {e.msg}{dsl_hint}"

        # Проверяем импорты
        import_pattern = re.compile(
            r'(?:^|\n)\s*(?:import|from)\s+([\w.]+)', re.MULTILINE
        )
        for match in import_pattern.finditer(code):
            full_module = match.group(1)
            root_module = full_module.split('.')[0]
            # Разрешаем явно указанные безопасные подмодули
            if full_module in cls.ALLOWED_SUBMODULES:
                continue
            if root_module in cls.BLOCKED_IMPORTS:
                return False, (
                    f"Запрещённый импорт: '{root_module}'. "
                    f"Для системных команд используй ```bash блок. "
                    f"Для информации о системе используй psutil (разрешён)."
                )

        # Проверяем опасные вызовы
        for pattern in cls.BLOCKED_CALLS:
            found = pattern.search(code)
            if found:
                return False, f"Запрещённый вызов: '{found.group()}'"

        return True, "OK"


# ── Инструменты ───────────────────────────────────────────────────────────────

class TerminalTool(BaseTool):
    """
    Terminal — запуск команд в терминале.
    SECURITY: Whitelist-подход. shell=False. Валидация через CommandValidator.
    Используется совместно с Execution System (Слой 8).
    """

    def __init__(self, working_dir=None, timeout=30):
        super().__init__('terminal', 'Запуск команд в терминале/оболочке')
        self.working_dir = working_dir
        self.timeout = timeout

    def run(self, command: str, timeout: int | None = None) -> dict:
        """
        Args:
            command -- строка команды
            timeout -- таймаут в секундах
        Returns:
            {'stdout': ..., 'stderr': ..., 'returncode': ..., 'success': bool}
        """
        # SECURITY: Валидация команды через whitelist
        allowed, reason = CommandValidator.validate(command)
        if not allowed:
            return {
                'error': f'Команда заблокирована: {reason}',
                'blocked_command': command,
                'success': False,
            }

        try:
            # SECURITY: все вызовы идут через центральный CommandGateway
            from execution.command_gateway import CommandGateway
            gw = CommandGateway.get_instance()
            args = shlex.split(command)
            r = gw.execute(
                args,
                timeout=timeout or self.timeout,
                cwd=self.working_dir,
                caller='TerminalTool.run',
            )
            if not r.allowed:
                return {
                    'error': f'Команда заблокирована: {r.reject_reason}',
                    'blocked_command': command,
                    'success': False,
                }
            return {
                'stdout': r.stdout,
                'stderr': r.stderr,
                'returncode': r.returncode,
                'success': r.returncode == 0,
            }
        except Exception as e:
            return {'error': str(e), 'success': False}


class FileSystemTool(BaseTool):
    """
    Filesystem — работа с файловой системой:
    чтение, запись, удаление, список файлов, проверка существования.
    SECURITY: Path jail + блокировка чувствительных файлов.
    """

    def __init__(self, base_dir=None):
        super().__init__('filesystem', 'Чтение, запись и управление файлами')
        self.base_dir = os.path.realpath(base_dir or os.getcwd())

    def run(self, action: str, path: str, content: str | None = None, encoding: str = 'utf-8') -> dict:
        """
        Args:
            action  -- 'read' | 'write' | 'append' | 'delete' | 'list' | 'exists'
            path    -- путь к файлу или директории
            content -- содержимое для записи (для action='write'/'append')
        """
        # SECURITY: Валидация пути через PathValidator (jail + blocked files)
        allowed, reason = PathValidator.validate(path, self.base_dir, action)
        if not allowed:
            return {
                'error': f'Доступ заблокирован: {reason}',
                'blocked_path': path,
                'success': False,
            }

        full_path = os.path.realpath(
            path if os.path.isabs(path) else os.path.join(self.base_dir, path)
        )
        try:
            if action == 'read':
                if not os.path.exists(full_path):
                    return {'error': f'Файл не найден: {full_path}', 'success': False}
                if os.path.isdir(full_path):
                    return {'error': f'Путь является директорией: {full_path}', 'success': False}
                # SECURITY: ограничение размера читаемого файла (10 MB)
                file_size = os.path.getsize(full_path)
                if file_size > 10 * 1024 * 1024:
                    return {'error': 'Файл слишком большой (>10 MB)', 'success': False}
                # Пробуем utf-8, при ошибке — latin-1 (читает любой байт без исключений)
                try:
                    with open(full_path, 'r', encoding=encoding, errors='replace') as f:
                        return {'content': f.read(), 'success': True}
                except (UnicodeDecodeError, LookupError):
                    with open(full_path, 'r', encoding='latin-1') as f:
                        return {'content': f.read(), 'success': True}

            elif action == 'write':
                os.makedirs(os.path.dirname(full_path), exist_ok=True)
                with open(full_path, 'w', encoding=encoding) as f:
                    f.write(content or '')
                return {'written': full_path, 'success': True}

            elif action == 'append':
                os.makedirs(os.path.dirname(full_path), exist_ok=True)
                with open(full_path, 'a', encoding=encoding) as f:
                    f.write(content or '')
                return {'appended': full_path, 'success': True}

            elif action == 'delete':
                if not os.path.exists(full_path):
                    return {'error': f'Файл не найден: {full_path}', 'success': False}
                os.remove(full_path)
                return {'deleted': full_path, 'success': True}

            elif action == 'list':
                if not os.path.isdir(full_path):
                    return {'error': f'Директория не найдена: {full_path}', 'success': False}
                entries = os.listdir(full_path)
                return {'entries': entries, 'success': True}

            elif action == 'exists':
                return {'exists': os.path.exists(full_path), 'success': True}

            else:
                return {'error': f"Неизвестный action: {action}", 'success': False}

        except (FileNotFoundError, PermissionError, IsADirectoryError,
                NotADirectoryError, OSError, UnicodeDecodeError) as e:
            return {'error': str(e), 'success': False}

    # Удобные методы
    def read(self, path: str) -> str | None:
        result = self.run('read', path)
        return result.get('content') if result.get('success') else None

    def write(self, path: str, content: str) -> bool:
        return self.run('write', path, content=content).get('success', False)

    def exists(self, path: str) -> bool:
        return self.run('exists', path).get('exists', False)

    def list_dir(self, path: str = '.') -> list:
        return self.run('list', path).get('entries', [])


class PythonRuntimeTool(BaseTool):
    """
    Python Runtime — выполнение Python-кода в изолированном namespace.
    SECURITY: Ограниченные builtins, запрет опасных импортов/вызовов.
    """

    # Безопасные builtins (без open, exec, eval, compile, __import__, etc.)
    SAFE_BUILTINS = {
        'abs': abs, 'all': all, 'any': any, 'bin': bin, 'bool': bool,
        'bytes': bytes, 'callable': callable, 'chr': chr, 'complex': complex,
        'dict': dict, 'divmod': divmod, 'enumerate': enumerate,
        'filter': filter, 'float': float, 'format': format,
        'frozenset': frozenset, 'hash': hash, 'hex': hex, 'id': id,
        'int': int, 'isinstance': isinstance, 'issubclass': issubclass,
        'iter': iter, 'len': len, 'list': list, 'map': map, 'max': max,
        'min': min, 'next': next, 'object': object, 'oct': oct, 'ord': ord,
        'pow': pow, 'print': print, 'range': range, 'repr': repr,
        'reversed': reversed, 'round': round, 'set': set, 'slice': slice,
        'sorted': sorted, 'str': str, 'sum': sum, 'super': super,
        'tuple': tuple, 'type': type, 'zip': zip,
        'True': True, 'False': False, 'None': None,
        # Необходимо для определения классов через class X: в sandbox
        '__build_class__': __build_class__,
        '__name__': '__main__',
        # Встроенные исключения — необходимы для try/except в sandbox-коде
        'Exception': Exception, 'BaseException': BaseException,
        'ValueError': ValueError, 'TypeError': TypeError,
        'KeyError': KeyError, 'IndexError': IndexError,
        'AttributeError': AttributeError, 'NameError': NameError,
        'RuntimeError': RuntimeError, 'StopIteration': StopIteration,
        'NotImplementedError': NotImplementedError, 'ArithmeticError': ArithmeticError,
        'ZeroDivisionError': ZeroDivisionError, 'OverflowError': OverflowError,
        'MemoryError': MemoryError, 'RecursionError': RecursionError,
        'IOError': IOError, 'OSError': OSError, 'PermissionError': PermissionError,
        'FileNotFoundError': FileNotFoundError, 'FileExistsError': FileExistsError,
        'IsADirectoryError': IsADirectoryError, 'NotADirectoryError': NotADirectoryError,
        'TimeoutError': TimeoutError, 'ConnectionError': ConnectionError,
        'UnicodeError': UnicodeError, 'UnicodeDecodeError': UnicodeDecodeError,
        'UnicodeEncodeError': UnicodeEncodeError,
        'AssertionError': AssertionError, 'ImportError': ImportError,
        'ModuleNotFoundError': ModuleNotFoundError,
        'LookupError': LookupError, 'BufferError': BufferError,
        'Warning': Warning, 'UserWarning': UserWarning,
    }

    # Модули, разрешённые для импорта в sandbox
    SAFE_MODULES = frozenset({
        # ── Ядро ──────────────────────────────────────────────────────────────
        # sys/platform — нужны многим пакетам внутри; доступ ограничен _safe_import
        'sys', 'platform',
        # ── Ядро ──────────────────────────────────────────────────────────────
        'psutil', 'json', 'math', 'datetime', 'time', 're',
        'collections', 'itertools', 'functools', 'io',
        'hashlib', 'base64', 'decimal', 'fractions',
        'statistics', 'random', 'string', 'textwrap', 'copy',
        'abc', 'typing', 'dataclasses', 'enum',
        'csv', 'configparser', 'struct', 'logging',
        # ── Файловая система (open() ограничен _safe_open) ───────────────────
        'os', 'os.path', 'pathlib', 'glob', 'shutil',
        # ── Структуры данных и алгоритмы──────────────────────────────────────
        'heapq', 'bisect', 'array', 'queue', 'pprint',
        'operator', 'weakref', 'contextlib', 'warnings',
        # ── Математика ───────────────────────────────────────────────────────
        'cmath', 'numbers',
        # ── Строки и локали ───────────────────────────────────────────────────
        'unicodedata', 'difflib', 'fnmatch', 'calendar', 'locale',
        # ── UUID, traceback, inspect ──────────────────────────────────────────
        'uuid', 'traceback', 'inspect',
        # ── Архивы и сжатие (read/write локально) ────────────────────────────
        'zipfile', 'gzip', 'bz2', 'lzma', 'zlib', 'rarfile',
        # ── БД и хранение ─────────────────────────────────────────────────────
        'sqlite3',
        # ── Дата и время ─────────────────────────────────────────────────────
        'zoneinfo',
        # ── Парсинг URL (без сети) ────────────────────────────────────────────
        'urllib.parse',
        # ── HTML/XML парсинг ──────────────────────────────────────────────────
        'html', 'html.parser', 'xml', 'xml.etree',
        'bs4', 'lxml',
        # ── Документы и офисные форматы ──────────────────────────────────────
        'pypdf', 'docx', 'pptx',
        'reportlab', 'reportlab.lib', 'reportlab.lib.pagesizes',
        'reportlab.lib.styles', 'reportlab.lib.units', 'reportlab.lib.colors',
        'reportlab.platypus', 'reportlab.pdfgen', 'reportlab.pdfgen.canvas',
        'openpyxl', 'openpyxl.styles', 'openpyxl.utils',
        'python_docx',
        # ── Визуализация (Agg backend, без GUI) ───────────────────────────────
        'matplotlib', 'matplotlib.pyplot', 'matplotlib.figure',
        'matplotlib.backends', 'seaborn',
        # ── Научные вычисления ────────────────────────────────────────────────
        'numpy', 'scipy',
        # ── Данные ────────────────────────────────────────────────────────────
        'pandas',
        # ── Изображения ──────────────────────────────────────────────────────
        'PIL', 'Pillow',
        # ── Конфигурация ─────────────────────────────────────────────────────
        'yaml', 'toml',
        # ── Прочие полезные ───────────────────────────────────────────────────
        'tabulate', 'tqdm', 'dateutil', 'pytz',
        # ── Windows Event Log API (pywin32) ───────────────────────────────────
        'win32evtlog', 'win32con', 'win32security', 'pywintypes',
        'win32evtlogutil',
    })

    def __init__(self, working_dir: str | None = None):
        super().__init__('python_runtime', 'Выполнение Python-кода')
        import os as _os
        self._working_dir = _os.path.abspath(working_dir or _os.getcwd())
        self._namespace = {}
        self._tool_layer = None  # устанавливается из agent.py после build_tool_layer()

    def _make_safe_builtins(self) -> dict:
        """Builds safe builtins with restricted import and open."""
        import builtins as _builtins
        safe_modules = self.SAFE_MODULES
        working_dir  = self._working_dir

        def _safe_import(name, glb=None, loc=None, fromlist=(), level=0):
            top = name.split('.')[0]
            # Разрешаем если:
            #   1) корневой пакет прямо в allowlist, или
            #   2) точное имя (напр. urllib.parse) в allowlist, или
            #   3) корень — часть разрешённого пути (urllib разрешён т.к. urllib.parse есть)
            if (top not in safe_modules
                    and name not in safe_modules
                    and not any(m.startswith(top + '.') for m in safe_modules)):
                raise ImportError(f"\u0418мпорт '{name}' запрещён в sandbox")
            # Дополнительно: если импортируется подмодуль, он сам должен быть разрешён
            if name != top and name not in safe_modules and top not in safe_modules:
                raise ImportError(f"\u0418мпорт '{name}' запрещён в sandbox")
            return _builtins.__import__(name, glb, loc, fromlist, level)

        # Паттерны путей-плейсхолдеров, которые LLM вставляет как шаблоны
        _PLACEHOLDER_SEGMENTS = (
            '/path/to/', '\\path\\to\\', r'\path\to\\',
            '/путь/к/', '\\путь\\к\\',
            '/your/path/', '\\your\\path\\',
            '/log/file', '\\log\\file',
            'путь/к/файлу', 'path/to/file',
            '/var/log/', 'c:/path/', 'c:\\path\\',
        )

        # Корневые системные директории Windows — запись туда перенаправляется в outputs/
        _WIN_SYSTEM_ROOTS = (
            'c:/windows/', 'c:/program files', 'c:/programdata',
            'c:/logs/', 'c:/log/', 'c:/temp/', 'c:/tmp/',
            'c:/users/', 'c:/system',
        )

        def _safe_open(*_open_args, **kwargs):
            import os as _os
            file = _open_args[0]
            mode = _open_args[1] if len(_open_args) > 1 else kwargs.pop('mode', 'r')
            extra_args = _open_args[2:]
            if isinstance(file, (int,)):          # fd allowed
                return _builtins.open(file, mode, *extra_args, **kwargs)
            # На Windows по умолчанию ставим UTF-8 для текстовых операций
            if 'b' not in mode and 'encoding' not in kwargs:
                kwargs['encoding'] = 'utf-8'

            file_str = str(file)
            file_lower = file_str.replace('\\', '/').lower()

            # Перенаправляем пути-плейсхолдеры в outputs/
            for seg in _PLACEHOLDER_SEGMENTS:
                if seg.replace('\\', '/').lower() in file_lower:
                    basename = _os.path.basename(file_str) or 'output.txt'
                    redirected = _os.path.join(working_dir, 'outputs', basename)
                    _os.makedirs(_os.path.dirname(redirected), exist_ok=True)
                    return _builtins.open(redirected, mode, *extra_args, **kwargs)

            resolved = _os.path.realpath(
                file_str if _os.path.isabs(file_str) else _os.path.join(working_dir, file_str)
            )
            wd_resolved = _os.path.realpath(working_dir)
            if not resolved.startswith(wd_resolved + _os.sep) and resolved != wd_resolved:
                # Специальная подсказка для Windows Event Log (.evtx)
                if file_str.lower().endswith('.evtx'):
                    raise PermissionError(
                        f"Файл '{file}' недоступен через open(). "
                        "Используйте win32evtlog для чтения журналов событий Windows: "
                        "import win32evtlog; h = win32evtlog.OpenEventLog(None, 'Application'); "
                        "events = win32evtlog.ReadEventLog(h, win32evtlog.EVENTLOG_BACKWARDS_READ|win32evtlog.EVENTLOG_SEQUENTIAL_READ, 0)"
                    )
                # Чтение или запись вне рабочей директории — перенаправляем в outputs/
                basename = _os.path.basename(file_str) or 'output.txt'
                redirected = _os.path.join(working_dir, 'outputs', basename)
                is_write = any(c in mode for c in ('w', 'a'))
                if is_write:
                    _os.makedirs(_os.path.dirname(redirected), exist_ok=True)
                    return _builtins.open(redirected, mode, *extra_args, **kwargs)
                # Чтение: если файл есть в outputs/ — читаем оттуда;
                # иначе возвращаем FileNotFoundError (не PermissionError)
                if _os.path.exists(redirected):
                    return _builtins.open(redirected, mode, *extra_args, **kwargs)
                raise FileNotFoundError(
                    f"Файл '{basename}' не найден ни по пути '{file}', "
                    f"ни в outputs/. Сначала создайте его через WRITE."
                )
            # Запрещаем опасные режимы
            if any(c in mode for c in ('x',)):
                raise PermissionError("Режим 'x' не разрешён")
            # Если файл пишется прямо в корень проекта (не в подпапку) — redirect в outputs/
            is_write = any(c in mode for c in ('w', 'a'))
            parent_dir = _os.path.realpath(_os.path.dirname(resolved))
            if is_write and parent_dir == wd_resolved:
                basename = _os.path.basename(resolved)
                # Не трогаем служебные файлы корня
                _root_keep = ('.env', 'agent.py', 'requirements.txt',
                              'log.txt', 'task_log.txt', 'system_log.txt')
                if basename not in _root_keep:
                    redirected = _os.path.join(working_dir, 'outputs', basename)
                    _os.makedirs(_os.path.dirname(redirected), exist_ok=True)
                    return _builtins.open(redirected, mode, *extra_args, **kwargs)
            # Создаём директорию если нужна (write/append внутри working_dir)
            if is_write:
                _os.makedirs(_os.path.dirname(resolved) or working_dir, exist_ok=True)
            return _builtins.open(resolved, mode, *extra_args, **kwargs)

        result = dict(self.SAFE_BUILTINS)
        result['__import__'] = _safe_import
        result['open'] = _safe_open
        return result

    def run(self, code: str, reset_namespace: bool = False) -> dict:
        """
        Args:
            code            -- Python код для выполнения
            reset_namespace -- сбросить namespace перед выполнением
        Returns:
            {'output': ..., 'result': ..., 'success': bool, 'error': ...}
        """
        # SECURITY: Валидация кода через CodeValidator
        allowed, reason = CodeValidator.validate(code)
        if not allowed:
            return {
                'error': f'Код заблокирован: {reason}',
                'success': False,
            }

        # Нормализуем отступы: LLM часто генерирует код с общим отступом
        import textwrap as _textwrap
        code = _textwrap.dedent(code).strip()
        if not code:
            return {'output': '', 'result': None, 'success': True}

        if reset_namespace:
            self._namespace = {}

        import io
        import sys
        stdout_capture = io.StringIO()
        old_stdout = sys.stdout
        sys.stdout = stdout_capture

        try:
            # SECURITY: Ограниченные builtins + безопасные import/open
            safe_globals: dict[str, Any] = {'__builtins__': self._make_safe_builtins()}
            safe_globals.update(self._namespace)
            # Инжектируем tool_layer чтобы код мог вызывать инструменты напрямую
            if self._tool_layer is not None:
                safe_globals['tool_layer'] = self._tool_layer

            # Инжектируем патченый psutil: win_service_iter пропускает недоступные сервисы
            try:
                import psutil as _psutil
                import types as _types
                _psutil_proxy = _types.ModuleType('psutil')
                _psutil_proxy.__dict__.update({
                    k: v for k, v in _psutil.__dict__.items()
                    if not k.startswith('__')
                })

                _orig_svc_iter = getattr(_psutil, 'win_service_iter', None)
                _orig_svc_get  = getattr(_psutil, 'win_service_get', None)

                # Враппер всегда определён — используется и iter, и get
                class _SafeService:
                    """Враппер psutil.WindowsService: перехватывает OSError при QSC."""
                    def __init__(self, svc):
                        self._svc = svc
                    def name(self):
                        try:
                            return self._svc.name()
                        except OSError:
                            return '<недоступен>'
                    def display_name(self):
                        try:
                            return self._svc.display_name()
                        except OSError:
                            return '<недоступен>'
                    def status(self):
                        try:
                            return self._svc.status()
                        except OSError:
                            return 'unknown'
                    def start_type(self):
                        try:
                            return self._svc.start_type()
                        except OSError:
                            return 'unknown'
                    def pid(self):
                        try:
                            return self._svc.pid()
                        except OSError:
                            return None
                    def binpath(self):
                        try:
                            return self._svc.binpath()
                        except OSError:
                            return ''
                    def username(self):
                        try:
                            return self._svc.username()
                        except OSError:
                            return ''
                    def as_dict(self, fields=None, attrs=None):
                        all_attrs = ('name', 'display_name', 'status',
                                     'start_type', 'pid', 'binpath', 'username')
                        # psutil оригинальный API принимает `attrs=` или `fields=`
                        _filter = fields or attrs
                        attrs = (tuple(f for f in _filter if f in all_attrs)
                                 if _filter else all_attrs)
                        result = {}
                        for attr in attrs:
                            try:
                                result[attr] = getattr(self, attr)()
                            except OSError:
                                result[attr] = None
                        return result
                    def __getattr__(self, item):
                        attr = getattr(self._svc, item)
                        if callable(attr):
                            def _safe(*a, **kw):
                                try:
                                    return attr(*a, **kw)
                                except OSError:
                                    return None
                            return _safe
                        return attr

                if _orig_svc_iter is not None:
                    def _safe_win_service_iter():
                        for svc in _orig_svc_iter():
                            try:
                                yield _SafeService(svc)
                            except OSError:
                                continue
                    _psutil_proxy.__dict__['win_service_iter'] = _safe_win_service_iter

                if _orig_svc_get is not None:
                    def _safe_win_service_get(name):
                        try:
                            return _SafeService(_orig_svc_get(name))
                        except OSError as _e:
                            raise OSError(f"Сервис '{name}' недоступен: {_e}") from _e
                    _psutil_proxy.__dict__['win_service_get'] = _safe_win_service_get

                safe_globals['psutil'] = _psutil_proxy

                # Перехватываем `import psutil` внутри sandbox — возвращаем прокси,
                # а не оригинальный модуль из builtins.__import__
                _orig_import = safe_globals['__builtins__']['__import__']
                def _patched_import(name, glb=None, loc=None, fromlist=(), level=0):
                    if name == 'psutil' or name.startswith('psutil.'):
                        return _psutil_proxy
                    return _orig_import(name, glb, loc, fromlist, level)
                safe_globals['__builtins__']['__import__'] = _patched_import

            except ImportError:
                pass

            # Принудительно устанавливаем Agg-бэкенд если код использует matplotlib
            # (предотвращает "Starting a Matplotlib GUI outside of the main thread")
            if 'matplotlib' in code:
                try:
                    import matplotlib
                    matplotlib.use('Agg')
                except Exception:
                    pass

            try:
                compiled = compile(code, '<agent>', 'exec')
            except IndentationError:
                # LLM нередко генерирует: строка 1 без отступа, строки 2+ с отступом.
                # textwrap.dedent не помогает (min_indent=0), поэтому убираем
                # минимальный отступ строк 2+ и повторяем попытку.
                lines = code.splitlines()
                rest = [ln for ln in lines[1:] if ln.strip()]
                min_rest = min((len(ln) - len(ln.lstrip()) for ln in rest), default=0)
                if min_rest > 0:
                    fixed_lines = [lines[0]] + [
                        ln[min_rest:] if len(ln) >= min_rest and ln[:min_rest] == ' ' * min_rest
                        else ln
                        for ln in lines[1:]
                    ]
                    code = '\n'.join(fixed_lines)
                compiled = compile(code, '<agent>', 'exec')

            exec(compiled, safe_globals)  # pylint: disable=exec-used

            # Обновляем namespace (без __builtins__)
            self._namespace = {k: v for k, v in safe_globals.items()
                               if k != '__builtins__'}
            output = stdout_capture.getvalue()
            result = self._namespace.get('__result__')
            return {'output': output, 'result': result, 'success': True}
        except Exception as e:
            import traceback as _traceback
            import re as _re
            _tb = _traceback.format_exc()
            _agent_files = _re.findall(r'File "([^"]+\.py)"', _tb)
            _inner = [
                f for f in _agent_files
                if self._working_dir in f and '<agent>' not in f and '<string>' not in f
            ]
            _hint = f' [FILE:{_inner[-1]}]' if _inner else ''
            return {'error': str(e) + _hint, 'output': stdout_capture.getvalue(), 'success': False}
        finally:
            sys.stdout = old_stdout


class SearchTool(BaseTool):
    """
    Search Engine — поиск информации через поисковые API.
    Адаптер: подключается к любому поисковому backend.
    """

    def __init__(self, backend=None):
        """
        Args:
            backend — объект с методом search(query, **kwargs) → list[dict]
                      Если None, возвращает заглушку.
        """
        super().__init__('search', 'Поиск информации через поисковые системы')
        self.backend = backend

    def run(self, query: str, num_results: int = 5) -> dict:
        if not self.backend:
            return {
                'results': [],
                'warning': 'Search backend не подключён',
                'success': False,
            }
        try:
            results = self.backend.search(query, num_results=num_results)
            return {'results': results, 'query': query, 'success': True}
        except Exception as e:
            return {'error': str(e), 'success': False}


class GitHubTool(BaseTool):
    """
    GitHub API — работа с репозиториями, issues, PR, файлами.
    """

    def __init__(self, token: str | None = None, backend=None):
        """
        Args:
            token   — GitHub Personal Access Token
            backend — готовый клиент с методами get/post/... (PyGithub и др.)
        """
        super().__init__('github', 'Работа с GitHub API: репозитории, issues, PR, файлы')
        self.token = token
        self.backend = backend

    def run(self, action: str, **kwargs) -> dict:
        """
        action: 'get_repo' | 'list_issues' | 'create_issue' | 'get_file' | 'create_pr'
        """
        if not self.backend:
            return {'warning': 'GitHub backend не подключён', 'success': False}
        try:
            method = getattr(self.backend, action, None)
            if not method:
                return {'error': f"Action '{action}' не поддерживается", 'success': False}
            result = method(**kwargs)
            return {'result': result, 'success': True}
        except Exception as e:
            return {'error': str(e), 'success': False}


class DockerTool(BaseTool):
    """
    Docker — управление контейнерами через CLI.
    """

    def __init__(self, terminal: TerminalTool | None = None):
        super().__init__('docker', 'Управление Docker-контейнерами')
        self.terminal = terminal or TerminalTool()

    def run(self, action: str, **kwargs) -> dict:
        """
        action: 'ps' | 'run' | 'stop' | 'rm' | 'build' | 'pull' | 'logs'
        """
        commands = {
            'ps':    'docker ps',
            'run':   f"docker run {kwargs.get('args', '')} {kwargs.get('image', '')}",
            'stop':  f"docker stop {kwargs.get('container', '')}",
            'rm':    f"docker rm {kwargs.get('container', '')}",
            'build': f"docker build {kwargs.get('context', '.')} -t {kwargs.get('tag', 'agent')}",
            'pull':  f"docker pull {kwargs.get('image', '')}",
            'logs':  f"docker logs {kwargs.get('container', '')}",
        }
        cmd = commands.get(action)
        if not cmd:
            return {'error': f"Action '{action}' не поддерживается", 'success': False}
        return self.terminal.run(cmd)


class DatabaseTool(BaseTool):
    """
    Database — выполнение SQL-запросов.
    Адаптер: подключается к любому DB backend (sqlite3, psycopg2, pymysql и др.).
    """

    def __init__(self, connection=None):
        super().__init__('database', 'Выполнение SQL-запросов к базам данных')
        self.connection = connection

    def run(self, query: str, params=None) -> dict:
        if not self.connection:
            return {'warning': 'DB connection не подключён', 'success': False}
        # SECURITY: блокируем опасные DDL/DCL операции от LLM
        _BLOCKED = ('DROP', 'TRUNCATE', 'DELETE', 'ALTER', 'CREATE',
                    'GRANT', 'REVOKE', 'ATTACH', 'DETACH', 'PRAGMA', '--', ';--')
        first_token = query.strip().upper().split()[0] if query.strip() else ''
        if first_token in _BLOCKED:
            return {'error': f'SQL операция {first_token!r} заблокирована (политика безопасности)', 'success': False}
        # Блокируем inline-комментарии и множественные операторы
        if '--' in query or query.count(';') > 1:
            return {'error': 'Запрос содержит запрещённые конструкции (--) или несколько операторов', 'success': False}
        try:
            cursor = self.connection.cursor()
            cursor.execute(query, params or [])
            if query.strip().upper().startswith('SELECT'):
                rows = cursor.fetchall()
                return {'rows': rows, 'success': True}
            else:
                self.connection.commit()
                return {'affected': cursor.rowcount, 'success': True}
        except Exception as e:
            return {'error': str(e), 'success': False}


class PackageManagerTool(BaseTool):
    """
    Package Manager — установка, обновление и управление пакетами.
    (Package & Dependency Manager, Слой 13)
    """

    def __init__(self, manager: str = 'pip', terminal: TerminalTool | None = None):
        super().__init__('package_manager', f'Управление пакетами через {manager}')
        self.manager = manager
        self.terminal = terminal or TerminalTool()

    def run(self, action: str, package: str = '', **kwargs) -> dict:
        """
        action: 'install' | 'uninstall' | 'upgrade' | 'list' | 'show'
        """
        if self.manager == 'pip':
            commands = {
                'install':   f"pip install {package}",
                'uninstall': f"pip uninstall -y {package}",
                'upgrade':   f"pip install --upgrade {package}",
                'list':      "pip list --format=json",
                'show':      f"pip show {package}",
            }
        elif self.manager == 'npm':
            commands = {
                'install':   f"npm install {package}",
                'uninstall': f"npm uninstall {package}",
                'upgrade':   f"npm update {package}",
                'list':      "npm list --json",
                'show':      f"npm show {package}",
            }
        else:
            return {'error': f"Package manager '{self.manager}' не поддерживается", 'success': False}

        cmd = commands.get(action)
        if not cmd:
            return {'error': f"Action '{action}' не поддерживается", 'success': False}
        return self.terminal.run(cmd)


class CloudAPITool(BaseTool):
    """
    Cloud API — адаптер для облачных провайдеров (OpenAI, Anthropic, AWS, GCP и др.).
    Подключается к любому backend с методом call(endpoint, payload).
    """

    def __init__(self, provider: str = 'generic', client=None, api_key: str | None = None):
        super().__init__('cloud_api', f'Вызовы облачных API ({provider})')
        self.provider = provider
        self.client = client
        self.api_key = api_key

    def run(self, endpoint: str, payload: dict | None = None, method: str = 'POST') -> dict:
        if not self.client:
            return {'warning': f'Cloud client ({self.provider}) не подключён', 'success': False}
        try:
            result = self.client.call(endpoint, payload=payload or {}, method=method)
            return {'result': result, 'success': True}
        except Exception as e:
            return {'error': str(e), 'success': False}


# ── Группа 1: Управление компьютером ─────────────────────────────────────────

class ScreenshotTool(BaseTool):
    """
    Screenshot — снимок экрана.
    Использует mss (быстрый, мультиплатформенный) или PIL.ImageGrab как fallback.
    """

    def __init__(self):
        super().__init__('screenshot', 'Снимок экрана или отдельного монитора')

    def run(self, monitor: int = 0, save_path: str | None = None,
            output: str | None = None, path: str | None = None) -> dict:
        """
        Args:
            monitor   — 0=весь рабочий стол, 1..N=конкретный монитор
            save_path — путь для сохранения PNG; если None — tempfile
            output    — алиас для save_path
            path      — алиас для save_path
        Returns:
            {'path': str, 'size': (w,h), 'success': bool}
        """
        save_path = save_path or output or path
        import tempfile
        if not save_path:
            fd, save_path = tempfile.mkstemp(suffix='.png')
            os.close(fd)

        try:
            import mss
            import mss.tools
            with mss.mss() as sct:
                mon = sct.monitors[monitor] if monitor < len(sct.monitors) else sct.monitors[0]
                img = sct.grab(mon)
                mss.tools.to_png(img.rgb, img.size, output=save_path)
                return {'path': save_path, 'size': img.size, 'success': True}
        except ImportError:
            pass

        try:
            from PIL import ImageGrab  # pyright: ignore[reportMissingImports]
            img = ImageGrab.grab()
            img.save(save_path)
            return {'path': save_path, 'size': img.size, 'success': True}
        except Exception as e:
            return {'error': f'mss и PIL недоступны: {e}', 'success': False}


class MouseKeyboardTool(BaseTool):
    """
    Mouse & Keyboard — управление мышью и клавиатурой через pyautogui.
    FAILSAFE включён: отвести мышь в верхний левый угол = экстренная остановка.
    """

    def __init__(self):
        super().__init__('mouse_keyboard',
                         'Управление мышью (click/move/scroll) и клавиатурой (type/hotkey)')

    def run(self, action: str, **kwargs) -> dict:
        """
        action: 'click' | 'move' | 'type' | 'hotkey' | 'scroll' | 'position' | 'screenshot'
        click:      x, y, button='left'
        move:       x, y, duration=0.2
        type:       text, interval=0.05
        hotkey:     keys=['ctrl','c'] или keys='ctrl+c'
        scroll:     clicks=3, x=None, y=None
        position:   — возвращает текущие x,y мыши
        screenshot: — снимок экрана через pyautogui
        """
        try:
            import pyautogui
            pyautogui.FAILSAFE = True

            if action == 'click':
                x, y = kwargs.get('x'), kwargs.get('y')
                btn = kwargs.get('button', 'left')
                if x is not None and y is not None:
                    pyautogui.click(x, y, button=btn)
                else:
                    pyautogui.click(button=btn)
                return {'success': True, 'x': x, 'y': y, 'button': btn}

            elif action == 'double_click':
                x, y = kwargs.get('x'), kwargs.get('y')
                if x is not None and y is not None:
                    pyautogui.doubleClick(x, y)
                else:
                    pyautogui.doubleClick()
                return {'success': True, 'x': x, 'y': y}

            elif action == 'move':
                x, y = kwargs['x'], kwargs['y']
                pyautogui.moveTo(x, y, duration=kwargs.get('duration', 0.2))
                return {'success': True, 'x': x, 'y': y}

            elif action == 'drag':
                pyautogui.dragTo(kwargs['x'], kwargs['y'],
                                 duration=kwargs.get('duration', 0.3),
                                 button=kwargs.get('button', 'left'))
                return {'success': True}

            elif action == 'type':
                text = kwargs['text']
                pyautogui.typewrite(text, interval=kwargs.get('interval', 0.05))
                return {'success': True, 'length': len(text)}

            elif action == 'hotkey':
                keys = kwargs.get('keys', [])
                if isinstance(keys, str):
                    keys = [k.strip() for k in keys.split('+')]
                pyautogui.hotkey(*keys)
                return {'success': True, 'keys': keys}

            elif action == 'key':
                pyautogui.press(kwargs['key'])
                return {'success': True, 'key': kwargs['key']}

            elif action == 'scroll':
                x, y = kwargs.get('x'), kwargs.get('y')
                clicks = kwargs.get('clicks', 3)
                if x is not None and y is not None:
                    pyautogui.scroll(clicks, x=x, y=y)
                else:
                    pyautogui.scroll(clicks)
                return {'success': True, 'clicks': clicks}

            elif action == 'position':
                pos = pyautogui.position()
                return {'success': True, 'x': pos.x, 'y': pos.y}

            elif action == 'screenshot':
                import tempfile
                fd, path = tempfile.mkstemp(suffix='.png')
                os.close(fd)
                pyautogui.screenshot(path)
                return {'success': True, 'path': path}

            else:
                return {'error': f"Неизвестный action: {action}", 'success': False}

        except ImportError:
            return {'error': 'pyautogui не установлен: pip install pyautogui', 'success': False}
        except Exception as e:
            return {'error': str(e), 'success': False}


class ClipboardTool(BaseTool):
    """Clipboard — чтение и запись буфера обмена."""

    def __init__(self):
        super().__init__('clipboard', 'Чтение и запись буфера обмена')

    def run(self, action: str = 'read', text: str | None = None) -> dict:
        """
        action: 'read' | 'write'
        text — обязателен для action='write'
        """
        try:
            import pyperclip
            if action == 'read':
                content = pyperclip.paste()
                return {'content': content, 'length': len(content), 'success': True}
            elif action == 'write':
                if text is None:
                    return {'error': 'text обязателен', 'success': False}
                pyperclip.copy(text)
                return {'success': True, 'length': len(text)}
            return {'error': f"Неизвестный action: {action}", 'success': False}
        except ImportError:
            pass

        # Fallback — tkinter (встроен в Python)
        try:
            import tkinter as tk
            root = tk.Tk()
            root.withdraw()
            if action == 'read':
                content = root.clipboard_get()
                root.destroy()
                return {'content': content, 'success': True}
            elif action == 'write':
                root.clipboard_clear()
                root.clipboard_append(text or '')
                root.update()
                root.destroy()
                return {'success': True}
            root.destroy()
        except Exception as e:
            return {'error': str(e), 'success': False}
        return {'error': f"Неизвестный action: {action}", 'success': False}


class WindowManagerTool(BaseTool):
    """Window Manager — список, фокус, закрытие, минимизация окон."""

    def __init__(self):
        super().__init__('window_manager',
                         'Управление окнами: список, фокус, закрыть, свернуть/развернуть')

    def run(self, action: str, title: str | None = None, **kwargs) -> dict:
        """
        action: 'list' | 'active' | 'focus' | 'close' | 'minimize' | 'maximize'
        title — (часть) заголовка окна для focus/close/minimize/maximize
        """
        try:
            import pygetwindow as gw

            if action == 'list':
                titles = [t for t in gw.getAllTitles() if t.strip()]
                return {'windows': titles, 'count': len(titles), 'success': True}

            elif action == 'active':
                win = gw.getActiveWindow()
                if win:
                    return {'title': win.title, 'left': win.left, 'top': win.top,
                            'width': win.width, 'height': win.height, 'success': True}
                return {'title': None, 'success': True}

            elif action in ('focus', 'close', 'minimize', 'maximize'):
                wins = gw.getWindowsWithTitle(title or '')
                if not wins:
                    return {'error': f"Окно '{title}' не найдено", 'success': False}
                w = wins[0]
                if action == 'focus':
                    w.activate()
                elif action == 'close':
                    w.close()
                elif action == 'minimize':
                    w.minimize()
                elif action == 'maximize':
                    w.maximize()
                return {'success': True, 'title': w.title}

            return {'error': f"Неизвестный action: {action}", 'success': False}

        except ImportError:
            return {'error': 'pygetwindow не установлен: pip install pygetwindow',
                    'success': False}
        except Exception as e:
            return {'error': str(e), 'success': False}


# ── Группа 2: Системные инструменты ──────────────────────────────────────────

class ProcessManagerTool(BaseTool):
    """
    Process Manager — управление процессами через psutil.
    Действия: list, find, kill, info, system.
    """

    def __init__(self):
        super().__init__('process_manager',
                         'Процессы: список, поиск, завершение, системные метрики CPU/RAM')

    def run(self, action: str, **kwargs) -> dict:
        try:
            import psutil

            if action == 'list':
                procs = []
                for p in psutil.process_iter(
                        ['pid', 'name', 'cpu_percent', 'memory_percent', 'status']):
                    try:
                        procs.append(p.info)
                    except (psutil.NoSuchProcess, psutil.AccessDenied):
                        pass
                return {'processes': procs[:100], 'total': len(procs), 'success': True}

            elif action == 'find':
                name = (kwargs.get('name') or '').lower()
                found = []
                for p in psutil.process_iter(['pid', 'name', 'status', 'cmdline']):
                    try:
                        if name in (p.info['name'] or '').lower():
                            found.append(p.info)
                    except (psutil.NoSuchProcess, psutil.AccessDenied):
                        pass
                return {'processes': found, 'count': len(found), 'success': True}

            elif action == 'kill':
                pid  = kwargs.get('pid')
                name = kwargs.get('name')
                killed = []
                if pid:
                    try:
                        psutil.Process(int(pid)).terminate()
                        killed.append(pid)
                    except Exception as e:
                        return {'error': str(e), 'success': False}
                elif name:
                    for p in psutil.process_iter(['pid', 'name']):
                        try:
                            if name.lower() in (p.info['name'] or '').lower():
                                p.terminate()
                                killed.append(p.info['pid'])
                        except (psutil.NoSuchProcess, psutil.AccessDenied):
                            pass
                return {'killed': killed, 'count': len(killed), 'success': bool(killed)}

            elif action == 'info':
                pid = kwargs.get('pid')
                if not pid:
                    return {'error': 'pid обязателен', 'success': False}
                p = psutil.Process(int(pid))
                return {
                    'pid': p.pid, 'name': p.name(), 'status': p.status(),
                    'cpu_percent': p.cpu_percent(interval=0.1),
                    'memory_mb': round(p.memory_info().rss / 1024 / 1024, 1),
                    'cmdline': p.cmdline(),
                    'success': True,
                }

            elif action == 'system':
                vm = psutil.virtual_memory()
                if os.name == 'nt':
                    drive = os.path.splitdrive(os.getcwd())[0] or 'C:'
                    disk_path = f"{drive}\\"
                else:
                    disk_path = '/'
                du = psutil.disk_usage(disk_path)
                return {
                    'cpu_percent':        psutil.cpu_percent(interval=1),
                    'memory_percent':     vm.percent,
                    'memory_available_mb': round(vm.available / 1024 / 1024, 1),
                    'disk_percent':       du.percent,
                    'disk_free_gb':       round(du.free / 1024 ** 3, 1),
                    'success': True,
                }

            return {'error': f"Неизвестный action: {action}", 'success': False}

        except ImportError:
            return {'error': 'psutil не установлен: pip install psutil', 'success': False}
        except Exception as e:
            return {'error': str(e), 'success': False}


class GitTool(BaseTool):
    """
    Git — полноценные git-операции через subprocess.
    Работает напрямую с git CLI, не требует GitHub токена.
    """

    def __init__(self, working_dir: str | None = None):
        super().__init__('git', 'Git: status, log, diff, add, commit, push, pull, branch, clone')
        self.working_dir = working_dir or os.getcwd()

    def run(self, action: str, **kwargs) -> dict:
        """
        action: status | log | diff | add | commit | push | pull |
                checkout | branch | stash | clone | init | reset
        """
        action = (action or '').strip().lower()
        cmd: list[str] | None = None

        if action == 'status':
            cmd = ['git', 'status', '--porcelain']
        elif action == 'log':
            n = int(kwargs.get('n', 10) or 10)
            cmd = ['git', 'log', '--oneline', f'-{max(1, n)}']
        elif action == 'diff':
            path = str(kwargs.get('path', '') or '').strip()
            cmd = ['git', 'diff'] + ([path] if path else [])
        elif action == 'diff_staged':
            cmd = ['git', 'diff', '--cached']
        elif action == 'add':
            path = str(kwargs.get('path', '.') or '.').strip()
            cmd = ['git', 'add', path]
        elif action == 'commit':
            message = str(kwargs.get('message', 'update') or 'update')
            cmd = ['git', 'commit', '-m', message]
        elif action == 'push':
            cmd = ['git', 'push', str(kwargs.get('remote', 'origin')), str(kwargs.get('branch', 'HEAD'))]
        elif action == 'pull':
            remote = str(kwargs.get('remote', 'origin'))
            branch = str(kwargs.get('branch', '') or '').strip()
            cmd = ['git', 'pull', remote] + ([branch] if branch else [])
        elif action == 'fetch':
            cmd = ['git', 'fetch', str(kwargs.get('remote', 'origin'))]
        elif action == 'checkout':
            branch = str(kwargs.get('branch', '') or '').strip()
            if not branch:
                return {'error': 'branch обязателен для checkout', 'success': False}
            cmd = ['git', 'checkout', branch]
        elif action == 'branch':
            cmd = ['git', 'branch', '-a']
        elif action == 'new_branch':
            branch = str(kwargs.get('branch', 'new-branch') or 'new-branch').strip()
            cmd = ['git', 'checkout', '-b', branch]
        elif action == 'stash':
            subcommand = str(kwargs.get('subcommand', '') or '').strip()
            cmd = ['git', 'stash'] + ([subcommand] if subcommand else [])
        elif action == 'clone':
            url = str(kwargs.get('url', '') or '').strip()
            dest = str(kwargs.get('dest', '') or '').strip()
            if not url:
                return {'error': 'url обязателен для clone', 'success': False}
            cmd = ['git', 'clone', url] + ([dest] if dest else [])
        elif action == 'init':
            cmd = ['git', 'init']
        elif action == 'reset':
            mode = str(kwargs.get('mode', '--soft') or '--soft').strip()
            ref = str(kwargs.get('ref', 'HEAD~1') or 'HEAD~1').strip()
            cmd = ['git', 'reset', mode, ref]
        elif action == 'show':
            ref = str(kwargs.get('ref', 'HEAD') or 'HEAD').strip()
            cmd = ['git', 'show', ref, '--stat']

        if not cmd:
            return {'error': f"Action '{action}' не поддерживается", 'success': False}

        try:
            proc = subprocess.run(
                cmd, shell=False, capture_output=True, text=True,
                encoding='utf-8', errors='replace',
                timeout=60, cwd=self.working_dir, check=False,
            )
            return {
                'stdout':     proc.stdout.strip(),
                'stderr':     proc.stderr.strip(),
                'returncode': proc.returncode,
                'success':    proc.returncode == 0,
            }
        except Exception as e:
            return {'error': str(e), 'success': False}


class NetworkTool(BaseTool):
    """
    Network — HTTP-запросы, ping, проверка портов, DNS-разрешение.
    """

    def __init__(self):
        super().__init__('network', 'Сеть: HTTP-запросы, ping, проверка портов, DNS')

    def run(self, action: str, **kwargs) -> dict:
        """
        action: 'http' | 'ping' | 'port_check' | 'dns'
        http:       url, method='GET', headers={}, body=None, timeout=15
        ping:       host, count=3
        port_check: host, port, timeout=3.0
        dns:        hostname
        """
        if action == 'http':
            return self._http(**kwargs)
        elif action == 'ping':
            return self._ping(**kwargs)
        elif action == 'port_check':
            return self._port_check(**kwargs)
        elif action == 'dns':
            return self._dns(**kwargs)
        return {'error': f"Неизвестный action: {action}", 'success': False}

    def _http(self, url: str, method: str = 'GET', headers: dict | None = None,
              body=None, params: dict | None = None, timeout: int = 15) -> dict:
        try:
            import requests as rq
            r = rq.request(
                method.upper(), url,
                headers=headers or {},
                json=body if isinstance(body, dict) else None,
                data=body if isinstance(body, str) else None,
                params=params or {},
                timeout=timeout,
            )
            try:
                content = r.json()
            except Exception:
                content = r.text[:5000]
            return {
                'status_code': r.status_code,
                'content':     content,
                'headers':     dict(r.headers),
                'success':     r.status_code < 400,
            }
        except Exception as e:
            return {'error': str(e), 'success': False}

    def _ping(self, host: str, count: int = 3) -> dict:
        import platform
        # SECURITY: validate host — only hostname/IP, no injection
        host = str(host).strip()
        if not re.match(r'^[a-zA-Z0-9.\-:]+$', host) or len(host) > 253:
            return {'error': f'Invalid host: {host[:50]}', 'success': False}
        count = max(1, min(int(count), 10))
        flag = '-n' if platform.system() == 'Windows' else '-c'
        try:
            from execution.command_gateway import CommandGateway
            gw = CommandGateway.get_instance()
            r = gw.execute(
                ['ping', flag, str(count), host],
                timeout=15, caller='NetworkTool._ping',
                allow_network=True,
            )
            return {
                'stdout':    r.stdout,
                'reachable': r.returncode == 0,
                'success':   r.allowed,
                'error':     r.reject_reason if not r.allowed else None,
            }
        except Exception as e:
            return {'error': str(e), 'success': False}

    def _port_check(self, host: str, port: int, timeout: float = 3.0) -> dict:
        try:
            sock = socket.create_connection((host, int(port)), timeout=timeout)
            sock.close()
            return {'host': host, 'port': port, 'open': True, 'success': True}
        except (socket.timeout, ConnectionRefusedError, OSError):
            return {'host': host, 'port': port, 'open': False, 'success': True}
        except Exception as e:
            return {'error': str(e), 'success': False}

    def _dns(self, hostname: str) -> dict:
        try:
            ip = socket.gethostbyname(hostname)
            return {'hostname': hostname, 'ip': ip, 'success': True}
        except Exception as e:
            return {'error': str(e), 'success': False}


class CronSchedulerTool(BaseTool):
    """
    Cron Scheduler — периодический запуск задач по интервалу (в памяти процесса).
    """

    def __init__(self):
        super().__init__('cron_scheduler', 'Планировщик: запуск функций по интервалу в секундах')
        self._jobs: dict[str, dict] = {}

    def run(self, action: str, **kwargs) -> dict:
        """
        action: 'schedule' | 'cancel' | 'list'
        schedule: job_id, interval_seconds, callback (callable), **callback_kwargs
        cancel:   job_id
        list:     —
        """
        if action == 'schedule':
            return self._schedule(
                kwargs['job_id'],
                float(kwargs['interval_seconds']),
                kwargs['callback'],
                **{k: v for k, v in kwargs.items()
                   if k not in ('job_id', 'interval_seconds', 'callback')},
            )
        elif action == 'cancel':
            return self._cancel(kwargs.get('job_id', ''))
        elif action == 'list':
            return {
                'jobs': [
                    {'job_id': jid, 'interval': j['interval']}
                    for jid, j in self._jobs.items()
                ],
                'success': True,
            }
        return {'error': f"Неизвестный action: {action}", 'success': False}

    def _schedule(self, job_id: str, interval: float, callback, **cb_kwargs) -> dict:
        self._cancel(job_id)
        self._jobs[job_id] = {'interval': interval, 'active': True}

        def _loop():
            while self._jobs.get(job_id, {}).get('active'):
                try:
                    callback(**cb_kwargs)
                except Exception:
                    pass
                time.sleep(interval)

        t = threading.Thread(target=_loop, daemon=True, name=f'cron_{job_id}')
        t.start()
        return {'job_id': job_id, 'interval': interval, 'success': True}

    def _cancel(self, job_id: str) -> dict:
        if job_id in self._jobs:
            self._jobs[job_id]['active'] = False
            del self._jobs[job_id]
            return {'cancelled': job_id, 'success': True}
        return {'error': f"Job '{job_id}' не найден", 'success': False}


class HTTPClientTool(BaseTool):
    """
    HTTP Client — произвольные HTTP-запросы к любым API и сайтам.
    """

    def __init__(self):
        super().__init__('http_client', 'Произвольные HTTP-запросы: GET/POST/PUT/DELETE + headers')

    def run(self, url: str, method: str = 'GET', headers: dict | None = None,
            body=None, params: dict | None = None, timeout: int = 15) -> dict:
        # Переиспользуем реализацию NetworkTool, чтобы избежать дублирования.
        network = NetworkTool()
        return network._http(
            url=url,
            method=method,
            headers=headers,
            body=body,
            params=params,
            timeout=timeout,
        )


# ── Группа 3: AI-инструменты ──────────────────────────────────────────────────

class ImageGeneratorTool(BaseTool):
    """
    Image Generator — генерация изображений через DALL-E 3.
    Опционально скачивает результат в файл.
    """

    def __init__(self, openai_client=None):
        super().__init__('image_generator', 'Генерация изображений через DALL-E 3')
        self.client = openai_client

    def run(self, prompt: str, size: str = '1024x1024',
            quality: str = 'standard', save_path: str | None = None) -> dict:
        """
        size:    '1024x1024' | '1792x1024' | '1024x1792'
        quality: 'standard' | 'hd'
        """
        if not self.client:
            return {'error': 'OpenAI client не подключён', 'success': False}
        oa = getattr(self.client, '_client', None)
        if not oa:
            return {'error': 'openai._client недоступен', 'success': False}
        try:
            response = oa.images.generate(
                model='dall-e-3', prompt=prompt,
                size=size, quality=quality, n=1,
            )
            url = response.data[0].url
            revised = getattr(response.data[0], 'revised_prompt', None)

            if save_path:
                import requests as rq
                r = rq.get(url, timeout=30)
                with open(save_path, 'wb') as f:
                    f.write(r.content)

            return {'url': url, 'revised_prompt': revised,
                    'save_path': save_path, 'success': True}
        except Exception as e:
            return {'error': str(e), 'success': False}


class OCRTool(BaseTool):
    """
    OCR — извлечение текста с изображений.
    Приоритет: GPT-4o Vision (точнее, не требует установки).
    Fallback: pytesseract (локально, нужна установка tesseract).
    """

    def __init__(self, openai_client=None):
        super().__init__('ocr', 'OCR: извлечение текста с изображений (GPT-4o Vision / tesseract)')
        self.client = openai_client

    def run(self, path: str, language: str = 'rus+eng') -> dict:
        if self.client:
            result = self._ocr_vision(path)
            if result.get('success'):
                return result

        return self._ocr_tesseract(path, language)

    def _ocr_vision(self, path: str) -> dict:
        oa = getattr(self.client, '_client', None)
        if not oa:
            return {'success': False, 'error': 'no oa client'}
        try:
            import base64
            ext = os.path.splitext(path)[1].lstrip('.') or 'png'
            with open(path, 'rb') as f:
                b64 = base64.b64encode(f.read()).decode()
            response = oa.chat.completions.create(
                model='gpt-4o',
                messages=[{
                    'role': 'user',
                    'content': [
                        {'type': 'text',
                         'text': 'Извлеки весь текст с изображения. Верни только текст.'},
                        {'type': 'image_url',
                         'image_url': {'url': f'data:image/{ext};base64,{b64}'}},
                    ],
                }],
                max_tokens=2000,
            )
            text = response.choices[0].message.content or ''
            return {'text': text, 'method': 'gpt4o_vision', 'success': True}
        except Exception as e:
            return {'error': str(e), 'success': False}

    def _ocr_tesseract(self, path: str, language: str) -> dict:
        try:
            import pytesseract
            from PIL import Image
            text = pytesseract.image_to_string(Image.open(path), lang=language)
            return {'text': text.strip(), 'method': 'tesseract', 'success': True}
        except ImportError:
            return {'error': 'Требуется: pip install pytesseract pillow + tesseract-ocr',
                    'success': False}
        except Exception as e:
            return {'error': str(e), 'success': False}


class TranslateTool(BaseTool):
    """
    Translate — перевод текста через LLM (любой язык).
    Не требует отдельного API — использует уже подключённый OpenAIClient.
    """

    def __init__(self, openai_client=None):
        super().__init__('translate', 'Перевод текста на любой язык через LLM')
        self.client = openai_client

    def run(self, text: str, target: str = 'ru', source: str = 'auto') -> dict:
        """
        text   — текст для перевода
        target — целевой язык (ru / en / de / zh / ...)
        source — исходный язык или 'auto'
        Приоритет: OpenAI LLM → deep_translator (offline fallback)
        """
        # Попытка 1: через OpenAI LLM (лучшее качество)
        if self.client:
            try:
                src_hint = f'с языка {source} ' if source != 'auto' else ''
                prompt = (
                    f'Переведи {src_hint}на {target}. '
                    f'Верни только перевод, без пояснений:\n\n{text}'
                )
                translated = self.client.infer(
                    prompt,
                    system='Ты профессиональный переводчик. Переводи точно и естественно.',
                )
                return {
                    'original':   text[:300],
                    'translated': translated,
                    'target':     target,
                    'backend':    'openai',
                    'success':    True,
                }
            except Exception:
                pass  # Fallback ниже

        # Попытка 2: deep_translator (бесплатно, без ключей)
        try:
            from deep_translator import GoogleTranslator
            src = source if source != 'auto' else 'auto'
            translated = GoogleTranslator(source=src, target=target).translate(text)
            return {
                'original':   text[:300],
                'translated': translated,
                'target':     target,
                'backend':    'deep_translator',
                'success':    True,
            }
        except Exception as e:
            return {'error': f'Перевод недоступен (OpenAI и deep_translator): {e}', 'success': False}


# ── Группа 4: Данные, документы, визуализация ────────────────────────────────

class SpreadsheetTool(BaseTool):
    """
    Spreadsheet — работа с Excel (.xlsx) и CSV файлами через openpyxl + pandas.

    Действия:
        read        — читает файл, возвращает строки как список словарей
        write       — создаёт/перезаписывает файл из списка словарей
        append      — добавляет строки в существующий файл
        summary     — статистика: shape, columns, dtypes, describe
        filter      — фильтрует строки по условию (column, op, value)
        sort        — сортировка по колонке
        csv_to_xlsx — конвертирует CSV в XLSX
        xlsx_to_csv — конвертирует XLSX в CSV
    """

    def __init__(self):
        super().__init__('spreadsheet',
                         'Excel/CSV: чтение, запись, фильтрация, сортировка, конвертация')

    def run(self, action: str, path: str, **kwargs) -> dict:
        try:
            import pandas as pd

            if action == 'read':
                df = self._load(pd, path, kwargs.get('sheet', 0))
                rows = df.head(kwargs.get('limit', 100)).to_dict(orient='records')
                return {
                    'rows':    rows,
                    'total':   len(df),
                    'columns': list(df.columns),
                    'success': True,
                }

            elif action == 'write':
                rows = kwargs.get('rows', [])
                df = pd.DataFrame(rows)
                self._save(df, path, kwargs.get('sheet', 'Sheet1'))
                return {'path': path, 'rows_written': len(df), 'success': True}

            elif action == 'append':
                rows = kwargs.get('rows', [])
                new_df = pd.DataFrame(rows)
                if os.path.exists(path):
                    existing = self._load(pd, path, kwargs.get('sheet', 0))
                    df = pd.concat([existing, new_df], ignore_index=True)
                else:
                    df = new_df
                self._save(df, path, kwargs.get('sheet', 'Sheet1'))
                return {'path': path, 'total_rows': len(df), 'success': True}

            elif action == 'summary':
                df = self._load(pd, path, kwargs.get('sheet', 0))
                return {
                    'shape':   list(df.shape),
                    'columns': list(df.columns),
                    'dtypes':  {c: str(t) for c, t in df.dtypes.items()},
                    'describe': df.describe(include='all').to_dict(),
                    'success': True,
                }

            elif action == 'filter':
                df = self._load(pd, path, kwargs.get('sheet', 0))
                col = kwargs['column']
                op  = kwargs.get('op', '==')
                val = kwargs['value']
                ops = {
                    '==': df[col] == val, '!=': df[col] != val,
                    '>':  df[col] >  val, '>=': df[col] >= val,
                    '<':  df[col] <  val, '<=': df[col] <= val,
                    'contains': df[col].astype(str).str.contains(str(val), na=False),
                }
                mask = ops.get(op)
                if mask is None:
                    return {'error': f"Неизвестный оператор: {op}", 'success': False}
                filtered = df[mask]
                return {
                    'rows':    filtered.head(kwargs.get('limit', 100)).to_dict(orient='records'),
                    'total':   len(filtered),
                    'success': True,
                }

            elif action == 'sort':
                df = self._load(pd, path, kwargs.get('sheet', 0))
                col = kwargs['column']
                asc = kwargs.get('ascending', True)
                df = df.sort_values(col, ascending=asc)
                out = kwargs.get('output', path)
                self._save(df, out, kwargs.get('sheet', 'Sheet1'))
                return {'path': out, 'rows': len(df), 'success': True}

            elif action == 'csv_to_xlsx':
                df = pd.read_csv(path)
                out = kwargs.get('output', path.replace('.csv', '.xlsx'))
                df.to_excel(out, index=False)
                return {'path': out, 'rows': len(df), 'success': True}

            elif action == 'xlsx_to_csv':
                df = pd.read_excel(path, sheet_name=kwargs.get('sheet', 0))
                out = kwargs.get('output', path.replace('.xlsx', '.csv'))
                df.to_csv(out, index=False, encoding='utf-8-sig')
                return {'path': out, 'rows': len(df), 'success': True}

            else:
                return {'error': f"Неизвестный action: {action}", 'success': False}

        except ImportError as e:
            return {'error': f'Зависимости не установлены: pip install pandas openpyxl — {e}',
                    'success': False}
        except Exception as e:
            return {'error': str(e), 'success': False}

    @staticmethod
    def _load(pd, path: str, sheet):
        if path.endswith('.csv'):
            return pd.read_csv(path)
        return pd.read_excel(path, sheet_name=sheet)

    @staticmethod
    def _save(df, path: str, sheet: str):
        if path.endswith('.csv'):
            df.to_csv(path, index=False, encoding='utf-8-sig')
        else:
            df.to_excel(path, index=False, sheet_name=str(sheet))


class PDFGeneratorTool(BaseTool):
    """
    PDF Generator — создание PDF-документов через reportlab.

    Действия:
        create      — создаёт PDF из списка элементов (параграфы, таблицы, изображения)
        from_text   — простое создание PDF из строки текста
        from_html   — PDF из HTML (требует xhtml2pdf)
        merge       — объединяет несколько PDF в один (требует pypdf)
    """

    def __init__(self):
        super().__init__('pdf_generator',
                         'Создание PDF: из текста, таблиц, изображений, HTML; объединение файлов')

    def run(self, action: str, output: str, **kwargs) -> dict:
        try:
            if action == 'from_text':
                return self._from_text(output, kwargs.get('text', ''),
                                       kwargs.get('title', ''),
                                       kwargs.get('font_size', 12))

            elif action == 'create':
                return self._create(output, kwargs.get('elements', []),
                                    kwargs.get('title', ''),
                                    kwargs.get('page_size', 'A4'))

            elif action == 'from_html':
                return self._from_html(output, kwargs.get('html', ''))

            elif action == 'merge':
                return self._merge(output, kwargs.get('files', []))

            else:
                return {'error': f"Неизвестный action: {action}", 'success': False}

        except ImportError as e:
            return {'error': f'Установи зависимости: pip install reportlab — {e}',
                    'success': False}
        except Exception as e:
            return {'error': str(e), 'success': False}

    def _from_text(self, output: str, text: str,
                   title: str = '', _font_size: int = 12) -> dict:
        from reportlab.lib.pagesizes import A4  # pyright: ignore[reportMissingImports]
        from reportlab.lib.styles import getSampleStyleSheet  # pyright: ignore[reportMissingImports]
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer  # pyright: ignore[reportMissingImports]
        from reportlab.lib.units import cm  # pyright: ignore[reportMissingImports]

        doc = SimpleDocTemplate(output, pagesize=A4,
                                rightMargin=2*cm, leftMargin=2*cm,
                                topMargin=2*cm, bottomMargin=2*cm)
        styles = getSampleStyleSheet()
        story = []

        if title:
            story.append(Paragraph(title, styles['Title']))
            story.append(Spacer(1, 0.5*cm))

        # Разбиваем текст на параграфы
        for para in text.split('\n\n'):
            para = para.strip()
            if para:
                # Экранируем HTML-символы для reportlab
                para = para.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                story.append(Paragraph(para, styles['Normal']))
                story.append(Spacer(1, 0.3*cm))

        doc.build(story)
        return {'path': output, 'success': True}

    def _create(self, output: str, elements: list,
                title: str = '', page_size: str = 'A4') -> dict:
        """
        elements — список dict:
            {'type': 'text',  'content': '...', 'style': 'Normal'/'Heading1'/...}
            {'type': 'table', 'data': [[row1], [row2], ...], 'headers': [...]}
            {'type': 'image', 'path': '...', 'width': 400, 'height': 300}
            {'type': 'spacer','height': 0.5}
        """
        from reportlab.lib.pagesizes import A4, letter  # pyright: ignore[reportMissingImports]
        from reportlab.lib.styles import getSampleStyleSheet  # pyright: ignore[reportMissingImports]
        from reportlab.lib.units import cm  # pyright: ignore[reportMissingImports]
        from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                        Table, TableStyle, Image as RLImage)  # pyright: ignore[reportMissingImports]
        from reportlab.lib import colors  # pyright: ignore[reportMissingImports]

        psize = A4 if page_size.upper() == 'A4' else letter
        doc = SimpleDocTemplate(output, pagesize=psize,
                                rightMargin=2*cm, leftMargin=2*cm,
                                topMargin=2*cm, bottomMargin=2*cm)
        styles = getSampleStyleSheet()
        story = []

        if title:
            story.append(Paragraph(title, styles['Title']))
            story.append(Spacer(1, 0.5*cm))

        for el in elements:
            t = el.get('type', 'text')

            if t == 'text':
                content = str(el.get('content', '')).replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                style   = styles.get(el.get('style', 'Normal'), styles['Normal'])
                story.append(Paragraph(content, style))
                story.append(Spacer(1, 0.2*cm))

            elif t == 'table':
                data    = el.get('data', [])
                headers = el.get('headers')
                if headers:
                    data = [headers] + data
                if data:
                    tbl = Table(data)
                    tbl.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                        ('TEXTCOLOR',  (0, 0), (-1, 0), colors.whitesmoke),
                        ('FONTNAME',   (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('GRID',       (0, 0), (-1, -1), 0.5, colors.black),
                        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.lightgrey]),
                    ]))
                    story.append(tbl)
                    story.append(Spacer(1, 0.3*cm))

            elif t == 'image':
                img_path = el.get('path', '')
                if os.path.exists(img_path):
                    w = el.get('width',  400)
                    h = el.get('height', 300)
                    story.append(RLImage(img_path, width=w, height=h))
                    story.append(Spacer(1, 0.3*cm))

            elif t == 'spacer':
                story.append(Spacer(1, el.get('height', 0.5)*cm))

        doc.build(story)
        return {'path': output, 'success': True}

    def _from_html(self, output: str, html: str) -> dict:
        try:
            from xhtml2pdf import pisa  # pyright: ignore[reportMissingImports]
            with open(output, 'wb') as f:
                result = pisa.CreatePDF(html, dest=f)
            if result.err:  # type: ignore[union-attr]
                return {'error': 'Ошибка конвертации HTML в PDF', 'success': False}
            return {'path': output, 'success': True}
        except ImportError:
            return {'error': 'pip install xhtml2pdf', 'success': False}

    def _merge(self, output: str, files: list) -> dict:
        try:
            from pypdf import PdfWriter
            writer = PdfWriter()
            for f in files:
                writer.append(f)
            with open(output, 'wb') as out_f:
                writer.write(out_f)
            return {'path': output, 'pages': len(writer.pages), 'success': True}
        except ImportError:
            return {'error': 'pip install pypdf', 'success': False}


class PowerPointTool(BaseTool):
    """
    PowerPoint — создание и редактирование .pptx презентаций (python-pptx).

    Действия:
        create      — создать презентацию из списка слайдов
        add_slide   — добавить слайд в существующий файл
        from_text   — быстро создать презентацию из текста (каждый абзац = слайд)
        read        — прочитать текст из существующей презентации
    """

    def __init__(self):
        super().__init__('powerpoint', 'Создание и редактирование PowerPoint (.pptx): слайды, текст, таблицы')

    def run(self, action: str, output: str = '', **kwargs) -> dict:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN

            if action == 'create':
                return self._create(output, kwargs.get('slides', []),
                                    kwargs.get('title', ''),
                                    kwargs.get('theme', 'dark'))

            elif action == 'add_slide':
                return self._add_slide(output, kwargs.get('title', ''),
                                       kwargs.get('content', ''),
                                       kwargs.get('notes', ''))

            elif action == 'from_text':
                return self._from_text(output, kwargs.get('text', ''),
                                       kwargs.get('title', 'Presentation'))

            elif action == 'read':
                return self._read(kwargs.get('path', output))

            else:
                return {'error': f'Неизвестный action: {action}', 'success': False}

        except ImportError:
            return {'error': 'pip install python-pptx', 'success': False}
        except Exception as e:
            return {'error': str(e), 'success': False}

    def _create(self, output: str, slides: list, _title: str, theme: str) -> dict:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        dark = theme == 'dark'
        bg_color = RGBColor(0x0d, 0x11, 0x17) if dark else RGBColor(0xFF, 0xFF, 0xFF)
        title_color = RGBColor(0x58, 0xA6, 0xFF) if dark else RGBColor(0x1F, 0x3C, 0x88)
        text_color = RGBColor(0xC9, 0xD1, 0xD9) if dark else RGBColor(0x24, 0x29, 0x2F)

        blank_layout = prs.slide_layouts[6]  # blank

        for slide_data in slides:
            slide = prs.slides.add_slide(blank_layout)

            # Background
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = bg_color

            sl_title = slide_data.get('title', '')
            sl_content = slide_data.get('content', '')
            sl_notes = slide_data.get('notes', '')

            if sl_title:
                txb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.2))
                tf = txb.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = sl_title
                p.font.size = Pt(32)
                p.font.bold = True
                p.font.color.rgb = title_color

            if sl_content:
                txb2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.3))
                tf2 = txb2.text_frame
                tf2.word_wrap = True
                for i, line in enumerate(sl_content.split('\n')):
                    if i == 0:
                        p2 = tf2.paragraphs[0]
                    else:
                        p2 = tf2.add_paragraph()
                    p2.text = line
                    p2.font.size = Pt(18)
                    p2.font.color.rgb = text_color

            if sl_notes:
                _ntf = getattr(slide.notes_slide, 'notes_text_frame', None)
                if _ntf is not None:
                    _ntf.text = sl_notes

        prs.save(output)
        return {'path': output, 'slides': len(slides), 'success': True}

    def _from_text(self, output: str, text: str, title: str) -> dict:
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        slides = []
        for i, para in enumerate(paragraphs):
            lines = para.split('\n')
            slides.append({
                'title': lines[0] if lines else f'Slide {i+1}',
                'content': '\n'.join(lines[1:]) if len(lines) > 1 else '',
            })
        if not slides:
            slides = [{'title': title, 'content': text}]
        return self._create(output, slides, title, 'dark')

    def _add_slide(self, path: str, title: str, content: str, notes: str) -> dict:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        prs = Presentation(path)
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        if title:
            txb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.2))
            p = txb.text_frame.paragraphs[0]
            p.text = title
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x58, 0xA6, 0xFF)

        if content:
            txb2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.3))
            txb2.text_frame.word_wrap = True
            txb2.text_frame.paragraphs[0].text = content
            txb2.text_frame.paragraphs[0].font.size = Pt(18)
            txb2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xC9, 0xD1, 0xD9)

        if notes:
            _ntf = getattr(slide.notes_slide, 'notes_text_frame', None)
            if _ntf is not None:
                _ntf.text = notes

        prs.save(path)
        return {'path': path, 'total_slides': len(prs.slides), 'success': True}

    def _read(self, path: str) -> dict:
        from pptx import Presentation
        prs = Presentation(path)
        result = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    _tf = getattr(shape, 'text_frame', None)
                    if _tf is None:
                        continue
                    for para in _tf.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            result.append({'slide': i + 1, 'text': texts})
        return {'slides': result, 'total': len(prs.slides), 'success': True}


class DataVizTool(BaseTool):
    """
    Data Visualization — построение графиков через matplotlib и plotly.

    Действия (matplotlib — сохраняет PNG):
        bar         — столбчатая диаграмма
        line        — линейный график
        pie         — круговая диаграмма
        scatter     — точечный график
        histogram   — гистограмма
        heatmap     — тепловая карта

    Действия (plotly — сохраняет HTML):
        plotly_bar
        plotly_line
        plotly_scatter
        plotly_pie
    """

    def __init__(self):
        super().__init__('data_viz',
                         'Графики: bar, line, pie, scatter, histogram, heatmap (matplotlib/plotly)')

    def run(self, action: str, output: str, **kwargs) -> dict:
        """
        Общие kwargs для matplotlib-графиков:
            data     — {'label': [values], ...} или просто список values
            labels   — подписи по оси X (список)
            title    — заголовок
            xlabel   — подпись оси X
            ylabel   — подпись оси Y
            figsize  — (width_inch, height_inch), по умолчанию (10, 6)
        """
        try:
            if action.startswith('plotly_'):
                return self._plotly(action[7:], output, **kwargs)

            import matplotlib
            matplotlib.use('Agg')   # без GUI — рисуем в файл
            import matplotlib.pyplot as plt

            figsize = kwargs.get('figsize', (10, 6))
            fig, ax = plt.subplots(figsize=figsize)

            title  = kwargs.get('title', '')
            xlabel = kwargs.get('xlabel', '')
            ylabel = kwargs.get('ylabel', '')

            if action == 'bar':
                data   = kwargs['data']
                # Автоконвертация плоского словаря {'A': 10, 'B': 20} → labels + list
                if isinstance(data, dict) and data and not isinstance(next(iter(data.values())), (list, tuple)):
                    if not kwargs.get('labels'):
                        kwargs['labels'] = list(data.keys())
                    data = list(data.values())
                labels = kwargs.get('labels', list(range(len(list(data.values())[0]
                                                           if isinstance(data, dict)
                                                           else data))))
                if isinstance(data, dict):
                    x = range(len(labels))
                    width = 0.8 / max(len(data), 1)
                    for i, (name, values) in enumerate(data.items()):
                        offset = [xi + i * width for xi in x]
                        ax.bar(offset, values, width=width, label=name)
                    ax.set_xticks([xi + width * len(data) / 2 for xi in x])
                    ax.set_xticklabels(labels, rotation=45, ha='right')
                    ax.legend()
                else:
                    ax.bar(labels, data)

            elif action == 'line':
                data = kwargs['data']
                labels = kwargs.get('labels')
                if isinstance(data, dict):
                    for name, values in data.items():
                        x = labels if labels else range(len(values))
                        ax.plot(x, values, marker='o', label=name)
                    ax.legend()
                else:
                    x = labels if labels else range(len(data))
                    ax.plot(x, data, marker='o')
                if labels:
                    ax.set_xticks(range(len(labels)))
                    ax.set_xticklabels(labels, rotation=45, ha='right')

            elif action == 'pie':
                values = kwargs['data']
                labels = kwargs.get('labels', [str(i) for i in range(len(values))])
                ax.pie(values, labels=labels, autopct='%1.1f%%', startangle=90)
                ax.axis('equal')

            elif action == 'scatter':
                x = kwargs['x']
                y = kwargs['y']
                label = kwargs.get('label')
                ax.scatter(x, y, label=label, alpha=0.7)
                if label:
                    ax.legend()

            elif action == 'histogram':
                data = kwargs['data']
                bins = kwargs.get('bins', 20)
                if isinstance(data, dict):
                    for name, values in data.items():
                        ax.hist(values, bins=bins, alpha=0.6, label=name)
                    ax.legend()
                else:
                    ax.hist(data, bins=bins, edgecolor='black')

            elif action == 'heatmap':
                import numpy as np  # pylint: disable=unused-import
                matrix = kwargs['data']   # 2D список
                im = ax.imshow(matrix, cmap=kwargs.get('cmap', 'coolwarm'), aspect='auto')
                fig.colorbar(im, ax=ax)
                xlabels = kwargs.get('xlabels')
                ylabels = kwargs.get('ylabels')
                if xlabels:
                    ax.set_xticks(range(len(xlabels)))
                    ax.set_xticklabels(xlabels, rotation=45, ha='right')
                if ylabels:
                    ax.set_yticks(range(len(ylabels)))
                    ax.set_yticklabels(ylabels)

            else:
                plt.close(fig)
                return {'error': f"Неизвестный action: {action}", 'success': False}

            if title:
                ax.set_title(title)
            if xlabel:
                ax.set_xlabel(xlabel)
            if ylabel:
                ax.set_ylabel(ylabel)

            ax.grid(kwargs.get('grid', True), alpha=0.3)
            plt.tight_layout()
            plt.savefig(output, dpi=kwargs.get('dpi', 150), bbox_inches='tight')
            plt.close(fig)
            return {'path': output, 'success': True}

        except ImportError as e:
            return {'error': f'pip install matplotlib — {e}', 'success': False}
        except Exception as e:
            return {'error': str(e), 'success': False}

    def _plotly(self, chart_type: str, output: str, **kwargs) -> dict:
        try:
            import plotly.graph_objects as go
            import plotly.express as px  # pylint: disable=unused-import

            title = kwargs.get('title', '')

            if chart_type == 'bar':
                data = kwargs['data']
                labels = kwargs.get('labels', list(range(len(
                    list(data.values())[0] if isinstance(data, dict) else data
                ))))
                if isinstance(data, dict):
                    fig = go.Figure([
                        go.Bar(name=name, x=labels, y=values)
                        for name, values in data.items()
                    ])
                    fig.update_layout(barmode='group')
                else:
                    fig = go.Figure(go.Bar(x=labels, y=data))

            elif chart_type == 'line':
                data = kwargs['data']
                labels = kwargs.get('labels')
                if isinstance(data, dict):
                    fig = go.Figure([
                        go.Scatter(name=name,
                                   x=labels if labels else list(range(len(values))),
                                   y=values, mode='lines+markers')
                        for name, values in data.items()
                    ])
                else:
                    x = labels if labels else list(range(len(data)))
                    fig = go.Figure(go.Scatter(x=x, y=data, mode='lines+markers'))

            elif chart_type == 'scatter':
                fig = go.Figure(go.Scatter(
                    x=kwargs['x'], y=kwargs['y'],
                    mode='markers',
                    text=kwargs.get('labels'),
                ))

            elif chart_type == 'pie':
                fig = go.Figure(go.Pie(
                    values=kwargs['data'],
                    labels=kwargs.get('labels'),
                ))

            else:
                return {'error': f"plotly_{chart_type} не поддерживается", 'success': False}

            if title:
                fig.update_layout(title=title)

            fig.write_html(output)
            return {'path': output, 'success': True}

        except ImportError as e:
            return {'error': f'pip install plotly — {e}', 'success': False}
        except Exception as e:
            return {'error': str(e), 'success': False}


# ── Группа 5: Коммуникации ───────────────────────────────────────────────────

class EmailTool(BaseTool):
    """
    Email — отправка и чтение почты через SMTP / IMAP.

    Работает с Gmail, Яндекс, Mail.ru и любым другим ящиком.
    Для Gmail: включи двухфакторку → создай «Пароль приложения» в настройках аккаунта.

    Действия:
        send      — отправить письмо (кому, тема, текст, вложения)
        read      — прочитать входящие (N последних писем)
        search    — поиск писем по теме/отправителю
        unread    — количество непрочитанных
        reply     — ответить на письмо по message_id
        delete    — удалить письмо по message_id
    """

    def __init__(self, smtp_host='smtp.gmail.com', smtp_port=587,
                 imap_host='imap.gmail.com', imap_port=993,
                 username=None, password=None):
        super().__init__('email',
                         'Email: отправка (SMTP) и чтение (IMAP) — Gmail, Яндекс, Mail.ru')
        self.smtp_host = smtp_host
        self.smtp_port = smtp_port
        self.imap_host = imap_host
        self.imap_port = imap_port
        self.username  = username
        self.password  = (password or '').replace(' ', '') or None

    def run(self, action: str, **kwargs) -> dict:
        if not self.username or not self.password:
            return {
                'error': (
                    'Email не настроен. Укажи email_username и email_password '
                    'в agent.py (или .env). Для Gmail нужен «Пароль приложения».'
                ),
                'success': False,
            }

        if action == 'send':
            return self._send(**kwargs)
        elif action == 'read':
            return self._read(**kwargs)
        elif action == 'search':
            return self._search(**kwargs)
        elif action == 'unread':
            return self._unread()
        elif action == 'reply':
            return self._reply(**kwargs)
        elif action == 'delete':
            return self._delete(**kwargs)
        else:
            return {'error': f"Неизвестный action: {action}", 'success': False}

    # ── Отправка ─────────────────────────────────────────────────────────────

    def _send(self, to, subject: str = '', body: str = '',
              attachments: list | None = None, html: bool = False) -> dict:
        """
        to          — строка или список адресов
        attachments — список путей к файлам
        html        — True если body содержит HTML
        """
        import smtplib
        from email.mime.multipart import MIMEMultipart
        from email.mime.text import MIMEText
        from email.mime.base import MIMEBase
        from email import encoders

        recipients = [to] if isinstance(to, str) else list(to)

        msg = MIMEMultipart('alternative' if html else 'mixed')
        msg['From']    = self.username or ''
        msg['To']      = ', '.join(recipients)
        msg['Subject'] = subject

        content_type = 'html' if html else 'plain'
        msg.attach(MIMEText(body, content_type, 'utf-8'))

        # Прикрепляем файлы
        for path in (attachments or []):
            if os.path.exists(path):
                with open(path, 'rb') as f:
                    part = MIMEBase('application', 'octet-stream')
                    part.set_payload(f.read())
                encoders.encode_base64(part)
                part.add_header('Content-Disposition',
                                f'attachment; filename="{os.path.basename(path)}"')
                msg.attach(part)

        try:
            with smtplib.SMTP(self.smtp_host, self.smtp_port, timeout=15) as server:
                server.ehlo()
                server.starttls()
                server.login(self.username or '', self.password or '')
                server.sendmail(self.username or '', recipients, msg.as_string())
            return {
                'sent_to': recipients,
                'subject': subject,
                'success': True,
            }
        except Exception as e:
            return {'error': str(e), 'success': False}

    # ── Чтение ───────────────────────────────────────────────────────────────

    def _imap_connect(self):
        import imaplib
        mail = imaplib.IMAP4_SSL(self.imap_host, self.imap_port)
        mail.login(self.username or '', self.password or '')
        return mail

    def _read(self, folder: str = 'INBOX', limit: int = 10,
              unread_only: bool = False) -> dict:
        import email as email_lib  # pylint: disable=unused-import
        from email.header import decode_header  # pylint: disable=unused-import

        try:
            mail = self._imap_connect()
            mail.select(folder)

            criteria = '(UNSEEN)' if unread_only else 'ALL'
            _, data = mail.search(None, criteria)
            ids = data[0].split()
            ids = ids[-limit:]   # последние N

            messages = []
            for uid in reversed(ids):
                _, msg_data = mail.fetch(uid, '(RFC822)')  # type: ignore[assignment]
                msg = email_lib.message_from_bytes(msg_data[0][1])  # type: ignore[index]

                subject = self._decode_header(msg.get('Subject', ''))
                sender  = self._decode_header(msg.get('From', ''))
                date    = msg.get('Date', '')
                body    = self._extract_body(msg)

                messages.append({
                    'uid':     uid.decode(),
                    'from':    sender,
                    'subject': subject,
                    'date':    date,
                    'body':    body[:500],
                })

            mail.logout()
            return {'messages': messages, 'count': len(messages), 'success': True}
        except Exception as e:
            return {'error': str(e), 'success': False}

    def _search(self, subject: str | None = None, sender: str | None = None,
                limit: int = 10) -> dict:
        import email as email_lib

        criteria_parts = []
        if subject:
            criteria_parts.append(f'SUBJECT "{subject}"')
        if sender:
            criteria_parts.append(f'FROM "{sender}"')
        if not criteria_parts:
            return {'error': 'Укажи subject или sender', 'success': False}

        criteria = '(' + ' '.join(criteria_parts) + ')'

        try:
            mail = self._imap_connect()
            mail.select('INBOX')
            _, data = mail.search(None, criteria)
            ids = data[0].split()[-limit:]

            messages = []
            for uid in reversed(ids):
                _, msg_data = mail.fetch(uid, '(RFC822)')  # type: ignore[assignment]
                msg = email_lib.message_from_bytes(msg_data[0][1])  # type: ignore[index]
                messages.append({
                    'uid':     uid.decode(),
                    'from':    self._decode_header(msg.get('From', '')),
                    'subject': self._decode_header(msg.get('Subject', '')),
                    'date':    msg.get('Date', ''),
                    'body':    self._extract_body(msg)[:300],
                })

            mail.logout()
            return {'messages': messages, 'count': len(messages), 'success': True}
        except Exception as e:
            return {'error': str(e), 'success': False}

    def _unread(self) -> dict:
        try:
            mail = self._imap_connect()
            mail.select('INBOX')
            _, data = mail.search(None, '(UNSEEN)')
            count = len(data[0].split()) if data[0] else 0
            mail.logout()
            return {'unread': count, 'success': True}
        except Exception as e:
            return {'error': str(e), 'success': False}

    def _reply(self, uid: str, body: str) -> dict:
        """Отвечает на письмо по uid (из _read / _search)."""
        import email as email_lib

        try:
            mail = self._imap_connect()
            mail.select('INBOX')
            _, msg_data = mail.fetch(uid, '(RFC822)')  # type: ignore[assignment]
            orig = email_lib.message_from_bytes(msg_data[0][1])  # type: ignore[index]
            mail.logout()

            to      = orig.get('From', '')
            subject = 'Re: ' + self._decode_header(orig.get('Subject', ''))
            return self._send(to=to, subject=subject, body=body)
        except Exception as e:
            return {'error': str(e), 'success': False}

    def _delete(self, uid: str) -> dict:
        try:
            mail = self._imap_connect()
            mail.select('INBOX')
            mail.store(uid, '+FLAGS', '\\Deleted')
            mail.expunge()
            mail.logout()
            return {'deleted': uid, 'success': True}
        except Exception as e:
            return {'error': str(e), 'success': False}

    # ── Helpers ──────────────────────────────────────────────────────────────

    @staticmethod
    def _decode_header(value: str) -> str:
        from email.header import decode_header
        parts = []
        for raw, enc in decode_header(value or ''):
            if isinstance(raw, bytes):
                parts.append(raw.decode(enc or 'utf-8', errors='replace'))
            else:
                parts.append(str(raw))
        return ' '.join(parts)

    @staticmethod
    def _extract_body(msg) -> str:
        body = ''
        if msg.is_multipart():
            for part in msg.walk():
                ct = part.get_content_type()
                if ct == 'text/plain':
                    try:
                        body = part.get_payload(decode=True).decode(
                            part.get_content_charset() or 'utf-8', errors='replace')
                        break
                    except Exception:
                        pass
        else:
            try:
                body = msg.get_payload(decode=True).decode(
                    msg.get_content_charset() or 'utf-8', errors='replace')
            except Exception:
                pass
        return body.strip()


class CalendarTool(BaseTool):
    """
    Calendar — работа с Google Calendar через Google Calendar API v3.

    Настройка (один раз):
        1. console.cloud.google.com → создай проект → включи «Google Calendar API»
        2. Credentials → OAuth 2.0 → Desktop App → скачай credentials.json
        3. Передай путь: CalendarTool(credentials_path='credentials.json')
        При первом запуске откроется браузер для авторизации → создастся token.json.

    Действия:
        list_events   — список событий (period: 'today'/'week'/'month' или start+end)
        create_event  — создать событие
        update_event  — изменить событие по event_id
        delete_event  — удалить событие по event_id
        find_event    — поиск по тексту в заголовке
        calendars     — список доступных календарей
    """

    _SCOPES = ['https://www.googleapis.com/auth/calendar']

    def __init__(self, credentials_path: str | None = None, token_path: str | None = None):
        super().__init__('calendar',
                         'Google Calendar: список, создание, изменение, удаление событий')
        self.credentials_path = credentials_path or 'credentials.json'
        self.token_path       = token_path or 'token.json'
        self._service         = None

    def run(self, action: str, **kwargs) -> dict:
        svc = self._get_service()
        if not svc:
            return {
                'error': (
                    'Google Calendar не настроен. '
                    'Скачай credentials.json из console.cloud.google.com '
                    'и положи рядом с агентом (или укажи путь через credentials_path)h).'
                ),
                'success': False,
            }

        try:
            if action == 'list_events':
                return self._list_events(svc, **kwargs)
            elif action == 'create_event':
                return self._create_event(svc, **kwargs)
            elif action == 'update_event':
                return self._update_event(svc, **kwargs)
            elif action == 'delete_event':
                return self._delete_event(svc, **kwargs)
            elif action == 'find_event':
                return self._find_event(svc, **kwargs)
            elif action == 'calendars':
                return self._list_calendars(svc)
            else:
                return {'error': f"Неизвестный action: {action}", 'success': False}
        except Exception as e:
            return {'error': str(e), 'success': False}

    # ── Google API ───────────────────────────────────────────────────────────

    def _get_service(self):
        if self._service:
            return self._service
        if not os.path.exists(self.credentials_path):
            return None
        try:
            from google.oauth2.credentials import Credentials
            from google_auth_oauthlib.flow import InstalledAppFlow
            from google.auth.transport.requests import Request
            from googleapiclient.discovery import build

            creds = None
            if os.path.exists(self.token_path):
                creds = Credentials.from_authorized_user_file(
                    self.token_path, self._SCOPES)

            if not creds or not creds.valid:
                if creds and creds.expired and creds.refresh_token:
                    creds.refresh(Request())
                else:
                    flow = InstalledAppFlow.from_client_secrets_file(
                        self.credentials_path, self._SCOPES)
                    creds = flow.run_local_server(port=0)
                with open(self.token_path, 'w', encoding='utf-8') as f:
                    f.write(creds.to_json())

            self._service = build('calendar', 'v3', credentials=creds)
            return self._service
        except ImportError:
            return None
        except Exception:
            return None

    def _list_events(self, svc, calendar_id: str = 'primary',
                     period: str = 'week', start: str | None = None,
                     end: str | None = None, limit: int = 20) -> dict:
        from datetime import datetime, timedelta, timezone

        now = datetime.now(timezone.utc)
        if start:
            t_min = start if start.endswith('Z') else start + 'Z'
        elif period == 'today':
            t_min = now.replace(hour=0, minute=0, second=0).isoformat()
        elif period == 'month':
            t_min = now.isoformat()
        else:   # week
            t_min = now.isoformat()

        if end:
            t_max = end if end.endswith('Z') else end + 'Z'
        elif period == 'today':
            t_max = now.replace(hour=23, minute=59, second=59).isoformat()
        elif period == 'month':
            t_max = (now + timedelta(days=30)).isoformat()
        else:
            t_max = (now + timedelta(days=7)).isoformat()

        result = svc.events().list(
            calendarId=calendar_id,
            timeMin=t_min, timeMax=t_max,
            maxResults=limit,
            singleEvents=True,
            orderBy='startTime',
        ).execute()

        events = []
        for e in result.get('items', []):
            start_val = e.get('start', {})
            events.append({
                'id':       e.get('id'),
                'summary':  e.get('summary', ''),
                'start':    start_val.get('dateTime', start_val.get('date', '')),
                'end':      e.get('end', {}).get('dateTime',
                             e.get('end', {}).get('date', '')),
                'location': e.get('location', ''),
                'description': e.get('description', '')[:200],
            })
        return {'events': events, 'count': len(events), 'success': True}

    def _create_event(self, svc, summary: str, start: str, end: str,
                      description: str = '', location: str = '',
                      calendar_id: str = 'primary',
                      attendees: list | None = None) -> dict:
        """
        start/end — ISO 8601: '2024-12-25T10:00:00+03:00'
                    или дата:  '2024-12-25' (весь день)
        """
        is_allday = 'T' not in start

        event_body: dict[str, Any] = {
            'summary':     summary,
            'description': description,
            'location':    location,
        }
        if is_allday:
            event_body['start'] = {'date': start}
            event_body['end']   = {'date': end}
        else:
            event_body['start'] = {'dateTime': start, 'timeZone': 'Europe/Moscow'}
            event_body['end']   = {'dateTime': end,   'timeZone': 'Europe/Moscow'}

        if attendees:
            event_body['attendees'] = [{'email': a} for a in attendees]

        created = svc.events().insert(
            calendarId=calendar_id, body=event_body).execute()
        return {
            'event_id': created.get('id'),
            'link':     created.get('htmlLink'),
            'summary':  created.get('summary'),
            'success':  True,
        }

    def _update_event(self, svc, event_id: str,
                      calendar_id: str = 'primary', **fields) -> dict:
        event = svc.events().get(
            calendarId=calendar_id, eventId=event_id).execute()
        for key, val in fields.items():
            event[key] = val
        updated = svc.events().update(
            calendarId=calendar_id, eventId=event_id, body=event).execute()
        return {'event_id': updated.get('id'), 'success': True}

    def _delete_event(self, svc, event_id: str,
                      calendar_id: str = 'primary') -> dict:
        svc.events().delete(
            calendarId=calendar_id, eventId=event_id).execute()
        return {'deleted': event_id, 'success': True}

    def _find_event(self, svc, query: str,
                    calendar_id: str = 'primary', limit: int = 10) -> dict:
        result = svc.events().list(
            calendarId=calendar_id,
            q=query,
            maxResults=limit,
            singleEvents=True,
            orderBy='startTime',
        ).execute()
        events = [
            {
                'id':      e.get('id'),
                'summary': e.get('summary', ''),
                'start':   e.get('start', {}).get('dateTime',
                            e.get('start', {}).get('date', '')),
            }
            for e in result.get('items', [])
        ]
        return {'events': events, 'count': len(events), 'success': True}

    def _list_calendars(self, svc) -> dict:
        result = svc.calendarList().list().execute()
        cals = [
            {'id': c.get('id'), 'summary': c.get('summary'),
             'primary': c.get('primary', False)}
            for c in result.get('items', [])
        ]
        return {'calendars': cals, 'count': len(cals), 'success': True}


class HuggingFaceTool(BaseTool):
    """
    HuggingFace Hub — поиск и анализ моделей, датасетов и Spaces.

    Использует HuggingFace Hub REST API v1 — только requests, без доп. зависимостей.
    Токен (HF_TOKEN) даёт доступ к приватным репозиториям и снимает rate-limit.

    Actions:
        search_models   — поиск моделей по задаче / запросу / языку
        model_info      — полная карточка модели: метрики, размер, лицензия
        search_datasets — поиск датасетов
        dataset_info    — информация о датасете
        search_spaces   — поиск Spaces (демо-приложений)
        trending        — трендовые модели / датасеты / Spaces
    """

    HF_API = 'https://huggingface.co/api'

    def __init__(self, token: str | None = None):
        super().__init__('huggingface',
                         'HuggingFace Hub: поиск моделей, датасетов и Spaces')
        self.token = token
        self._headers = {'Authorization': f'Bearer {token}'} if token else {}

    # ── Основной диспетчер ────────────────────────────────────────────────────

    def run(self, action: str, **kwargs) -> dict:
        """
        action: 'search_models' | 'model_info' | 'search_datasets' |
                'dataset_info'  | 'search_spaces' | 'trending'
        """
        try:
            import requests as _req
            self._req = _req  # pylint: disable=attribute-defined-outside-init
        except ImportError:
            return {'error': 'requests не установлен: pip install requests',
                    'success': False}

        _map = {
            'search_models':   self._search_models,
            'model_info':      self._model_info,
            'search_datasets': self._search_datasets,
            'dataset_info':    self._dataset_info,
            'search_spaces':   self._search_spaces,
            'trending':        self._trending,
        }
        fn = _map.get(action)
        if not fn:
            return {
                'error': f"Action '{action}' не поддерживается. "
                         f"Доступны: {list(_map)}",
                'success': False,
            }
        try:
            return fn(**kwargs)
        except Exception as e:
            return {'error': str(e), 'success': False}

    # ── Модели ────────────────────────────────────────────────────────────────

    def _search_models(self, query: str = '', task: str | None = None,
                       language: str | None = None, library: str | None = None,
                       limit: int = 10) -> dict:
        """
        Поиск моделей на HuggingFace Hub.

        Args:
            query    — ключевые слова (название, автор)
            task     — задача: text-classification | translation |
                       text-generation | image-classification | ...
            language — язык: ru | en | zh | ...
            library  — transformers | diffusers | timm | sentence-transformers | ...
            limit    — максимум результатов (1-100)
        """
        params = {'limit': min(limit, 100), 'sort': 'downloads', 'direction': -1,
                  'full': 'True'}
        if query:
            params['search'] = query
        if task:
            params['pipeline_tag'] = task
        if language:
            params['language'] = language
        if library:
            params['library'] = library

        raw = self._get(f'{self.HF_API}/models', params)
        if not isinstance(raw, list):
            return {'error': str(raw), 'success': False}

        models = []
        for m in raw:
            models.append({
                'id':            m.get('id', ''),
                'task':          m.get('pipeline_tag', ''),
                'downloads':     m.get('downloads', 0),
                'likes':         m.get('likes', 0),
                'last_modified': m.get('lastModified', '')[:10],
                'library':       m.get('library_name', ''),
                'tags':          m.get('tags', [])[:6],
                'private':       m.get('private', False),
                'gated':         m.get('gated', False),
            })
        return {'models': models, 'count': len(models),
                'query': query, 'task': task, 'success': True}

    def _model_info(self, model_id: str) -> dict:
        """
        Полная карточка модели: описание, метрики, файлы, размер, лицензия.

        Args:
            model_id — 'bert-base-uncased' или 'mistralai/Mistral-7B-v0.1'
        """
        data = self._get(f'{self.HF_API}/models/{model_id}')
        if isinstance(data, dict) and 'error' in data and 'id' not in data:
            return {'error': data.get('error', 'not found'), 'success': False}

        card_data: dict = data.get('cardData') or {}  # type: ignore[assignment]

        # Метрики eval
        metrics = []
        for er in card_data.get('eval_results', []):
            m_list = er.get('metrics', [{}])
            metrics.append({
                'task':    er.get('task', {}).get('name', ''),
                'dataset': er.get('dataset', {}).get('name', ''),
                'metric':  m_list[0].get('name', '') if m_list else '',
                'value':   m_list[0].get('value', '') if m_list else '',
            })

        # Файлы и размер
        siblings: list = data.get('siblings', [])  # type: ignore[assignment]
        files     = [s['rfilename'] for s in siblings if 'rfilename' in s]
        total_bytes = sum(s.get('size', 0) for s in siblings if s.get('size'))
        size_gb   = round(total_bytes / 1e9, 2) if total_bytes else None

        return {
            'id':            data.get('id', model_id),
            'task':          data.get('pipeline_tag', ''),
            'author':        data.get('author', ''),
            'downloads':     data.get('downloads', 0),
            'likes':         data.get('likes', 0),
            'license':       card_data.get('license', 'unknown'),
            'language':      card_data.get('language', []),
            'library':       data.get('library_name', ''),
            'tags':          data.get('tags', []),
            'metrics':       metrics,
            'files':         files[:25],
            'size_gb':       size_gb,
            'private':       data.get('private', False),
            'gated':         data.get('gated', False),
            'success': True,
        }

    # ── Датасеты ──────────────────────────────────────────────────────────────

    def _search_datasets(self, query: str = '', task: str | None = None,
                         language: str | None = None, limit: int = 10) -> dict:
        """
        Поиск датасетов.

        Args:
            query    — ключевые слова
            task     — категория задачи
            language — язык датасета
            limit    — максимум результатов
        """
        params = {'limit': min(limit, 100), 'sort': 'downloads', 'direction': -1}
        if query:
            params['search'] = query
        if task:
            params['task_categories'] = task
        if language:
            params['language'] = language

        raw = self._get(f'{self.HF_API}/datasets', params)
        if not isinstance(raw, list):
            return {'error': str(raw), 'success': False}

        datasets = []
        for d in raw:
            datasets.append({
                'id':            d.get('id', ''),
                'downloads':     d.get('downloads', 0),
                'likes':         d.get('likes', 0),
                'last_modified': d.get('lastModified', '')[:10],
                'tags':          d.get('tags', [])[:6],
            })
        return {'datasets': datasets, 'count': len(datasets),
                'query': query, 'success': True}

    def _dataset_info(self, dataset_id: str) -> dict:
        """Полная информация о датасете."""
        data = self._get(f'{self.HF_API}/datasets/{dataset_id}')
        if isinstance(data, dict) and 'error' in data and 'id' not in data:
            return {'error': data.get('error', 'not found'), 'success': False}

        card_data: dict = data.get('cardData') or {}  # type: ignore[assignment]
        return {
            'id':              data.get('id', dataset_id),
            'downloads':       data.get('downloads', 0),
            'likes':           data.get('likes', 0),
            'license':         card_data.get('license', 'unknown'),
            'language':        card_data.get('language', []),
            'task_categories': card_data.get('task_categories', []),
            'tags':            data.get('tags', []),
            'private':         data.get('private', False),
            'success': True,
        }

    # ── Spaces ────────────────────────────────────────────────────────────────

    def _search_spaces(self, query: str = '', limit: int = 10) -> dict:
        """Поиск Spaces (демо-приложений на HuggingFace)."""
        params = {'limit': min(limit, 100), 'sort': 'likes', 'direction': -1}
        if query:
            params['search'] = query

        raw = self._get(f'{self.HF_API}/spaces', params)
        if not isinstance(raw, list):
            return {'error': str(raw), 'success': False}

        spaces = []
        for s in raw:
            spaces.append({
                'id':    s.get('id', ''),
                'likes': s.get('likes', 0),
                'sdk':   s.get('sdk', ''),
                'tags':  s.get('tags', [])[:5],
            })
        return {'spaces': spaces, 'count': len(spaces),
                'query': query, 'success': True}

    # ── Тренды ────────────────────────────────────────────────────────────────

    def _trending(self, category: str = 'models', limit: int = 10) -> dict:
        """
        Трендовые позиции.

        Args:
            category — 'models' | 'datasets' | 'spaces'
            limit    — количество результатов
        """
        if category not in ('models', 'datasets', 'spaces'):
            return {'error': "category должен быть models/datasets/spaces",
                    'success': False}
        params = {'limit': min(limit, 50), 'sort': 'trending', 'direction': -1}
        raw = self._get(f'{self.HF_API}/{category}', params)
        if not isinstance(raw, list):
            return {'error': str(raw), 'success': False}

        items = []
        for item in raw:
            items.append({
                'id':        item.get('id', ''),
                'downloads': item.get('downloads', 0),
                'likes':     item.get('likes', 0),
                'task':      item.get('pipeline_tag', ''),
                'tags':      item.get('tags', [])[:3],
            })
        return {'items': items, 'count': len(items),
                'category': category, 'success': True}

    # ── HTTP helper ───────────────────────────────────────────────────────────

    def _get(self, url: str, params: dict | None = None) -> Any:
        try:
            r = self._req.get(url, headers=self._headers,
                              params=params, timeout=15)
            r.raise_for_status()
            return r.json()
        except Exception as e:
            return {'error': str(e)}


# ── Группа 6: Архивы, SSH, слежение за файлами ───────────────────────────────

class ArchiveTool(BaseTool):
    """
    Archive — работа с ZIP, TAR и RAR архивами.
    Действия: create, extract, list, add_file
    Форматы: zip, tar, tar.gz, tar.bz2, rar (требует WinRAR или unrar)
    """

    # Стандартные пути к WinRAR на Windows
    _WINRAR_PATHS = [
        r'C:\Program Files\WinRAR\Rar.exe',
        r'C:\Program Files (x86)\WinRAR\Rar.exe',
    ]
    _UNRAR_PATHS = [
        r'C:\Program Files\WinRAR\UnRAR.exe',
        r'C:\Program Files (x86)\WinRAR\UnRAR.exe',
    ]

    def __init__(self):
        super().__init__('archive', 'Работа с архивами: создание, распаковка, список файлов (zip/tar/rar)')

    def _find_exe(self, paths: list) -> str | None:
        import shutil
        for p in paths:
            if os.path.exists(p):
                return p
        # Попробуем через PATH
        name = os.path.basename(paths[0]).lower()
        found = shutil.which(name)
        return found

    def run(self, action: str, archive_path: str, **kwargs) -> dict:
        """
        action: 'create' | 'extract' | 'list' | 'add_file'
        create:   source_paths (list[str]), format='zip'|'tar'|'tar.gz'|'tar.bz2'|'rar'
                  password (str, optional) — для rar/zip
        extract:  dest_dir (str), password (str, optional)
        list:     —
        add_file: file_path (str), arcname (str, optional)
        """
        try:
            import zipfile
            import tarfile as _tarfile

            fmt = kwargs.get('format', 'zip')
            password = kwargs.get('password')

            # ── RAR: используем WinRAR CLI ──────────────────────────────────────
            if fmt == 'rar' or (action in ('extract', 'list') and
                                archive_path.lower().endswith('.rar')):
                return self._rar_action(action, archive_path, kwargs)

            if action == 'create':
                sources = kwargs.get('source_paths', [])
                if fmt == 'zip':
                    with zipfile.ZipFile(archive_path, 'w',
                                         compression=zipfile.ZIP_DEFLATED) as zf:
                        if password:
                            zf.setpassword(password.encode())
                        for src in sources:
                            if os.path.isdir(src):
                                for root, _, files in os.walk(src):
                                    for fname in files:
                                        fp = os.path.join(root, fname)
                                        zf.write(fp, os.path.relpath(fp, src))
                            elif os.path.exists(src):
                                zf.write(src, os.path.basename(src))
                    return {'path': archive_path, 'format': 'zip', 'success': True}
                else:
                    mode = 'w:gz' if 'gz' in fmt else ('w:bz2' if 'bz2' in fmt else 'w')
                    with _tarfile.open(archive_path, mode) as tf:
                        for src in sources:
                            tf.add(src, arcname=os.path.basename(src))
                    return {'path': archive_path, 'format': fmt, 'success': True}

            elif action == 'extract':
                dest = kwargs.get('dest_dir', '.')
                os.makedirs(dest, exist_ok=True)
                if zipfile.is_zipfile(archive_path):
                    with zipfile.ZipFile(archive_path, 'r') as zf:
                        if password:
                            zf.setpassword(password.encode())
                        zf.extractall(dest)
                    return {'dest': dest, 'success': True}
                elif _tarfile.is_tarfile(archive_path):
                    with _tarfile.open(archive_path, 'r:*') as tf:
                        tf.extractall(dest)
                    return {'dest': dest, 'success': True}
                return {'error': 'Неизвестный формат архива', 'success': False}

            elif action == 'list':
                if zipfile.is_zipfile(archive_path):
                    with zipfile.ZipFile(archive_path, 'r') as zf:
                        names = zf.namelist()
                    return {'files': names, 'count': len(names), 'success': True}
                elif _tarfile.is_tarfile(archive_path):
                    with _tarfile.open(archive_path, 'r:*') as tf:
                        names = tf.getnames()
                    return {'files': names, 'count': len(names), 'success': True}
                return {'error': 'Неизвестный формат архива', 'success': False}

            elif action == 'add_file':
                file_path = kwargs.get('file_path', '')
                arcname = kwargs.get('arcname', os.path.basename(file_path))
                if not os.path.exists(file_path):
                    return {'error': f'Файл не найден: {file_path}', 'success': False}
                with zipfile.ZipFile(archive_path, 'a',
                                     compression=zipfile.ZIP_DEFLATED) as zf:
                    zf.write(file_path, arcname)
                return {'added': arcname, 'success': True}

            return {'error': f"Неизвестный action: {action}", 'success': False}
        except Exception as e:
            return {'error': str(e), 'success': False}

    def _rar_action(self, action: str, archive_path: str, kwargs: dict) -> dict:
        """RAR операции через WinRAR CLI (Rar.exe / UnRAR.exe)."""
        # SECURITY: validate all paths through PathValidator
        _pv = PathValidator()
        _ok, _reason = _pv.validate(archive_path)
        if not _ok:
            return {'error': f'Path blocked: {_reason}', 'success': False}

        password = kwargs.get('password')
        # SECURITY: sanitize password — no shell metacharacters
        if password and not re.match(r'^[\w\-@#$%^&*()+=!.]+$', str(password)):
            return {'error': 'Password contains forbidden characters', 'success': False}
        pwd_flag = [f'-p{password}'] if password else []

        if action == 'create':
            rar_exe = self._find_exe(self._WINRAR_PATHS)
            if not rar_exe:
                return {'error': 'WinRAR не найден. Установи WinRAR или добавь Rar.exe в PATH.', 'success': False}
            sources = kwargs.get('source_paths', [])
            if not sources:
                return {'error': 'source_paths обязателен для create', 'success': False}
            # SECURITY: validate source paths
            for sp in sources:
                _s_ok, _s_reason = _pv.validate(sp)
                if not _s_ok:
                    return {'error': f'Source path blocked: {_s_reason}', 'success': False}
            cmd = [rar_exe, 'a', '-ep1'] + pwd_flag + [archive_path] + sources
            from execution.command_gateway import CommandGateway
            gw = CommandGateway.get_instance()
            r = gw.execute(cmd, timeout=120, caller='ArchiveTool.rar_create')
            if not r.allowed:
                return {'error': r.reject_reason, 'success': False}
            if r.returncode not in (0, 1):  # WinRAR: 0=ok, 1=warning
                return {'error': r.stderr or r.stdout, 'success': False}
            return {'path': archive_path, 'format': 'rar', 'success': True}

        elif action == 'extract':
            unrar_exe = self._find_exe(self._UNRAR_PATHS) or self._find_exe(self._WINRAR_PATHS)
            if not unrar_exe:
                # Попробуем через rarfile (требует unrar в PATH)
                try:
                    import rarfile
                    dest = kwargs.get('dest_dir', '.')
                    os.makedirs(dest, exist_ok=True)
                    with rarfile.RarFile(archive_path) as rf:
                        rf.extractall(dest, pwd=password)
                    return {'dest': dest, 'success': True}
                except Exception as e:
                    return {'error': f'Нет UnRAR.exe и rarfile не смог: {e}', 'success': False}
            dest = kwargs.get('dest_dir', '.')
            os.makedirs(dest, exist_ok=True)
            exe_name = os.path.basename(unrar_exe).lower()
            if 'unrar' in exe_name:
                cmd = ['UnRAR.exe' if 'unrar' in exe_name else unrar_exe,
                       'x', '-y'] + pwd_flag + [archive_path, dest + '\\']
                cmd[0] = unrar_exe
            else:
                cmd = [unrar_exe, 'x', '-y'] + pwd_flag + [archive_path, dest + '\\']
            from execution.command_gateway import CommandGateway
            gw = CommandGateway.get_instance()
            r = gw.execute(cmd, timeout=120, caller='ArchiveTool.rar_extract')
            if not r.allowed:
                return {'error': r.reject_reason, 'success': False}
            if r.returncode not in (0, 1):
                return {'error': r.stderr or r.stdout, 'success': False}
            return {'dest': dest, 'success': True}

        elif action == 'list':
            unrar_exe = self._find_exe(self._UNRAR_PATHS) or self._find_exe(self._WINRAR_PATHS)
            if not unrar_exe:
                try:
                    import rarfile
                    with rarfile.RarFile(archive_path) as rf:
                        names = rf.namelist()
                    return {'files': names, 'count': len(names), 'success': True}
                except Exception as e:
                    return {'error': str(e), 'success': False}
            cmd = [unrar_exe, 'l'] + pwd_flag + [archive_path]
            from execution.command_gateway import CommandGateway
            gw = CommandGateway.get_instance()
            r = gw.execute(cmd, timeout=60, caller='ArchiveTool.rar_list')
            if not r.allowed:
                return {'error': r.reject_reason, 'success': False}
            # Парсим вывод UnRAR
            lines = r.stdout.splitlines()
            files = [ln.split()[-1] for ln in lines if ln.strip() and
                     not ln.startswith('-') and not ln.startswith('RAR') and
                     len(ln.split()) >= 5]
            return {'files': files, 'count': len(files), 'raw': r.stdout, 'success': True}

        return {'error': f"Неизвестный action для RAR: {action}", 'success': False}


class SSHTool(BaseTool):
    """
    SSH — подключение к удалённым серверам через paramiko.
    Действия: connect, exec, upload, download, disconnect
    """

    def __init__(self):
        super().__init__('ssh', 'SSH: подключение к серверу, выполнение команд, передача файлов')
        self._client = None

    def run(self, action: str, **kwargs) -> dict:
        """
        connect:    host, port=22, username, password=None, key_path=None
        exec:       command (str)
        upload:     local_path, remote_path
        download:   remote_path, local_path
        disconnect: —
        """
        try:
            import paramiko  # type: ignore[import]

            if action == 'connect':
                host     = kwargs['host']
                port     = int(kwargs.get('port', 22))
                username = kwargs['username']
                password = kwargs.get('password')
                key_path = kwargs.get('key_path')

                client = paramiko.SSHClient()
                client.set_missing_host_key_policy(paramiko.AutoAddPolicy())
                connect_kwargs: dict = {'hostname': host, 'port': port, 'username': username}
                if key_path:
                    try:
                        connect_kwargs['pkey'] = paramiko.Ed25519Key.from_private_key_file(key_path)
                    except Exception:
                        connect_kwargs['pkey'] = paramiko.RSAKey.from_private_key_file(key_path)
                elif password:
                    connect_kwargs['password'] = password
                client.connect(**connect_kwargs)
                self._client = client
                return {'connected': host, 'port': port, 'success': True}

            elif action == 'exec':
                if not self._client:
                    return {'error': 'Нет активного SSH-соединения. Сначала вызови connect.',
                            'success': False}
                command = kwargs.get('command', '')
                _, stdout, stderr = self._client.exec_command(
                    command, timeout=kwargs.get('timeout', 30))
                out = stdout.read().decode('utf-8', errors='replace')
                err = stderr.read().decode('utf-8', errors='replace')
                rc  = stdout.channel.recv_exit_status()
                return {'stdout': out, 'stderr': err, 'returncode': rc,
                        'success': rc == 0}

            elif action == 'upload':
                if not self._client:
                    return {'error': 'Нет активного SSH-соединения', 'success': False}
                sftp = self._client.open_sftp()
                sftp.put(kwargs['local_path'], kwargs['remote_path'])
                sftp.close()
                return {'uploaded': kwargs['remote_path'], 'success': True}

            elif action == 'download':
                if not self._client:
                    return {'error': 'Нет активного SSH-соединения', 'success': False}
                sftp = self._client.open_sftp()
                sftp.get(kwargs['remote_path'], kwargs['local_path'])
                sftp.close()
                return {'downloaded': kwargs['local_path'], 'success': True}

            elif action == 'disconnect':
                if self._client:
                    self._client.close()
                    self._client = None
                return {'disconnected': True, 'success': True}

            return {'error': f"Неизвестный action: {action}", 'success': False}

        except ImportError:
            return {'error': 'paramiko не установлен: pip install paramiko', 'success': False}
        except Exception as e:
            return {'error': str(e), 'success': False}


class FileWatcherTool(BaseTool):
    """
    File Watcher — мониторинг изменений файлов и директорий через watchdog.
    Действия: watch, stop, list_events
    """

    def __init__(self):
        super().__init__('file_watcher',
                         'Слежение за изменениями файлов и папок (watchdog)')
        self._observers: dict = {}
        self._events: list = []

    def run(self, action: str, **kwargs) -> dict:
        """
        watch:        path (str), recursive=True, watch_id='default'
        stop:         watch_id='default'
        list_events:  limit=50 — последние события
        clear_events: —
        """
        try:
            import watchdog.observers as _obs  # pyright: ignore[reportMissingImports]
            import watchdog.events   as _evt   # pyright: ignore[reportMissingImports]

            if action == 'watch':
                path       = kwargs.get('path', '.')
                recursive  = kwargs.get('recursive', True)
                watch_id   = kwargs.get('watch_id', 'default')
                events_ref = self._events

                class _Handler(_evt.FileSystemEventHandler):
                    def on_any_event(self, event):
                        events_ref.append({
                            'type':    event.event_type,
                            'path':    event.src_path,
                            'is_dir':  event.is_directory,
                            'time':    time.time(),
                        })
                        if len(events_ref) > 1000:
                            events_ref.pop(0)

                observer = _obs.Observer()
                observer.schedule(_Handler(), path, recursive=recursive)
                observer.start()
                self._observers[watch_id] = observer
                return {'watching': path, 'watch_id': watch_id, 'success': True}

            elif action == 'stop':
                watch_id = kwargs.get('watch_id', 'default')
                obs = self._observers.pop(watch_id, None)
                if obs:
                    obs.stop()
                    obs.join(timeout=3)
                return {'stopped': watch_id, 'success': True}

            elif action == 'list_events':
                limit  = kwargs.get('limit', 50)
                events = self._events[-limit:]
                return {'events': events, 'total': len(self._events), 'success': True}

            elif action == 'clear_events':
                self._events.clear()
                return {'cleared': True, 'success': True}

            return {'error': f"Неизвестный action: {action}", 'success': False}

        except ImportError:
            return {'error': 'watchdog не установлен: pip install watchdog', 'success': False}
        except Exception as e:
            return {'error': str(e), 'success': False}


# ── Группа 7: Безопасность и уведомления ─────────────────────────────────────

class NotificationTool(BaseTool):
    """
    Notification — системные уведомления через plyer.
    Fallback на Windows: PowerShell BurntToast (если plyer недоступен).
    """

    def __init__(self):
        super().__init__('notification',
                         'Системные уведомления на рабочем столе (Windows/Linux/macOS)')

    def run(self, title: str, message: str, timeout: int = 10,
            app_name: str = 'Agent', **kwargs) -> dict:
        try:
            from plyer import notification as _notif  # pyright: ignore[reportMissingImports]
            _notif.notify(  # type: ignore[operator]
                title=title,
                message=message,
                app_name=app_name,
                timeout=timeout,
            )
            return {'sent': True, 'title': title, 'success': True}
        except ImportError:
            pass

        # Fallback: Windows PowerShell toast
        if os.name == 'nt':
            try:
                # SECURITY: sanitize title/message to prevent PS injection
                def _ps_escape(s: str) -> str:
                    return str(s).replace('"', '`"').replace("'", "`'").replace('`', '``')[:200]
                safe_title = _ps_escape(title)
                safe_msg = _ps_escape(message)
                safe_timeout = int(max(1, min(timeout, 30))) * 1000
                ps_cmd = (
                    f'Add-Type -AssemblyName System.Windows.Forms; '
                    f'$n = New-Object System.Windows.Forms.NotifyIcon; '
                    f'$n.Icon = [System.Drawing.SystemIcons]::Information; '
                    f'$n.Visible = $true; '
                    f'$n.ShowBalloonTip({safe_timeout}, '
                    f'"{safe_title}", "{safe_msg}", '
                    f'[System.Windows.Forms.ToolTipIcon]::None); '
                    f'Start-Sleep -Milliseconds {safe_timeout + 500}; '
                    f'$n.Dispose()'
                )
                subprocess.Popen(
                    ['powershell', '-NoProfile', '-NonInteractive',
                     '-Command', ps_cmd],
                    creationflags=subprocess.CREATE_NO_WINDOW
                )
                return {'sent': True, 'title': title,
                        'method': 'powershell', 'success': True}
            except Exception as e:
                return {'error': str(e), 'success': False}

        return {'error': 'plyer не установлен: pip install plyer', 'success': False}


class EncryptionTool(BaseTool):
    """
    Encryption — шифрование и дешифровка строк и файлов через cryptography (Fernet).
    """

    def __init__(self, key: str | None = None):
        """
        key — base64-urlsafe ключ Fernet (44 символа).
              Если None — генерируется новый при каждой инициализации.
              Сохрани ключ! Без него дешифровка невозможна.
        """
        super().__init__('encryption', 'Шифрование / дешифрование строк и файлов (Fernet/AES)')
        self._key_str = key

    def _get_fernet(self):
        from cryptography.fernet import Fernet  # pyright: ignore[reportMissingImports]
        if not self._key_str:
            self._key_str = Fernet.generate_key().decode()
        return Fernet(self._key_str.encode() if isinstance(self._key_str, str)
                      else self._key_str)

    def run(self, action: str, **kwargs) -> dict:
        """
        generate_key: — → {'key': str}
        encrypt_text: text (str) → {'encrypted': str}
        decrypt_text: encrypted (str) → {'text': str}
        encrypt_file: path (str), output (str, optional)
        decrypt_file: path (str), output (str, optional)
        """
        try:
            if action == 'generate_key':
                from cryptography.fernet import Fernet
                key = Fernet.generate_key().decode()
                self._key_str = key
                return {'key': key, 'success': True}

            elif action == 'encrypt_text':
                text = kwargs.get('text', '')
                f = self._get_fernet()
                encrypted = f.encrypt(text.encode('utf-8')).decode()
                return {'encrypted': encrypted, 'key': self._key_str, 'success': True}

            elif action == 'decrypt_text':
                encrypted = kwargs.get('encrypted', '')
                f = self._get_fernet()
                text = f.decrypt(encrypted.encode()).decode('utf-8')
                return {'text': text, 'success': True}

            elif action == 'encrypt_file':
                path = kwargs['path']
                output = kwargs.get('output', path + '.enc')
                f = self._get_fernet()
                with open(path, 'rb') as fp:
                    data = fp.read()
                with open(output, 'wb') as fp:
                    fp.write(f.encrypt(data))
                return {'output': output, 'key': self._key_str, 'success': True}

            elif action == 'decrypt_file':
                path = kwargs['path']
                output = kwargs.get('output', path.removesuffix('.enc'))
                f = self._get_fernet()
                with open(path, 'rb') as fp:
                    data = fp.read()
                with open(output, 'wb') as fp:
                    fp.write(f.decrypt(data))
                return {'output': output, 'success': True}

            return {'error': f"Неизвестный action: {action}", 'success': False}

        except ImportError:
            return {'error': 'cryptography не установлен: pip install cryptography',
                    'success': False}
        except Exception as e:
            return {'error': str(e), 'success': False}


# ── Группа 8: NLP и ML ────────────────────────────────────────────────────────

class NLPTool(BaseTool):
    """
    NLP — обработка естественного языка: токенизация, ключевые слова, сентимент.
    Использует NLTK (локально). OpenAI fallback для сентимента.
    """

    def __init__(self, openai_client=None):
        super().__init__('nlp',
                         'NLP: ключевые слова, сентимент, токенизация, именованные сущности (NLTK)')
        self.client = openai_client

    def run(self, action: str, text: str, **kwargs) -> dict:
        """
        action: 'keywords' | 'sentiment' | 'tokenize' | 'sentences' | 'summary_stats'
        """
        if action == 'sentiment' and self.client:
            return self._sentiment_llm(text)

        try:
            import nltk  # pyright: ignore[reportMissingImports]

            # Тихая загрузка нужных корпусов
            for corpus in ('punkt', 'stopwords', 'averaged_perceptron_tagger',
                           'maxent_ne_chunker', 'words', 'punkt_tab'):
                try:
                    nltk.data.find(f'tokenizers/{corpus}')
                except LookupError:
                    try:
                        nltk.download(corpus, quiet=True)
                    except Exception:
                        pass

            if action == 'tokenize':
                tokens = nltk.word_tokenize(text)
                return {'tokens': tokens, 'count': len(tokens), 'success': True}

            elif action == 'sentences':
                sents = nltk.sent_tokenize(text)
                return {'sentences': sents, 'count': len(sents), 'success': True}

            elif action == 'keywords':
                from nltk.corpus   import stopwords as _sw
                from nltk.tokenize import word_tokenize
                import string
                lang = kwargs.get('language', 'english')
                stop_words = set(_sw.words(lang)) if lang in _sw.fileids() else set()
                tokens = word_tokenize(text.lower())
                words  = [t for t in tokens
                          if t not in stop_words and t not in string.punctuation
                          and len(t) > 2]
                from collections import Counter
                freq = Counter(words)
                top  = freq.most_common(kwargs.get('top_n', 10))
                return {'keywords': [{'word': w, 'freq': f} for w, f in top],
                        'success': True}

            elif action == 'sentiment':
                # NLTK VADER (если доступен)
                try:
                    from nltk.sentiment import SentimentIntensityAnalyzer
                    nltk.download('vader_lexicon', quiet=True)
                    sia = SentimentIntensityAnalyzer()
                    scores = sia.polarity_scores(text)
                    label = ('positive' if scores['compound'] >= 0.05
                             else 'negative' if scores['compound'] <= -0.05
                             else 'neutral')
                    return {'label': label, 'scores': scores, 'success': True}
                except Exception:
                    pass
                if self.client:
                    return self._sentiment_llm(text)
                return {'error': 'VADER недоступен и нет LLM клиента', 'success': False}

            elif action == 'summary_stats':
                from nltk.tokenize import word_tokenize
                tokens = word_tokenize(text)
                sents  = nltk.sent_tokenize(text)
                from collections import Counter
                freq = Counter(t.lower() for t in tokens if t.isalpha())
                return {
                    'char_count':  len(text),
                    'word_count':  len(tokens),
                    'sent_count':  len(sents),
                    'unique_words': len(freq),
                    'top5':        freq.most_common(5),
                    'success': True,
                }

            return {'error': f"Неизвестный action: {action}", 'success': False}

        except ImportError:
            return {'error': 'nltk не установлен: pip install nltk', 'success': False}
        except Exception as e:
            return {'error': str(e), 'success': False}

    def _sentiment_llm(self, text: str) -> dict:
        if not self.client:
            return {'error': 'OpenAI client не подключён', 'success': False}
        try:
            result = self.client.infer(
                f'Определи тональность текста: positive, negative или neutral.\n'
                f'Ответь одним словом.\n\nТекст: {text[:500]}',
                system='Определяй тональность. Ответ только одним словом: positive/negative/neutral.',
            ).strip().lower()
            return {'label': result, 'method': 'llm', 'success': True}
        except Exception as e:
            return {'error': str(e), 'success': False}


class EmbeddingTool(BaseTool):
    """
    Embedding — векторизация текста через sentence-transformers (локально).
    Fallback: OpenAI text-embedding-3-small.
    """

    def __init__(self, model_name: str = 'paraphrase-multilingual-MiniLM-L12-v2',
                 openai_client=None):
        super().__init__('embedding',
                         'Векторизация текста: sentence-transformers (локально) или OpenAI')
        self.model_name  = model_name
        self.openai_client = openai_client
        self._model = None

    def _load_model(self):
        if self._model is None:
            from sentence_transformers import SentenceTransformer  # pyright: ignore[reportMissingImports]
            self._model = SentenceTransformer(self.model_name)
        return self._model

    def run(self, action: str, **kwargs) -> dict:
        """
        embed:    text (str) → {'vector': list[float]}
        embed_many: texts (list[str]) → {'vectors': list[list[float]]}
        similarity: text_a (str), text_b (str) → {'score': float}
        search:   query (str), texts (list[str]), top_k=5
        """
        if action == 'embed':
            return self._embed_single(kwargs.get('text', ''))

        elif action == 'embed_many':
            texts = kwargs.get('texts', [])
            vecs  = []
            for t in texts:
                r = self._embed_single(t)
                if not r.get('success'):
                    return r
                vecs.append(r['vector'])
            return {'vectors': vecs, 'count': len(vecs), 'success': True}

        elif action == 'similarity':
            import math
            r_a = self._embed_single(kwargs.get('text_a', ''))
            r_b = self._embed_single(kwargs.get('text_b', ''))
            if not r_a.get('success') or not r_b.get('success'):
                return {'error': 'Не удалось получить эмбеддинги', 'success': False}
            a, b = r_a['vector'], r_b['vector']
            dot  = sum(x * y for x, y in zip(a, b))
            norm = (math.sqrt(sum(x**2 for x in a)) *
                    math.sqrt(sum(x**2 for x in b)))
            score = dot / norm if norm else 0.0
            return {'score': round(score, 4), 'success': True}

        elif action == 'search':
            import math
            query = kwargs.get('query', '')
            texts = kwargs.get('texts', [])
            top_k = kwargs.get('top_k', 5)
            r_q = self._embed_single(query)
            if not r_q.get('success'):
                return r_q
            q_vec = r_q['vector']
            results = []
            for i, t in enumerate(texts):
                r = self._embed_single(t)
                if not r.get('success'):
                    continue
                v    = r['vector']
                dot  = sum(x * y for x, y in zip(q_vec, v))
                norm = (math.sqrt(sum(x**2 for x in q_vec)) *
                        math.sqrt(sum(x**2 for x in v)))
                score = dot / norm if norm else 0.0
                results.append({'index': i, 'text': t[:100], 'score': round(score, 4)})
            results.sort(key=lambda r: r['score'], reverse=True)
            return {'results': results[:top_k], 'success': True}

        return {'error': f"Неизвестный action: {action}", 'success': False}

    def _embed_single(self, text: str) -> dict:
        # Попытка локальной модели
        try:
            model = self._load_model()
            vec = model.encode(text).tolist()
            return {'vector': vec, 'dim': len(vec), 'success': True}
        except ImportError:
            pass

        # Fallback: OpenAI
        if self.openai_client:
            try:
                oa = getattr(self.openai_client, '_client', None)
                if oa:
                    r = oa.embeddings.create(
                        model='text-embedding-3-small', input=[text])
                    vec = r.data[0].embedding
                    return {'vector': vec, 'dim': len(vec),
                            'method': 'openai', 'success': True}
            except Exception as e:
                return {'error': str(e), 'success': False}

        return {'error': ('sentence-transformers не установлен: '
                          'pip install sentence-transformers'), 'success': False}


# ── Группа 9: Разработка и инфраструктура ────────────────────────────────────

class CodeAnalyzerTool(BaseTool):
    """
    Code Analyzer — статический анализ Python-кода: AST, сложность, pylint, black.
    """

    def __init__(self):
        super().__init__('code_analyzer',
                         'Статический анализ Python: AST, сложность, pylint, форматирование')

    def run(self, action: str, **kwargs) -> dict:
        """
        analyze:   code (str) | path (str) — AST-анализ: функции, классы, импорты
        complexity: code (str) | path (str) — цикломатическая сложность (ast)
        lint:      path (str) — запустить pylint (если установлен)
        format:    code (str) — форматировать через black (возвращает текст)
        check_format: code (str) — проверить форматирование без изменений
        """
        try:
            import ast

            if action in ('analyze', 'complexity'):
                code = kwargs.get('code')
                path = kwargs.get('path')
                if not code and path:
                    with open(path, 'r', encoding='utf-8') as f:
                        code = f.read()
                if not code:
                    return {'error': 'Нужен code или path', 'success': False}

                tree = ast.parse(code)

                if action == 'analyze':
                    functions = [n.name for n in ast.walk(tree)
                                 if isinstance(n, ast.FunctionDef)]
                    classes   = [n.name for n in ast.walk(tree)
                                 if isinstance(n, ast.ClassDef)]
                    imports   = []
                    for n in ast.walk(tree):
                        if isinstance(n, ast.Import):
                            imports += [a.name for a in n.names]
                        elif isinstance(n, ast.ImportFrom):
                            imports.append(n.module or '')
                    return {
                        'functions':   functions,
                        'classes':     classes,
                        'imports':     list(set(imports)),
                        'lines':       code.count('\n') + 1,
                        'success': True,
                    }

                elif action == 'complexity':
                    # Простой подсчёт ветвлений (if/for/while/with/try/except)
                    branch_nodes = (ast.If, ast.For, ast.While, ast.With,
                                    ast.Try, ast.ExceptHandler, ast.BoolOp)
                    branches = sum(1 for n in ast.walk(tree)
                                   if isinstance(n, branch_nodes))
                    functions = {n.name: sum(1 for c in ast.walk(n)
                                             if isinstance(c, branch_nodes))
                                 for n in ast.walk(tree)
                                 if isinstance(n, ast.FunctionDef)}
                    return {
                        'total_branches':    branches,
                        'function_complexity': functions,
                        'success': True,
                    }

            elif action == 'lint':
                path = kwargs.get('path', '')
                # SECURITY: validate path through PathValidator
                if path:
                    _pv = PathValidator()
                    _pv_ok, _pv_reason = _pv.validate(path)
                    if not _pv_ok:
                        return {'error': f'Path blocked: {_pv_reason}', 'success': False}
                try:
                    from execution.command_gateway import CommandGateway
                    gw = CommandGateway.get_instance()
                    r = gw.execute(
                        ['python', '-m', 'pylint', '--output-format=text', path],
                        timeout=60, caller='CodeAnalysisTool.lint',
                    )
                    if not r.allowed:
                        return {'error': r.reject_reason, 'success': False}
                    return {
                        'output':     r.stdout[:3000],
                        'returncode': r.returncode,
                        'success':    True,
                    }
                except FileNotFoundError:
                    return {'error': 'pylint не установлен: pip install pylint',
                            'success': False}

            elif action in ('format', 'check_format'):
                code = kwargs.get('code', '')
                try:
                    import black  # pyright: ignore[reportMissingImports]
                    mode = black.Mode()
                    if action == 'check_format':
                        try:
                            black.format_str(code, mode=mode)
                            return {'formatted': False, 'success': True}
                        except black.NothingChanged:
                            return {'formatted': True, 'success': True}
                    else:
                        result = black.format_str(code, mode=mode)
                        return {'code': result, 'success': True}
                except ImportError:
                    return {'error': 'black не установлен: pip install black',
                            'success': False}

            return {'error': f"Неизвестный action: {action}", 'success': False}

        except Exception as e:
            return {'error': str(e), 'success': False}


class WebFormTool(BaseTool):
    """
    Web Form — автоматизация браузера через playwright: заполнение форм, авторизация.
    """

    def __init__(self):
        super().__init__('web_form',
                         'Браузерная автоматизация: заполнение форм, клики, скриншоты (playwright)')

    def run(self, action: str, **kwargs) -> dict:
        """
        open:       url (str), headless=True → инициализирует браузер
        click:      selector (str)
        fill:       selector (str), value (str)
        submit:     selector (str, optional — нажать Enter или кликнуть кнопку)
        get_text:   selector (str)
        screenshot: path (str)
        navigate:   url (str)
        close:      —
        """
        try:
            from playwright.sync_api import sync_playwright  # pyright: ignore[reportMissingImports]

            # Переиспользуем браузер из kwargs или создаём новый при 'open'
            if action == 'open':
                url      = kwargs.get('url', 'about:blank')
                headless = kwargs.get('headless', True)
                pw_ctx   = sync_playwright().__enter__()
                browser  = pw_ctx.chromium.launch(headless=headless)
                page     = browser.new_page()
                page.goto(url, wait_until='networkidle', timeout=30000)
                self._pw  = pw_ctx   # pyright: ignore[reportAttributeAccessIssue]
                self._browser = browser  # pyright: ignore[reportAttributeAccessIssue]
                self._page    = page     # pyright: ignore[reportAttributeAccessIssue]
                return {'url': page.url, 'title': page.title(), 'success': True}

            page = getattr(self, '_page', None)
            if not page and action != 'close':
                return {
                    'error': 'Браузер не открыт. Сначала вызови open с url.',
                    'success': False,
                }
            assert page is not None  # narrowed for type checker

            if action == 'navigate':
                page.goto(kwargs['url'], wait_until='networkidle', timeout=30000)
                return {'url': page.url, 'title': page.title(), 'success': True}

            elif action == 'click':
                page.click(kwargs['selector'], timeout=10000)
                return {'clicked': kwargs['selector'], 'success': True}

            elif action == 'fill':
                page.fill(kwargs['selector'], kwargs.get('value', ''),
                          timeout=10000)
                return {'filled': kwargs['selector'], 'success': True}

            elif action == 'submit':
                sel = kwargs.get('selector')
                if sel:
                    page.click(sel, timeout=10000)
                else:
                    page.keyboard.press('Enter')
                return {'submitted': True, 'success': True}

            elif action == 'get_text':
                text = page.text_content(kwargs['selector'], timeout=10000)
                return {'text': text, 'success': True}

            elif action == 'screenshot':
                path = kwargs.get('path', 'screenshot.png')
                page.screenshot(path=path)
                return {'path': path, 'success': True}

            elif action == 'close':
                browser = getattr(self, '_browser', None)
                pw_ctx  = getattr(self, '_pw', None)
                if browser:
                    browser.close()
                if pw_ctx:
                    pw_ctx.__exit__(None, None, None)
                self._browser = None
                self._page    = None
                self._pw      = None
                return {'closed': True, 'success': True}

            return {'error': f"Неизвестный action: {action}", 'success': False}

        except ImportError:
            return {'error': ('playwright не установлен: '
                              'pip install playwright && playwright install chromium'),
                    'success': False}
        except Exception as e:
            return {'error': str(e), 'success': False}


# ── Группа 10: Очередь задач и самомониторинг ────────────────────────────────

class TaskQueueTool(BaseTool):
    """
    Task Queue — приоритетная очередь задач в памяти процесса.
    Действия: push, pop, peek, list, size, clear
    """

    def __init__(self):
        super().__init__('task_queue',
                         'Приоритетная очередь задач: push, pop, peek, list, clear')
        import heapq
        self._heap: list = []
        self._counter  = 0
        self._lock     = threading.Lock()
        self._heapq    = heapq

    def run(self, action: str, **kwargs) -> dict:
        """
        push:  task (str|dict), priority=5 (1=высший, 10=низший), metadata (dict)
        pop:   — → возвращает задачу с наивысшим приоритетом
        peek:  — → просмотр без удаления
        list:  limit=20 → список задач по приоритету
        size:  — → количество задач
        clear: —
        """
        with self._lock:
            if action == 'push':
                task     = kwargs.get('task', '')
                priority = int(kwargs.get('priority', 5))
                meta     = kwargs.get('metadata', {})
                entry = (priority, self._counter, {
                    'task':     task,
                    'priority': priority,
                    'metadata': meta,
                    'added_at': time.time(),
                })
                self._heapq.heappush(self._heap, entry)
                self._counter += 1
                return {'queued': task, 'priority': priority,
                        'queue_size': len(self._heap), 'success': True}

            elif action == 'pop':
                if not self._heap:
                    return {'error': 'Очередь пуста', 'success': False}
                _, _, item = self._heapq.heappop(self._heap)
                return {'item': item, 'remaining': len(self._heap), 'success': True}

            elif action == 'peek':
                if not self._heap:
                    return {'error': 'Очередь пуста', 'success': False}
                _, _, item = self._heap[0]
                return {'item': item, 'queue_size': len(self._heap), 'success': True}

            elif action == 'list':
                limit = int(kwargs.get('limit', 20))
                items = sorted(self._heap, key=lambda x: (x[0], x[1]))
                return {
                    'tasks': [item for _, _, item in items[:limit]],
                    'total': len(self._heap),
                    'success': True,
                }

            elif action == 'size':
                return {'size': len(self._heap), 'success': True}

            elif action == 'clear':
                self._heap.clear()
                return {'cleared': True, 'success': True}

        return {'error': f"Неизвестный action: {action}", 'success': False}


class SelfMonitorTool(BaseTool):
    """
    Self Monitor — мониторинг логов агента, аномалии, CPU/RAM, алерты.
    """

    _ANOMALY_PATTERNS = [
        re.compile(r'\bERROR\b|\bFAIL\b|\bCRITICAL\b', re.IGNORECASE),
        re.compile(r'exception|traceback|stacktrace', re.IGNORECASE),
        re.compile(r'timeout|timed out', re.IGNORECASE),
        re.compile(r'out of memory|oom|memoryerror', re.IGNORECASE),
    ]

    def __init__(self, log_path: str | None = None):
        super().__init__('self_monitor',
                         'Мониторинг агента: логи, аномалии, CPU/RAM, алерты')
        self.log_path = log_path
        self._alerts: list = []

    def run(self, action: str, **kwargs) -> dict:
        """
        health:        — системные метрики + краткий статус
        read_logs:     lines=50, level=None ('ERROR'|'WARNING'|...), pattern=None
        scan_anomalies: lines=200 → список аномальных строк
        alerts:        — список накопленных алертов
        clear_alerts:  —
        """
        try:
            import psutil

            if action == 'health':
                vm = psutil.virtual_memory()
                cpu = psutil.cpu_percent(interval=0.5)
                disk_path = 'C:\\' if os.name == 'nt' else '/'
                disk = psutil.disk_usage(disk_path)
                status = 'ok'
                warnings = []
                if cpu > 90:
                    warnings.append(f'Высокая нагрузка CPU: {cpu}%')
                    status = 'warning'
                if vm.percent > 90:
                    warnings.append(f'Мало RAM: {round(vm.available/1024**2)}MB свободно')
                    status = 'warning'
                return {
                    'status':         status,
                    'cpu_percent':    cpu,
                    'ram_percent':    vm.percent,
                    'ram_free_mb':    round(vm.available / 1024**2, 1),
                    'disk_percent':   disk.percent,
                    'disk_free_gb':   round(disk.free / 1024**3, 1),
                    'warnings':       warnings,
                    'alerts_count':   len(self._alerts),
                    'success': True,
                }

            elif action == 'read_logs':
                if not self.log_path or not os.path.exists(self.log_path):
                    return {'error': 'log_path не задан или файл не существует',
                            'success': False}
                lines_n = int(kwargs.get('lines', 50))
                level   = kwargs.get('level')
                pattern = kwargs.get('pattern')
                with open(self.log_path, 'r', encoding='utf-8', errors='replace') as f:
                    all_lines = f.readlines()
                tail = all_lines[-lines_n:]
                if level:
                    tail = [ln for ln in tail if level.upper() in ln.upper()]
                if pattern:
                    rx = re.compile(pattern, re.IGNORECASE)
                    tail = [ln for ln in tail if rx.search(ln)]
                return {'lines': [ln.rstrip() for ln in tail],
                        'count': len(tail), 'success': True}

            elif action == 'scan_anomalies':
                if not self.log_path or not os.path.exists(self.log_path):
                    return {'error': 'log_path не задан или файл не существует',
                            'success': False}
                lines_n = int(kwargs.get('lines', 200))
                with open(self.log_path, 'r', encoding='utf-8', errors='replace') as f:
                    tail = f.readlines()[-lines_n:]
                anomalies = []
                for ln in tail:
                    for rx in self._ANOMALY_PATTERNS:
                        if rx.search(ln):
                            anomalies.append(ln.rstrip())
                            # Сохраняем как алерт
                            self._alerts.append({
                                'line': ln.rstrip(),
                                'time': time.time(),
                            })
                            if len(self._alerts) > 500:
                                self._alerts.pop(0)
                            break
                return {'anomalies': anomalies, 'count': len(anomalies), 'success': True}

            elif action == 'alerts':
                limit = int(kwargs.get('limit', 50))
                return {'alerts': self._alerts[-limit:],
                        'total': len(self._alerts), 'success': True}

            elif action == 'clear_alerts':
                self._alerts.clear()
                return {'cleared': True, 'success': True}

            return {'error': f"Неизвестный action: {action}", 'success': False}

        except ImportError:
            return {'error': 'psutil не установлен: pip install psutil', 'success': False}
        except Exception as e:
            return {'error': str(e), 'success': False}


# ── GUI Agent — видит экран → понимает → кликает ─────────────────────────────

class GUIAgentTool(BaseTool):
    """
    GUI Agent — «видит экран → понимает → действует».

    Цикл:
        1. Снимает скриншот (mss / PIL fallback)
        2. Отправляет в GPT-4o Vision с описанием цели
        3. Получает следующее действие в JSON (click/type/hotkey/key/scroll/wait/done)
        4. Выполняет через pyautogui
        5. Повторяет до достижения цели или max_steps

    Позволяет агенту работать с любым GUI-приложением на экране:
    браузерами, Office, Adobe, 3D-редакторами, системными диалогами.
    """

    _SYSTEM_PROMPT = (
        'Ты — GUI-агент. Ты видишь скриншот экрана Windows. Твоя роль — выполнить задачу '
        'шаг за шагом, управляя мышью и клавиатурой.\n\n'
        'Отвечай ТОЛЬКО валидным JSON (без markdown-блоков, без пояснений):\n'
        '{\n'
        '  "action": "click"|"double_click"|"right_click"|"type"|"hotkey"|"key"|"scroll"|"wait"|"done",\n'
        '  "x": число,           // пиксельная координата X (для click/double_click/right_click/scroll)\n'
        '  "y": число,           // пиксельная координата Y\n'
        '  "text": "строка",     // для action=type\n'
        '  "keys": ["ctrl","s"], // для action=hotkey\n'
        '  "key": "enter",       // для action=key (одна клавиша)\n'
        '  "clicks": 3,          // для scroll (+ вверх, - вниз)\n'
        '  "seconds": 1.0,       // для wait\n'
        '  "reason": "..."       // краткое объяснение шага\n'
        '}\n\n'
        'Правила:\n'
        '- action="done" когда задача полностью выполнена или невозможна\n'
        '- Никакого текста кроме JSON\n'
        '- Будь точен с координатами — смотри на скриншот внимательно\n'
        '- Для ввода текста на русском используй action=type'
    )

    def __init__(self, openai_client=None):
        super().__init__(
            'gui_agent',
            'GUI-агент: видит экран через GPT-4o Vision → кликает, печатает, управляет любым окном'
        )
        self._llm = openai_client

    def run(self, goal: str, max_steps: int = 10, pause: float = 1.0,
            region: tuple | None = None) -> dict:
        """
        Args:
            goal      — что сделать, например «открой Блокнот и напиши Hello»
            max_steps — защита от зацикливания (макс. шагов)
            pause     — пауза между шагами в секундах
            region    — (x, y, w, h) снимать только часть экрана; None = весь экран
        Returns:
            {'success': bool, 'goal': str, 'steps': list, 'reason': str}
        """
        if not self._llm:
            return {'success': False, 'error': 'openai_client не подключён', 'goal': goal}

        # Проверяем зависимости один раз
        try:
            import pyautogui  # pyright: ignore[reportMissingImports]
            import openai as _openai  # pyright: ignore[reportMissingImports]
        except ImportError as e:
            return {'success': False, 'error': f'Зависимость не установлена: {e}. Выполни: pip install pyautogui openai', 'goal': goal}

        import base64
        import json
        import re as _re
        import time as _time

        pyautogui.FAILSAFE = True
        pyautogui.PAUSE = 0.1

        api_key = (getattr(self._llm, 'api_key', None)
                   or os.environ.get('OPENAI_API_KEY', ''))
        oa = _openai.OpenAI(api_key=api_key)

        steps_log: list[dict] = []

        for step_n in range(1, max_steps + 1):
            # ── 1. Скриншот ──────────────────────────────────────────────────
            img_path = self._screenshot(region)
            if not img_path:
                return {'success': False, 'error': 'Не удалось сделать скриншот',
                        'steps': steps_log, 'goal': goal}

            with open(img_path, 'rb') as _f:
                b64 = base64.b64encode(_f.read()).decode()
            try:
                os.unlink(img_path)
            except OSError:
                pass

            # ── 2. Vision API ─────────────────────────────────────────────────
            user_msg = (
                f'Задача: {goal}\n'
                f'Шаг {step_n} из максимум {max_steps}.\n'
                f'Что нужно сделать на этом скриншоте, чтобы продвинуться к цели?'
            )
            try:
                resp = oa.chat.completions.create(
                    model='gpt-4o',
                    messages=[  # type: ignore[arg-type]
                        {'role': 'system', 'content': self._SYSTEM_PROMPT},
                        {'role': 'user', 'content': [
                            {'type': 'text', 'text': user_msg},
                            {'type': 'image_url', 'image_url': {
                                'url': f'data:image/png;base64,{b64}',
                                'detail': 'high',
                            }},
                        ]},
                    ],
                    max_tokens=400,
                    temperature=0.1,
                )
                raw = resp.choices[0].message.content or ''
            except Exception as e:
                return {'success': False, 'error': f'Vision API: {e}',
                        'steps': steps_log, 'goal': goal}

            # ── 3. Парсим JSON ────────────────────────────────────────────────
            action_data = self._parse_json(raw)
            step_entry: dict = {'step': step_n, 'raw': raw[:300], 'action': action_data}
            steps_log.append(step_entry)

            if not action_data:
                step_entry['error'] = f'JSON не распознан: {raw[:200]}'
                continue

            act = action_data.get('action', '')

            # ── 4. Готово ─────────────────────────────────────────────────────
            if act == 'done':
                return {
                    'success': True,
                    'goal': goal,
                    'steps': steps_log,
                    'reason': action_data.get('reason', 'Задача выполнена'),
                }

            # ── 5. Выполняем действие ─────────────────────────────────────────
            result = self._do(pyautogui, action_data)
            step_entry['result'] = result
            _time.sleep(pause)

        return {
            'success': False,
            'error': f'Превышено max_steps={max_steps}',
            'steps': steps_log,
            'goal': goal,
        }

    # ── Helpers ───────────────────────────────────────────────────────────────

    def _screenshot(self, region: tuple | None) -> str | None:
        """Снимает скриншот, возвращает путь к PNG-файлу."""
        import tempfile
        fd, path = tempfile.mkstemp(suffix='.png')
        os.close(fd)
        try:
            import mss  # pyright: ignore[reportMissingImports]
            import mss.tools
            with mss.mss() as sct:
                if region:
                    x, y, w, h = region
                    mon: dict = {'left': x, 'top': y, 'width': w, 'height': h}
                else:
                    mon = sct.monitors[0]
                img = sct.grab(mon)
                mss.tools.to_png(img.rgb, img.size, output=path)
            return path
        except ImportError:
            pass
        try:
            from PIL import ImageGrab  # pyright: ignore[reportMissingImports]
            bbox = (region[0], region[1], region[0] + region[2], region[1] + region[3]) if region else None
            pil_img = ImageGrab.grab(bbox=bbox)
            pil_img.save(path)
            return path
        except Exception:
            return None

    @staticmethod
    def _parse_json(raw: str) -> dict | None:
        """Извлекает JSON из ответа GPT-4o."""
        import json
        # Убираем возможные markdown-блоки ```json ... ```
        cleaned = raw.strip()
        cleaned = re.sub(r'^```(?:json)?\s*', '', cleaned)
        cleaned = re.sub(r'\s*```$', '', cleaned)
        # Ищем первый {...}
        m = re.search(r'\{.*\}', cleaned, re.DOTALL)
        if not m:
            return None
        try:
            return json.loads(m.group())
        except (json.JSONDecodeError, ValueError):
            return None

    @staticmethod
    def _do(pg, data: dict) -> dict:
        """Выполняет одно действие через pyautogui."""
        act = data.get('action', '')
        try:
            if act in ('click', 'double_click', 'right_click'):
                x, y = int(data['x']), int(data['y'])
                if act == 'click':
                    pg.click(x, y)
                elif act == 'double_click':
                    pg.doubleClick(x, y)
                else:
                    pg.rightClick(x, y)
                return {'success': True, 'x': x, 'y': y}

            elif act == 'type':
                text = data.get('text', '')
                # pyautogui.typewrite не поддерживает Unicode/кириллицу
                # Используем clipboard paste если доступен pyperclip
                try:
                    import pyperclip  # pyright: ignore[reportMissingImports]
                    pyperclip.copy(text)
                    pg.hotkey('ctrl', 'v')
                except ImportError:
                    # Fallback: посимвольный ввод (только ASCII)
                    pg.typewrite(str(text), interval=0.04)
                return {'success': True, 'length': len(text)}

            elif act == 'hotkey':
                keys = data.get('keys', [])
                if isinstance(keys, str):
                    keys = [k.strip() for k in keys.split('+')]
                pg.hotkey(*keys)
                return {'success': True, 'keys': keys}

            elif act == 'key':
                pg.press(data.get('key', ''))
                return {'success': True, 'key': data.get('key')}

            elif act == 'scroll':
                x = int(data.get('x') or 0)
                y = int(data.get('y') or 0)
                clicks = int(data.get('clicks', 3))
                if x and y:
                    pg.scroll(clicks, x=x, y=y)
                else:
                    pg.scroll(clicks)
                return {'success': True, 'clicks': clicks}

            elif act == 'wait':
                secs = min(float(data.get('seconds', 1.0)), 10.0)
                time.sleep(secs)
                return {'success': True, 'seconds': secs}

            return {'success': False, 'error': f'Неизвестный action: {act}'}
        except Exception as e:
            return {'success': False, 'error': str(e)}


# ── Tool Registry — реестр всех инструментов ─────────────────────────────────

class ToolLayer:
    """
    Tool Layer — реестр и точка доступа ко всем инструментам (Слой 5).

    Cognitive Core и Agent System обращаются к инструментам через ToolLayer.
    Позволяет регистрировать, получать и вызывать инструменты по имени.
    """

    def __init__(self):
        self._tools: dict[str, BaseTool] = {}

    def register(self, tool: BaseTool):
        """Регистрирует инструмент в реестре."""
        self._tools[tool.name] = tool

    def get(self, name: str) -> BaseTool | None:
        """Возвращает инструмент по имени."""
        return self._tools.get(name)

    def use(self, tool_name: str, **kwargs):
        """Вызывает инструмент по имени с переданными аргументами."""
        tool = self._tools.get(tool_name)
        if not tool:
            raise KeyError(f"Инструмент '{tool_name}' не зарегистрирован. Доступны: {self.list()}")
        return tool(**kwargs)

    def list(self) -> list[str]:
        """Список имён всех зарегистрированных инструментов."""
        return list(self._tools.keys())

    def describe(self) -> list[dict]:
        """Описание всех инструментов (для Cognitive Core)."""
        return [{'name': t.name, 'description': t.description} for t in self._tools.values()]


class RedditTool(BaseTool):
    """
    Reddit — публикация постов, комментариев, мониторинг сабреддитов через PRAW.

    Переменные окружения:
        REDDIT_CLIENT_ID      — Client ID приложения
        REDDIT_CLIENT_SECRET  — Client Secret
        REDDIT_USERNAME       — логин аккаунта
        REDDIT_PASSWORD       — пароль аккаунта
        REDDIT_USER_AGENT     — строка вида 'MyApp/1.0 by u/YourUsername'

    Регистрация приложения (бесплатно):
        reddit.com → User Settings → Safety & Privacy → Manage third-party app access
        → Create another app... → тип: script → получишь client_id и secret
    """

    def __init__(
        self,
        client_id: str | None = None,
        client_secret: str | None = None,
        username: str | None = None,
        password: str | None = None,
        user_agent: str | None = None,
    ):
        super().__init__('reddit', 'Публикация постов и мониторинг Reddit через PRAW')
        import os as _os
        self._client_id     = client_id     or _os.environ.get('REDDIT_CLIENT_ID', '')
        self._client_secret = client_secret or _os.environ.get('REDDIT_CLIENT_SECRET', '')
        self._username      = username      or _os.environ.get('REDDIT_USERNAME', '')
        self._password      = password      or _os.environ.get('REDDIT_PASSWORD', '')
        self._user_agent    = user_agent    or _os.environ.get(
            'REDDIT_USER_AGENT', f'AgentBot/1.0 by u/{self._username or "user"}'
        )
        self._reddit = None

    def _get_reddit(self):
        if self._reddit is not None:
            return self._reddit
        if not all([self._client_id, self._client_secret, self._username, self._password]):
            raise ValueError(
                'Reddit credentials не заданы. '
                'Добавь в .env: REDDIT_CLIENT_ID, REDDIT_CLIENT_SECRET, '
                'REDDIT_USERNAME, REDDIT_PASSWORD'
            )
        import praw
        self._reddit = praw.Reddit(
            client_id=self._client_id,
            client_secret=self._client_secret,
            username=self._username,
            password=self._password,
            user_agent=self._user_agent,
        )
        return self._reddit

    def post(self, subreddit: str, title: str, text: str = '', url: str = '') -> dict:
        """
        Опубликовать пост в сабреддите.
        text — текстовый пост (selfpost)
        url  — ссылочный пост (link post); если задан, text игнорируется
        """
        try:
            r = self._get_reddit()
            sub = r.subreddit(subreddit)
            if url:
                submission = sub.submit(title=title, url=url)
            else:
                submission = sub.submit(title=title, selftext=text)
            if submission is None:
                return {'ok': False, 'error': 'Reddit API вернул None — вероятно read-only режим или нет прав на публикацию.', 'success': False}
            return {
                'ok': True,
                'post_id':  submission.id,
                'url':      f'https://reddit.com{submission.permalink}',
                'subreddit': subreddit,
                'title':    title,
                'success':  True,
            }
        except Exception as e:
            return {'ok': False, 'error': str(e), 'subreddit': subreddit, 'success': False}

    def comment(self, post_id: str, text: str) -> dict:
        """Оставить комментарий к посту по его ID."""
        try:
            r = self._get_reddit()
            submission = r.submission(id=post_id)
            comment = submission.reply(text)
            if comment is None:
                return {'ok': False, 'error': 'Reddit API вернул None — вероятно read-only режим или нет прав на комментирование.', 'success': False}
            return {'ok': True, 'comment_id': comment.id, 'success': True}
        except Exception as e:
            return {'ok': False, 'error': str(e), 'success': False}

    def search(self, subreddit: str, query: str, limit: int = 10) -> dict:
        """Поиск постов в сабреддите."""
        try:
            r = self._get_reddit()
            results = []
            for s in r.subreddit(subreddit).search(query, limit=limit):
                results.append({
                    'id':    s.id,
                    'title': s.title,
                    'score': s.score,
                    'url':   f'https://reddit.com{s.permalink}',
                    'author': str(s.author),
                })
            return {'ok': True, 'results': results, 'count': len(results), 'success': True}
        except Exception as e:
            return {'ok': False, 'error': str(e), 'success': False}

    def hot(self, subreddit: str, limit: int = 10) -> dict:
        """Топ постов (hot) из сабреддита."""
        try:
            r = self._get_reddit()
            results = []
            for s in r.subreddit(subreddit).hot(limit=limit):
                results.append({
                    'id':    s.id,
                    'title': s.title,
                    'score': s.score,
                    'url':   f'https://reddit.com{s.permalink}',
                })
            return {'ok': True, 'results': results, 'count': len(results), 'success': True}
        except Exception as e:
            return {'ok': False, 'error': str(e), 'success': False}

    def me(self) -> dict:
        """Информация о текущем аккаунте (проверка авторизации)."""
        try:
            r = self._get_reddit()
            me = r.user.me()
            if me is None:
                return {
                    'ok': False,
                    'error': 'Reddit API работает в read-only режиме — авторизация пользователя отсутствует. '
                             'Укажи reddit_username/reddit_password в config/credentials.json.',
                    'success': False,
                }
            return {
                'ok':        True,
                'username':  me.name,
                'karma':     me.link_karma + me.comment_karma,
                'verified':  me.verified,
                'success':   True,
            }
        except Exception as e:
            return {'ok': False, 'error': str(e), 'success': False}

    def run(self, action: str = 'me', **params) -> dict:
        """
        action: me | post | comment | search | hot
        """
        actions = {
            'me':      self.me,
            'post':    self.post,
            'comment': self.comment,
            'search':  self.search,
            'hot':     self.hot,
        }
        fn = actions.get(action)
        if not fn:
            return {'ok': False, 'error': f'Неизвестный action: {action}. Доступные: {list(actions)}', 'success': False}
        try:
            return fn(**params) or {'ok': True, 'success': True}
        except TypeError as e:
            return {'ok': False, 'error': f'Неверные параметры: {e}', 'success': False}


def build_tool_layer(
    working_dir: str | None = None,
    db_connection=None,
    github_token: str | None = None,
    github_backend=None,
    search_backend=None,
    cloud_clients: dict | None = None,
    openai_client=None,
    # Email
    email_username: str | None = None,
    email_password: str | None = None,
    email_smtp_host: str = 'smtp.gmail.com',
    email_smtp_port: int = 587,
    email_imap_host: str = 'imap.gmail.com',
    email_imap_port: int = 993,
    # Google Calendar
    google_credentials_path: str | None = None,
    # HuggingFace Hub
    hf_token: str | None = None,
    # Upwork API (OAuth 2.0)
    upwork_client_id: str | None = None,
    upwork_client_secret: str | None = None,
    upwork_access_token: str | None = None,
    upwork_refresh_token: str | None = None,
    # Figma REST API
    figma_token: str | None = None,
    # Twilio / Vonage (голосовые звонки)
    twilio_sid: str | None = None,
    twilio_token: str | None = None,
    twilio_from: str | None = None,
    vonage_key: str | None = None,
    vonage_secret: str | None = None,
    vonage_from: str | None = None,
    # Blender path
    blender_path: str | None = None,
    # OpenSCAD path (CAD)
    openscad_path: str | None = None,
    # ADB / Appium (Mobile)
    adb_path: str | None = None,
    appium_host: str | None = None,
    # Reddit (PRAW)
    reddit_client_id: str | None = None,
    reddit_client_secret: str | None = None,
    reddit_username: str | None = None,
    reddit_password: str | None = None,
    reddit_user_agent: str | None = None,
) -> ToolLayer:
    """
    Фабрика: создаёт ToolLayer с полным набором инструментов (28 шт.; + browser через agent.py = 29).

    Args:
        working_dir    — рабочая директория для terminal, filesystem, git
        db_connection  — подключение к БД (если None — автоматически SQLite)
        github_token   — токен GitHub
        github_backend — готовый GitHub-клиент
        search_backend — поисковый backend
        cloud_clients  — dict {provider_name: client}
        openai_client  — OpenAIClient для image_generator, ocr, translate
    """
    import os as _os
    terminal = TerminalTool(working_dir=working_dir)

    # ── Авто-подключение SQLite если нет внешнего DB-соединения ──────────────
    if db_connection is None:
        import sqlite3 as _sqlite3
        _db_path = _os.path.join(working_dir or '.', 'agent_memory.db')
        db_connection = _sqlite3.connect(_db_path, check_same_thread=False)

    # ── Адаптер OpenAI → CloudAPITool.client ─────────────────────────────────
    class _OpenAICloudAdapter:
        """Тонкая обёртка: CloudAPITool.client.call() → openai_client.infer()."""
        def __init__(self, client):
            self._llm = client
        def call(self, _endpoint: str, payload: dict | None = None, _method: str = 'POST'):
            prompt = (payload or {}).get('prompt') or str(payload)
            return self._llm.infer(prompt)

    layer = ToolLayer()

    # ── Оригинальные 9 ───────────────────────────────────────────────────────
    layer.register(terminal)
    layer.register(FileSystemTool(base_dir=working_dir))
    import os as _os
    _outputs_dir = _os.path.join(working_dir, 'outputs') if working_dir else 'outputs'
    _os.makedirs(_outputs_dir, exist_ok=True)
    layer.register(PythonRuntimeTool(working_dir=working_dir))
    layer.register(SearchTool(backend=search_backend))
    layer.register(GitHubTool(token=github_token, backend=github_backend))
    layer.register(DockerTool(terminal=terminal))
    layer.register(DatabaseTool(connection=db_connection))
    layer.register(PackageManagerTool(terminal=terminal))

    # CloudAPITool: всегда регистрируем openai-провайдер (+ любые внешние клиенты)
    if openai_client:
        layer.register(CloudAPITool(provider='openai', client=_OpenAICloudAdapter(openai_client)))
    for provider, client in (cloud_clients or {}).items():
        layer.register(CloudAPITool(provider=provider, client=client))

    # ── Управление компьютером (4 новых) ─────────────────────────────────────
    layer.register(ScreenshotTool())
    layer.register(MouseKeyboardTool())
    layer.register(ClipboardTool())
    layer.register(WindowManagerTool())

    # ── Системные (4 новых) ───────────────────────────────────────────────────
    layer.register(ProcessManagerTool())
    layer.register(GitTool(working_dir=working_dir))
    layer.register(NetworkTool())
    layer.register(CronSchedulerTool())
    layer.register(HTTPClientTool())

    # ── AI-инструменты (3) ────────────────────────────────────────────────────
    layer.register(ImageGeneratorTool(openai_client=openai_client))
    layer.register(OCRTool(openai_client=openai_client))
    layer.register(TranslateTool(openai_client=openai_client))

    # ── Данные, документы, визуализация (3) ──────────────────────────────────
    layer.register(SpreadsheetTool())
    layer.register(PDFGeneratorTool())
    layer.register(PowerPointTool())
    layer.register(DataVizTool())

    # ── Коммуникации (2) ─────────────────────────────────────────────────────
    layer.register(EmailTool(
        smtp_host=email_smtp_host,
        smtp_port=email_smtp_port,
        imap_host=email_imap_host,
        imap_port=email_imap_port,
        username=email_username,
        password=email_password,
    ))
    layer.register(CalendarTool(
        credentials_path=google_credentials_path,
        token_path=_os.path.join(_os.path.dirname(google_credentials_path or 'credentials.json'), 'token.json'),
    ))
    layer.register(HuggingFaceTool(token=hf_token))

    # ── Upwork ───────────────────────────────────────────────────────────────
    try:
        from tools.upwork_tool import UpworkTool
        layer.register(UpworkTool(
            client_id=upwork_client_id,
            client_secret=upwork_client_secret,
            access_token=upwork_access_token,
            refresh_token=upwork_refresh_token,
        ))
    except Exception:
        pass

    # ── Figma ─────────────────────────────────────────────────────────────
    try:
        from tools.figma_tool import FigmaTool
        layer.register(FigmaTool(access_token=figma_token))
    except Exception:
        pass

    # ── ImageEdit (Pillow + GIMP + Inkscape + ImageMagick) ───────────────
    try:
        from tools.image_edit_tool import ImageEditTool
        layer.register(ImageEditTool())
    except Exception:
        pass

    # ── VideoEdit (ffmpeg + AE/Premiere ExtendScript) ─────────────────────
    try:
        from tools.video_edit_tool import VideoEditTool
        layer.register(VideoEditTool())
    except Exception:
        pass

    # ── Blender (headless 3D) ─────────────────────────────────────────────
    try:
        from tools.blender_tool import BlenderTool
        layer.register(BlenderTool(blender_path=blender_path))
    except Exception:
        pass

    # ── CAD (ezdxf DXF + OpenSCAD + AutoCAD COM) ─────────────────────────
    try:
        from tools.cad_tool import CADTool
        layer.register(CADTool(openscad_path=openscad_path))
    except Exception:
        pass

    # ── VoiceCall (Twilio / Vonage) ───────────────────────────────────────
    try:
        from tools.voice_call_tool import VoiceCallTool
        layer.register(VoiceCallTool(
            twilio_sid=twilio_sid,
            twilio_token=twilio_token,
            twilio_from=twilio_from,
            vonage_key=vonage_key,
            vonage_secret=vonage_secret,
            vonage_from=vonage_from,
        ))
    except Exception:
        pass

    # ── Mobile (ADB + Appium + iOS) ───────────────────────────────────────
    try:
        from tools.mobile_tool import MobileTool
        layer.register(MobileTool(adb_path=adb_path, appium_host=appium_host))
    except Exception:
        pass

    # ── Архивы, SSH, слежение за файлами (3) ─────────────────────────────────
    layer.register(ArchiveTool())
    layer.register(SSHTool())
    layer.register(FileWatcherTool())

    # ── Безопасность и уведомления (2) ───────────────────────────────────────
    layer.register(NotificationTool())
    layer.register(EncryptionTool())

    # ── NLP и ML (2) ─────────────────────────────────────────────────────────
    layer.register(NLPTool(openai_client=openai_client))
    layer.register(EmbeddingTool(openai_client=openai_client))

    # ── Разработка и инфраструктура (2) ──────────────────────────────────────
    layer.register(CodeAnalyzerTool())
    layer.register(WebFormTool())

    # ── Очередь задач и самомониторинг (2) ───────────────────────────────────
    layer.register(TaskQueueTool())
    layer.register(SelfMonitorTool())

    # ── GUI-агент: видит экран → понимает → кликает ───────────────────────────
    layer.register(GUIAgentTool(openai_client=openai_client))

    # ── Reddit (PRAW) ─────────────────────────────────────────────────────────
    layer.register(RedditTool(
        client_id=reddit_client_id,
        client_secret=reddit_client_secret,
        username=reddit_username,
        password=reddit_password,
        user_agent=reddit_user_agent,
    ))

    return layer
